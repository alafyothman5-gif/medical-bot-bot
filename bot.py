#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Smart Medical Assistant Bot - OpenRouter Paid Economical Edition
OSCE SIMULATOR V4 (Step-by-Step Clinical Assessment)
Past Papers UI (Universities / Subjects)
Admin Staging Pipeline & Vision API Integration
"""

import asyncio
import hashlib
import io
import json
import logging
import os
import random
import re
import sqlite3
import time
import uuid
import base64
import zipfile
import mimetypes
import shutil
import tempfile
import subprocess
import contextvars
from pathlib import Path
from datetime import datetime, timedelta
from typing import Any, Dict, List, Optional, Tuple

import aiohttp
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from telegram import BotCommand, InlineKeyboardButton, InlineKeyboardMarkup, Update
from telegram.ext import (
    Application,
    ApplicationBuilder,
    CallbackQueryHandler,
    CommandHandler,
    ContextTypes,
    MessageHandler,
    ApplicationHandlerStop,
    filters,
)

try:
    from dotenv import load_dotenv
    load_dotenv()
except Exception:
    pass

# =============================================================================
# LOGGING & ENVIRONMENT
# =============================================================================

logging.basicConfig(format="%(asctime)s - %(name)s - %(levelname)s - %(message)s", level=logging.INFO)
logger = logging.getLogger("medical_bot_openrouter")

TOKEN = os.environ.get("TELEGRAM_BOT_TOKEN")
def _env_bool(name: str, default: str = "0") -> bool:
    return os.environ.get(name, default).strip().lower() in {"1", "true", "yes", "on"}

def _env_keys(*names: str) -> List[str]:
    raw_parts = []
    for name in names:
        val = os.environ.get(name, "")
        if val:
            raw_parts.append(val)
    raw = "\n".join(raw_parts)
    return [k.strip() for k in re.split(r"[,;\s]+", raw) if k.strip()]

OPENROUTER_API_KEY = os.environ.get("OPENROUTER_API_KEY")
OPENROUTER_PAID_KEYS = _env_keys("OPENROUTER_PAID_KEYS", "OPENROUTER_API_KEYS", "OPENROUTER_API_KEY")
OPENROUTER_FREE_KEYS = _env_keys("OPENROUTER_FREE_KEYS", "OPENROUTER_FREE_API_KEYS", "OPENROUTER_FREE_API_KEY", "OPENROUTER_API_KEY")
OPENROUTER_FREE_MODEL = os.environ.get("OPENROUTER_FREE_MODEL", "openrouter/free")
GROQ_API_KEYS = _env_keys("GROQ_API_KEYS", "GROQ_API_KEY")
GROQ_MODEL = os.environ.get("GROQ_MODEL", "llama-3.1-8b-instant")
CLOUDFLARE_ACCOUNT_ID = os.environ.get("CLOUDFLARE_ACCOUNT_ID", "").strip()
CLOUDFLARE_API_TOKENS = _env_keys("CLOUDFLARE_API_TOKENS", "CLOUDFLARE_API_TOKEN")
CLOUDFLARE_MODEL = os.environ.get("CLOUDFLARE_MODEL", "@cf/meta/llama-3.1-8b-instruct").strip()
ALLOW_PAID_FALLBACK_FOR_USERS = _env_bool("ALLOW_PAID_FALLBACK_FOR_USERS", "0")
ALLOW_PAID_FALLBACK_FOR_PREMIUM = _env_bool("ALLOW_PAID_FALLBACK_FOR_PREMIUM", "1")
ALLOW_PAID_FALLBACK_FOR_ADMIN = _env_bool("ALLOW_PAID_FALLBACK_FOR_ADMIN", "1")
PREMIUM_ENABLED = _env_bool("PREMIUM_ENABLED", "1")
FREE_PROVIDER_MODE = _env_bool("FREE_PROVIDER_MODE", "1")
GOOGLE_API_KEYS = _env_keys("GOOGLE_API_KEYS", "GEMINI_API_KEYS", "GOOGLE_API_KEY", "GEMINI_API_KEY")
GEMINI_DIRECT_MODEL = os.environ.get("GEMINI_DIRECT_MODEL", os.environ.get("GOOGLE_GEMINI_MODEL", "gemini-2.5-flash-lite"))
GEMINI_MODEL_NAME = os.environ.get("GEMINI_MODEL", "google/gemini-2.5-flash-lite")
ADMIN_DOCUMENT_REPAIR_FIX_V2_OPENROUTER_VISION_FALLBACK = True
_AI_USER_ID: contextvars.ContextVar[int] = contextvars.ContextVar("ai_user_id", default=0)

WELCOME_GIF = os.environ.get("WELCOME_GIF", os.environ.get("WELCOME_GIF_URL", os.environ.get("WELCOME_GIF_FILE_ID", ""))).strip()
DATABASE_PATH = os.environ.get("DATABASE_PATH", "/workspace/bot_cache.sqlite3")
MAX_FILE_SIZE_BYTES = 30 * 1024 * 1024
ADMIN_PRELOAD_ZIP_MAX_FILES = int(os.environ.get("ADMIN_PRELOAD_ZIP_MAX_FILES", "100"))
ADMIN_PRELOAD_ZIP_MAX_TOTAL_BYTES = int(os.environ.get("ADMIN_PRELOAD_ZIP_MAX_TOTAL_MB", "200")) * 1024 * 1024
MAX_TEXT_CHARS_FOR_AI = int(os.environ.get("MAX_TEXT_CHARS_FOR_AI", "25000"))
ADMIN_VISION_MAX_PAGES = int(os.environ.get("ADMIN_VISION_MAX_PAGES", "18"))
OFFICE_CONVERT_TIMEOUT = int(os.environ.get("OFFICE_CONVERT_TIMEOUT", "90"))
ECONOMY_MODE = os.environ.get("ECONOMY_MODE", "1").strip().lower() in {"1", "true", "yes", "on"}
FREE_USAGE_INTERVALS = {"summary": 24*3600, "basic": 24*3600, "cases": 48*3600, "challenge": 48*3600, "osce": 7*24*3600}
PLAN_LABELS = {"free": "Free", "premium": "Premium", "premium_plus": "Premium Plus", "admin": "Admin"}
PLAN_PRICES_AR = {
    "premium_month": "15 د.ل / شهر",
    "premium_semester": "40 د.ل / فصل كامل 4 شهور",
    "premium_plus_semester": "60 د.ل / فصل كامل 4 شهور",
}
PLAN_USAGE_LIMITS = {
    "free": {
        "summary": {"limit": 1, "window": 24*3600},
        "basic": {"limit": 1, "window": 24*3600},
        "cases": {"limit": 1, "window": 48*3600},
        "challenge": {"limit": 1, "window": 48*3600},
        "osce": {"limit": 1, "window": 7*24*3600},
    },
    "premium": {
        "summary": {"limit": 3, "window": 24*3600},
        "basic": {"limit": 3, "window": 24*3600},
        "cases": {"limit": 2, "window": 24*3600},
        "challenge": {"limit": 2, "window": 24*3600},
        "osce": {"limit": 2, "window": 7*24*3600},
    },
    "premium_plus": {
        "summary": {"limit": 5, "window": 24*3600},
        "basic": {"limit": 5, "window": 24*3600},
        "cases": {"limit": 3, "window": 24*3600},
        "challenge": {"limit": 3, "window": 24*3600},
        "osce": {"limit": 1, "window": 24*3600},
    },
    "admin": {},
}
MAX_CONCURRENT_GEMINI = int(os.environ.get("MAX_CONCURRENT_GEMINI", "4"))
COOLDOWN_SECONDS = int(os.environ.get("COOLDOWN_SECONDS", "20"))
USE_QUEUE = os.environ.get("USE_QUEUE", "1").strip().lower() in {"1", "true", "yes", "on"}
ADMIN_USER_IDS = {1692255414}
ADMIN_LOG_CHAT_ID = int(os.environ.get("ADMIN_LOG_CHAT_ID", "-1003987051773"))

_GEMINI_SEMAPHORE = asyncio.Semaphore(MAX_CONCURRENT_GEMINI)
_ACTIVE_AI_JOBS = 0
_WAITING_AI_JOBS = 0
_AI_QUEUE_LOCK = asyncio.Lock()
_GOOGLE_KEY_INDEX = 0
_generation_locks: Dict[str, asyncio.Lock] = {}

def get_generation_lock(key: str) -> asyncio.Lock:
    if key not in _generation_locks:
        _generation_locks[key] = asyncio.Lock()
    return _generation_locks[key]

# =============================================================================
# TRANSLATIONS
# =============================================================================

SUPPORTED_LANGS = ("en", "ar")

TRANSLATIONS = {
    "en": {
        "choose_lang": "🌐 Choose your language\nاختر لغتك",
        "welcome": "Welcome to your Smart Medical Assistant 🩺\n\nSend a medical PDF, DOCX, PPTX, or TXT file.",
        "help": "Send a medical file up to 30MB. Choose Quiz or Summary. Commands: /start /stop /stats /review /language /support",
        "choose_mode": "📂 Choose mode:",
        "mode_quiz": "📝 Quiz",
        "mode_summary": "📚 Summary",
        "mode_osce": "🎭 OSCE Simulator",
        "mode_past_papers": "📚 Past Papers",
        "mode_progress": "📊 Progress & XP",
        "mode_leaderboard": "🏆 Leaderboard",
        "mode_subscription": "💎 Subscription",
        "mode_quiz_prompt": "📝 Quiz mode selected. Send your file.",
        "mode_summary_prompt": "📚 Summary mode selected. Choose style, then send your file.",
        "mode_osce_prompt": "🎭 OSCE mode selected. Send a medical lecture file, and I will create an interactive OSCE case.",
        "summary_style_prompt": "📋 Choose summary style:",
        "style_disease": "🩺 Disease Profile",
        "style_highyield": "🧠 High-Yield",
        "style_osce": "📋 OSCE",
        "file_received": "📄 File received — extracting text...",
        "file_too_large": "⚠️ File is larger than 30MB. Please send a smaller file.",
        "file_unsupported": "⚠️ Unsupported file type. Send PDF, DOCX, PPTX, or TXT.",
        "file_no_text": "⚠️ I could not extract readable text from this file.",
        "file_empty": "⚠️ File is empty or contains only images (no readable text).",
        "file_ready_quiz": "✅ File ready. Choose question level:",
        "file_ready_summary": "✅ File ready. Preparing summary...",
        "cached_file": "♻️ Same file found in cache. No extra AI cost.",
        "level_basic": "🟢 Basic",
        "level_cases": "🩺 Cases",
        "level_challenge": "🧠 Challenge",
        "basic_count_prompt": "How many Basic questions do you want?",
        "generating_pool": "⏳ Preparing question bank. This happens once per file...",
        "generating_summary": "⏳ Preparing summary. This happens once per file/style...",
        "not_enough": "⚠️ This lecture only has {n} available questions for this level. Starting with available questions.",
        "quiz_ready": "✅ Quiz ready: {n} questions.",
        "no_questions": "❌ No questions could be generated from this file.",
        "summary_header": "📚 Summary\n\n",
        "stale": "⚠️ This button is outdated. Send /start or upload the file again.",
        "cooldown": "⚠️ Please wait {wait} seconds before sending another request.",
        "busy": "⚠️ You already have a request being processed. Please wait.",
        "q_meta": "Question {i}/{n}",
        "correct": "✅ Correct",
        "wrong": "❌ Wrong — you chose {chosen}",
        "answer": "Answer:",
        "explanation": "Explanation:",
        "complete": "🎓 Quiz Complete\n\nScore: {score}/{total} ({pct}%)",
        "review": "📋 Review Mistakes",
        "stats": "📊 My Stats",
        "new_quiz": "➕ New Quiz",
        "review_none": "✅ No wrong answers to review.",
        "review_done": "✅ Review complete.",
        "stopped": "🛑 Stopped.",
        "no_active": "No active quiz.",
        "quiz_in_progress": "A quiz is in progress. Use /stop to end it.",
        "send_file_hint": "Send a medical file to start.",
        "stats_text": "📊 Your Statistics\n\nQuizzes: {quizzes}\nCorrect: {correct}/{total}\nAccuracy: {accuracy}%\nBest score: {best}%",
        "error": "⚠️ Error: {err}",
        "summary_failed": "⚠️ Failed to generate summary. Please try again or use another file.",
        "osce_generating": "🎭 Creating your OSCE station...",
        "osce_ready": "🎭 OSCE Station Ready\n\nYou are the doctor. I am the patient. Start by greeting me and taking the history.",
        "osce_no_session": "No active OSCE session.",
        "osce_finished": "✅ OSCE finished. Here is your examiner feedback:",
    },
    "ar": {
        "choose_lang": "🌐 اختر لغتك",
        "welcome": "مرحباً بك في المساعد الطبي الذكي 🩺\n\nأرسل ملف PDF أو DOCX أو PPTX أو TXT.",
        "help": "أرسل ملفاً طبياً حتى 30MB. اختر اختبار أو ملخص. الأوامر: /start /stop /stats /review /language /support",
        "choose_mode": "📂 اختر الوضع:",
        "mode_quiz": "📝 اختبار",
        "mode_summary": "📚 ملخص",
        "mode_osce": "🎭 محاكي OSCE",
        "mode_past_papers": "📚 أسئلة سنوات سابقة",
        "mode_progress": "📊 تقدمي ونقاطي",
        "mode_leaderboard": "🏆 لوحة الصدارة",
        "mode_subscription": "💎 الاشتراك والدفع",
        "mode_quiz_prompt": "📝 تم اختيار الاختبار. أرسل الملف.",
        "mode_summary_prompt": "📚 تم اختيار الملخص. اختر نوع الملخص ثم أرسل الملف.",
        "mode_osce_prompt": "🎭 تم اختيار محاكي OSCE. أرسل ملف محاضرة طبية وسأصنع لك حالة أوسكي تفاعلية.",
        "summary_style_prompt": "📋 اختر نوع الملخص:",
        "style_disease": "🩺 ملف المرض",
        "style_highyield": "🧠 نقاط مهمة",
        "style_osce": "📋 OSCE",
        "file_received": "📄 تم استلام الملف — جاري استخراج النص...",
        "file_too_large": "⚠️ الملف أكبر من 30MB. أرسل ملفاً أصغر.",
        "file_unsupported": "⚠️ نوع الملف غير مدعوم. أرسل PDF, DOCX, PPTX أو TXT فقط.",
        "file_no_text": "⚠️ لم أستطع استخراج نص مقروء من الملف.",
        "file_empty": "⚠️ الملف فارغ أو يحتوي على صور فقط.",
        "file_ready_quiz": "✅ الملف جاهز. اختر مستوى الأسئلة:",
        "file_ready_summary": "✅ الملف جاهز. جاري تجهيز الملخص...",
        "cached_file": "♻️ الملف موجود مسبقاً في الكاش.",
        "level_basic": "🟢 Basic مباشر",
        "level_cases": "🩺 Cases حالات",
        "level_challenge": "🧠 Challenge تحدي",
        "basic_count_prompt": "كم سؤال Basic تريد؟",
        "generating_pool": "⏳ جاري تجهيز بنك الأسئلة. يحدث مرة واحدة فقط...",
        "generating_summary": "⏳ جاري تجهيز الملخص...",
        "not_enough": "⚠️ هذه المحاضرة فيها {n} سؤال فقط لهذا المستوى.",
        "quiz_ready": "✅ الاختبار جاهز: {n} سؤال.",
        "no_questions": "❌ لم أستطع توليد أسئلة من هذا الملف.",
        "summary_header": "📚 الملخص\n\n",
        "stale": "⚠️ هذا الزر قديم. أرسل /start أو ارفع الملف من جديد.",
        "cooldown": "⚠️ انتظر {wait} ثانية قبل إرسال طلب جديد.",
        "busy": "⚠️ لديك طلب قيد المعالجة. الرجاء الانتظار.",
        "q_meta": "السؤال {i}/{n}",
        "correct": "✅ صحيح",
        "wrong": "❌ خطأ — اخترت {chosen}",
        "answer": "الإجابة:",
        "explanation": "الشرح:",
        "complete": "🎓 انتهى الاختبار\n\nالنتيجة: {score}/{total} ({pct}%)",
        "review": "📋 مراجعة الأخطاء",
        "stats": "📊 إحصائياتي",
        "new_quiz": "➕ اختبار جديد",
        "review_none": "✅ لا توجد أخطاء للمراجعة.",
        "review_done": "✅ انتهت المراجعة.",
        "stopped": "🛑 تم الإيقاف.",
        "no_active": "لا يوجد اختبار نشط.",
        "quiz_in_progress": "يوجد اختبار نشط. استخدم /stop لإنهائه.",
        "send_file_hint": "أرسل ملفاً طبياً للبدء.",
        "stats_text": "📊 إحصائياتك\n\nالاختبارات: {quizzes}\nالصحيح: {correct}/{total}\nالدقة: {accuracy}%\nأفضل نتيجة: {best}%",
        "error": "⚠️ خطأ: {err}",
        "summary_failed": "⚠️ فشل توليد الملخص.",
        "osce_generating": "🎭 جاري إنشاء محطة OSCE...",
        "osce_ready": "🎭 محطة OSCE جاهزة\n\nأنت الطبيب، وأنا المريض. ابدأ بالسلام وأخذ التاريخ المرضي.",
        "osce_no_session": "لا توجد جلسة OSCE نشطة.",
        "osce_finished": "✅ انتهت محطة OSCE. هذا تقييم الممتحن:",
    },
}

def get_lang(user_data: dict) -> str:
    """Return language safely even if corrupted user_data is not a dict."""
    if not isinstance(user_data, dict):
        user_data = {}
    lang = user_data.get("lang", "ar")
    return lang if lang in SUPPORTED_LANGS else "ar"

def t(user_data: dict, key: str, **kwargs) -> str:
    """Safe translation helper."""
    if not isinstance(user_data, dict):
        user_data = {}
    lang = get_lang(user_data)
    base = TRANSLATIONS.get(lang, TRANSLATIONS.get("ar", {}))
    text = base.get(key, key) if isinstance(base, dict) else str(key)
    try:
        return text.format(**kwargs)
    except Exception:
        return text

# =============================================================================
# SQLITE CACHE & PROGRESS DB
# =============================================================================

_db_lock = asyncio.Lock()

def _db_connect() -> sqlite3.Connection:
    db_dir = os.path.dirname(DATABASE_PATH)
    if db_dir:
        os.makedirs(db_dir, exist_ok=True)
    conn = sqlite3.connect(DATABASE_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_db() -> None:
    conn = _db_connect()
    cur = conn.cursor()
    cur.execute("PRAGMA journal_mode=WAL")
    cur.execute("PRAGMA synchronous=NORMAL")
    cur.execute("CREATE TABLE IF NOT EXISTS files (file_hash TEXT PRIMARY KEY, file_name TEXT, file_size INTEGER, mime_type TEXT, extracted_text TEXT, created_at TEXT)")
    cur.execute("CREATE TABLE IF NOT EXISTS question_banks (file_hash TEXT NOT NULL, level TEXT NOT NULL, questions_json TEXT NOT NULL, created_at TEXT, PRIMARY KEY (file_hash, level))")
    cur.execute("CREATE TABLE IF NOT EXISTS summaries (file_hash TEXT NOT NULL, style TEXT NOT NULL, summary_text TEXT NOT NULL, created_at TEXT, PRIMARY KEY (file_hash, style))")
    cur.execute("CREATE TABLE IF NOT EXISTS users (user_id INTEGER PRIMARY KEY, full_name TEXT, username TEXT, lang TEXT DEFAULT 'ar', plan TEXT DEFAULT 'free', plan_expires_at REAL DEFAULT 0, created_at TEXT, last_seen TEXT)")
    cur.execute("CREATE TABLE IF NOT EXISTS user_progress (user_id INTEGER PRIMARY KEY, xp INTEGER DEFAULT 0, level INTEGER DEFAULT 1, streak INTEGER DEFAULT 0, last_activity_day TEXT, quizzes INTEGER DEFAULT 0, correct INTEGER DEFAULT 0, total INTEGER DEFAULT 0, osce_done INTEGER DEFAULT 0, summaries_done INTEGER DEFAULT 0, updated_at TEXT)")
    cur.execute("CREATE TABLE IF NOT EXISTS xp_events (id INTEGER PRIMARY KEY AUTOINCREMENT, user_id INTEGER, points INTEGER, reason TEXT, created_at TEXT)")
    cur.execute("CREATE TABLE IF NOT EXISTS payments (id INTEGER PRIMARY KEY AUTOINCREMENT, user_id INTEGER, status TEXT DEFAULT 'pending', note TEXT, created_at TEXT)")

    # الجداول الإضافية لميزة الـ Past Papers والتجميع العشوائي  
    cur.execute("CREATE TABLE IF NOT EXISTS past_papers (id INTEGER PRIMARY KEY AUTOINCREMENT, uni TEXT NOT NULL, subject TEXT NOT NULL, questions_json TEXT NOT NULL, created_at TEXT)")  
    cur.execute("CREATE TABLE IF NOT EXISTS admin_staging (user_id INTEGER NOT NULL, uni TEXT NOT NULL, subject TEXT NOT NULL, raw_content TEXT NOT NULL, created_at TEXT)")  
    cur.execute("CREATE TABLE IF NOT EXISTS user_usage_limits (user_id INTEGER NOT NULL, kind TEXT NOT NULL, last_used_at REAL NOT NULL, count INTEGER DEFAULT 0, window_start REAL DEFAULT 0, PRIMARY KEY (user_id, kind))")
    for alter_sql in [
        "ALTER TABLE users ADD COLUMN plan_expires_at REAL DEFAULT 0",
        "ALTER TABLE user_usage_limits ADD COLUMN count INTEGER DEFAULT 0",
        "ALTER TABLE user_usage_limits ADD COLUMN window_start REAL DEFAULT 0",
    ]:
        try:
            cur.execute(alter_sql)
        except Exception:
            pass
    cur.execute("CREATE TABLE IF NOT EXISTS past_progress (user_id INTEGER NOT NULL, uni TEXT NOT NULL, subject TEXT NOT NULL, attempts INTEGER DEFAULT 0, correct INTEGER DEFAULT 0, wrong INTEGER DEFAULT 0, wrong_json TEXT DEFAULT '[]', updated_at TEXT, PRIMARY KEY (user_id, uni, subject))")  
      
    conn.commit()  
    conn.close()

async def db_get_file(file_hash: str) -> Optional[dict]:
    def work():
        conn = _db_connect()
        row = conn.execute("SELECT * FROM files WHERE file_hash=?", (file_hash,)).fetchone()
        conn.close()
        return dict(row) if row else None
    return await asyncio.to_thread(work)

async def db_save_file(file_hash: str, file_name: str, file_size: int, mime_type: str, extracted_text: str) -> None:
    async with _db_lock:
        def work():
            conn = _db_connect()
            conn.execute(
                "INSERT OR REPLACE INTO files(file_hash, file_name, file_size, mime_type, extracted_text, created_at) VALUES(?,?,?,?,?,?)",
                (file_hash, file_name, file_size, mime_type, extracted_text, datetime.utcnow().isoformat()),
            )
            conn.commit()
            conn.close()
        await asyncio.to_thread(work)

async def db_get_questions(file_hash: str, level: str) -> Optional[List[dict]]:
    def work():
        conn = _db_connect()
        row = conn.execute("SELECT questions_json FROM question_banks WHERE file_hash=? AND level=?", (file_hash, level)).fetchone()
        conn.close()
        if not row: return None
        try: return json.loads(row["questions_json"])
        except Exception: return None
    return await asyncio.to_thread(work)

async def db_save_questions(file_hash: str, level: str, questions: List[dict]) -> None:
    async with _db_lock:
        def work():
            conn = _db_connect()
            conn.execute(
                "INSERT OR REPLACE INTO question_banks(file_hash, level, questions_json, created_at) VALUES(?,?,?,?)",
                (file_hash, level, json.dumps(questions, ensure_ascii=False), datetime.utcnow().isoformat()),
            )
            conn.commit()
            conn.close()
        await asyncio.to_thread(work)

async def db_get_summary(file_hash: str, style: str) -> Optional[str]:
    def work():
        conn = _db_connect()
        row = conn.execute("SELECT summary_text FROM summaries WHERE file_hash=? AND style=?", (file_hash, style)).fetchone()
        conn.close()
        return row["summary_text"] if row else None
    return await asyncio.to_thread(work)

async def db_save_summary(file_hash: str, style: str, summary: str) -> None:
    async with _db_lock:
        def work():
            conn = _db_connect()
            conn.execute(
                "INSERT OR REPLACE INTO summaries(file_hash, style, summary_text, created_at) VALUES(?,?,?,?)",
                (file_hash, style, summary, datetime.utcnow().isoformat()),
            )
            conn.commit()
            conn.close()
        await asyncio.to_thread(work)

def calc_level(xp: int) -> int: return max(1, int(xp // 250) + 1)
def today_local_key() -> str: return datetime.utcnow().strftime("%Y-%m-%d")

async def db_register_user(user) -> None:
    if not user: return
    async with _db_lock:
        def work():
            now = datetime.utcnow().isoformat()
            conn = _db_connect()
            conn.execute(
                "INSERT INTO users(user_id, full_name, username, created_at, last_seen) VALUES(?,?,?,?,?) "
                "ON CONFLICT(user_id) DO UPDATE SET full_name=excluded.full_name, username=excluded.username, last_seen=excluded.last_seen",
                (user.id, user.full_name or "", user.username or "", now, now),
            )
            conn.execute("INSERT OR IGNORE INTO user_progress(user_id, updated_at) VALUES(?, ?)", (user.id, now))
            conn.commit()
            conn.close()
        await asyncio.to_thread(work)

async def db_get_user_plan(user_id: int) -> str:
    if not user_id:
        return "free"
    def work():
        conn = _db_connect()
        row = conn.execute("SELECT plan, plan_expires_at FROM users WHERE user_id=?", (int(user_id),)).fetchone()
        if not row:
            conn.close()
            return "free"
        plan = (row["plan"] or "free").strip().lower()
        expires = float(row["plan_expires_at"] or 0)
        if plan not in {"free", "premium", "premium_plus", "admin"}:
            plan = "free"
        if plan not in {"free", "admin"} and expires and expires < time.time():
            conn.execute("UPDATE users SET plan='free', plan_expires_at=0 WHERE user_id=?", (int(user_id),))
            conn.commit()
            plan = "free"
        conn.close()
        return plan
    try:
        return str(await asyncio.to_thread(work)).strip().lower() or "free"
    except Exception:
        return "free"

async def db_get_user_plan_info(user_id: int) -> dict:
    def work():
        conn = _db_connect()
        row = conn.execute("SELECT plan, plan_expires_at FROM users WHERE user_id=?", (int(user_id),)).fetchone()
        conn.close()
        if not row:
            return {"plan": "free", "plan_expires_at": 0}
        return {"plan": (row["plan"] or "free"), "plan_expires_at": float(row["plan_expires_at"] or 0)}
    try:
        info = await asyncio.to_thread(work)
        plan = await db_get_user_plan(user_id)
        info["plan"] = plan
        return info
    except Exception:
        return {"plan": "free", "plan_expires_at": 0}

async def db_set_user_plan(user_id: int, plan: str, days: Optional[int] = None) -> None:
    plan = (plan or "free").strip().lower()
    if plan not in {"free", "premium", "premium_plus", "admin"}:
        plan = "free"
    expires = 0.0
    if plan in {"premium", "premium_plus"}:
        expires = time.time() + max(1, int(days or 30)) * 86400
    async with _db_lock:
        def work():
            now = datetime.utcnow().isoformat()
            conn = _db_connect()
            conn.execute(
                "INSERT INTO users(user_id, plan, plan_expires_at, created_at, last_seen) VALUES(?,?,?,?,?) "
                "ON CONFLICT(user_id) DO UPDATE SET plan=excluded.plan, plan_expires_at=excluded.plan_expires_at, last_seen=excluded.last_seen",
                (int(user_id), plan, expires, now, now),
            )
            conn.commit(); conn.close()
        await asyncio.to_thread(work)

def set_ai_user_id(user_id: Optional[int]) -> None:
    try: _AI_USER_ID.set(int(user_id or 0))
    except Exception: _AI_USER_ID.set(0)

async def _paid_fallback_allowed(user_id: Optional[int] = None) -> bool:
    uid = int(user_id or _AI_USER_ID.get() or 0)
    if is_admin_user(uid):
        return ALLOW_PAID_FALLBACK_FOR_ADMIN
    if PREMIUM_ENABLED and uid:
        plan = await db_get_user_plan(uid)
        if plan in {"premium", "premium_plus", "admin"}:
            return ALLOW_PAID_FALLBACK_FOR_PREMIUM
    return ALLOW_PAID_FALLBACK_FOR_USERS

async def db_add_xp(user_id: int, points: int, reason: str, correct: int = 0, total: int = 0, quiz_done: bool = False, osce_done: bool = False, summary_done: bool = False) -> dict:
    async with _db_lock:
        def work():
            now = datetime.utcnow().isoformat()
            today = today_local_key()
            conn = _db_connect()
            row = conn.execute("SELECT * FROM user_progress WHERE user_id=?", (user_id,)).fetchone()
            if not row:
                conn.execute("INSERT OR IGNORE INTO user_progress(user_id, updated_at) VALUES(?, ?)", (user_id, now))
                row = conn.execute("SELECT * FROM user_progress WHERE user_id=?", (user_id,)).fetchone()

            old_xp = int(row["xp"] or 0)  
            old_level = int(row["level"] or 1)  
            last_day = row["last_activity_day"]  
            streak = int(row["streak"] or 0)  

            if last_day != today:  
                try:  
                    prev = datetime.strptime(last_day, "%Y-%m-%d") if last_day else None  
                    if prev and (datetime.strptime(today, "%Y-%m-%d") - prev).days == 1:  
                        streak += 1  
                    else: streak = 1  
                except Exception: streak = 1  

            new_xp = old_xp + max(0, int(points))  
            new_level = calc_level(new_xp)  
            conn.execute(  
                "UPDATE user_progress SET xp=?, level=?, streak=?, last_activity_day=?, quizzes=quizzes+?, "  
                "correct=correct+?, total=total+?, osce_done=osce_done+?, summaries_done=summaries_done+?, updated_at=? WHERE user_id=?",  
                (new_xp, new_level, streak, today, 1 if quiz_done else 0, int(correct), int(total), 1 if osce_done else 0, 1 if summary_done else 0, now, user_id),  
            )  
            conn.execute("INSERT INTO xp_events(user_id, points, reason, created_at) VALUES(?,?,?,?)", (user_id, int(points), reason, now))  
            conn.commit()  
            out = dict(conn.execute("SELECT * FROM user_progress WHERE user_id=?", (user_id,)).fetchone())  
            conn.close()  
            out["leveled_up"] = new_level > old_level  
            return out  
        return await asyncio.to_thread(work)

async def db_get_progress(user_id: int) -> dict:
    def work():
        conn = _db_connect()
        row = conn.execute("SELECT * FROM user_progress WHERE user_id=?", (user_id,)).fetchone()
        conn.close()
        return dict(row) if row else {"xp": 0, "level": 1, "streak": 0, "quizzes": 0, "correct": 0, "total": 0, "osce_done": 0, "summaries_done": 0}
    return await asyncio.to_thread(work)

async def db_get_leaderboard(limit: int = 10) -> list:
    def work():
        conn = _db_connect()
        rows = conn.execute("SELECT u.full_name, u.username, p.xp, p.level, p.streak FROM user_progress p LEFT JOIN users u ON u.user_id=p.user_id ORDER BY p.xp DESC, p.level DESC LIMIT ?", (limit,)).fetchall()
        conn.close()
        return [dict(r) for r in rows]
    return await asyncio.to_thread(work)

async def db_create_payment_request(user_id: int, note: str = "") -> None:
    async with _db_lock:
        def work():
            conn = _db_connect()
            conn.execute("INSERT INTO payments(user_id, status, note, created_at) VALUES(?,?,?,?)", (user_id, "pending", note, datetime.utcnow().isoformat()))
            conn.commit()
            conn.close()
        await asyncio.to_thread(work)

def ai_queue_status_text() -> str:
    return f"⚙️ AI Load: active {_ACTIVE_AI_JOBS}/{MAX_CONCURRENT_GEMINI} | waiting {_WAITING_AI_JOBS}"

async def cleanup_old_cache(days: int = 30) -> None:
    cutoff = (datetime.utcnow() - timedelta(days=days)).isoformat()
    async with _db_lock:
        def work():
            conn = _db_connect()
            cur = conn.cursor()
            old_hashes = [r["file_hash"] for r in cur.execute("SELECT file_hash FROM files WHERE created_at < ?", (cutoff,)).fetchall()]
            for h in old_hashes:
                cur.execute("DELETE FROM question_banks WHERE file_hash=?", (h,))
                cur.execute("DELETE FROM summaries WHERE file_hash=?", (h,))
                cur.execute("DELETE FROM files WHERE file_hash=?", (h,))
            cur.execute("DELETE FROM question_banks WHERE created_at < ?", (cutoff,))
            cur.execute("DELETE FROM summaries WHERE created_at < ?", (cutoff,))
            conn.commit()
            conn.close()
            return len(old_hashes)
        deleted = await asyncio.to_thread(work)
        logger.info("Cleanup done. Deleted old files: %s", deleted)

# =============================================================================
# HELPERS & USAGE LIMITS
# =============================================================================

user_last_request: Dict[int, float] = {}

def sha256_bytes(data: bytes) -> str: return hashlib.sha256(data).hexdigest()

def escape_md(text: object) -> str:
    value = "" if text is None else str(text)
    for ch in ("*", "_", "`", "["): value = value.replace(ch, f"\{ch}")
    return value

def progress_bar(current: int, total: int, width: int = 12) -> str:
    if total <= 0: return "░" * width + " 0%"
    filled = round(current / total * width); pct = round(current / total * 100)
    return "█" * filled + "░" * (width - filled) + f" {pct}%"

async def cooldown_guard(update: Update, user_data: dict, user_id: int) -> bool:
    now = time.time()
    last = user_last_request.get(user_id, 0)
    if now - last < COOLDOWN_SECONDS:
        wait = max(1, int(COOLDOWN_SECONDS - (now - last) + 0.999))
        await update.effective_message.reply_text(t(user_data, "cooldown", wait=wait))
        return True
    user_last_request[user_id] = now
    return False

def get_stats(user_data: dict) -> dict:
    if "stats" not in user_data: user_data["stats"] = {"quizzes": 0, "correct": 0, "total": 0, "best": 0}
    return user_data["stats"]

def update_stats(user_data: dict, score: int, total: int) -> None:
    stats = get_stats(user_data)
    pct = round(score / total * 100) if total else 0
    stats["quizzes"] += 1; stats["correct"] += score; stats["total"] += total
    if pct > stats["best"]: stats["best"] = pct

def today_key() -> str: return datetime.utcnow().strftime("%Y-%m-%d")

def _format_wait_ar(seconds: int) -> str:
    seconds = max(1, int(seconds))
    days, rem = divmod(seconds, 86400)
    hours, rem = divmod(rem, 3600)
    minutes = max(1, rem // 60) if not days and not hours else rem // 60
    parts = []
    if days: parts.append(f"{days} يوم")
    if hours: parts.append(f"{hours} ساعة")
    if minutes and not days: parts.append(f"{minutes} دقيقة")
    return " و ".join(parts) or "دقيقة"

def usage_limit_for(kind: str) -> int:
    # Kept for backward compatibility. Economy mode uses time windows instead of counters.
    return 1

def is_admin_user(user_id: Optional[int]) -> bool: return bool(user_id and user_id in ADMIN_USER_IDS)

def _usage_rule(plan: str, kind: str) -> dict:
    plan = (plan or "free").strip().lower()
    if plan == "admin":
        return {"limit": None, "window": 0}
    return PLAN_USAGE_LIMITS.get(plan, PLAN_USAGE_LIMITS["free"]).get(kind, {"limit": 1, "window": 24*3600})

async def check_daily_limit(update: Update, context: ContextTypes.DEFAULT_TYPE, kind: str) -> bool:
    """Economy limits. Cached results should call this only when AI generation is needed."""
    user = update.effective_user
    user_id = user.id if user else 0
    if is_admin_user(user_id):
        return True
    if not ECONOMY_MODE:
        return True
    plan = await db_get_user_plan(user_id)
    rule = _usage_rule(plan, kind)
    limit = rule.get("limit")
    window = int(rule.get("window") or 0)
    if not limit or window <= 0:
        return True
    now = time.time()

    def work():
        conn = _db_connect()
        row = conn.execute("SELECT last_used_at, count, window_start FROM user_usage_limits WHERE user_id=? AND kind=?", (int(user_id), kind)).fetchone()
        conn.close()
        if not row:
            return {"last": 0.0, "count": 0, "start": 0.0}
        return {"last": float(row["last_used_at"] or 0), "count": int(row["count"] or 0), "start": float(row["window_start"] or 0)}

    usage = await asyncio.to_thread(work)
    start = usage["start"] or usage["last"] or 0.0
    if not start or now - start >= window:
        return True
    if usage["count"] < int(limit):
        return True
    remaining = int(window - (now - start))
    labels = {"summary": "الملخص", "basic": "توليد أسئلة Basic", "cases": "Cases", "challenge": "Challenge", "osce": "OSCE"}
    await update.effective_message.reply_text(
        f"⚠️ وصلت لحد خطة {PLAN_LABELS.get(plan, plan)} في: {labels.get(kind, kind)}.\n"
        f"يرجع متاح بعد: {_format_wait_ar(remaining)}.\n\n"
        "ملاحظة: إذا كان الملف محفوظ مسبقاً في مخزون البوت فلن يستهلك حدك ولا API."
    )
    return False

async def increment_daily_usage(user_data: dict, kind: str, user_id: Optional[int] = None) -> None:
    if not user_id or is_admin_user(user_id) or not ECONOMY_MODE:
        return
    plan = await db_get_user_plan(int(user_id))
    rule = _usage_rule(plan, kind)
    limit = rule.get("limit")
    window = int(rule.get("window") or 0)
    if not limit or window <= 0:
        return
    now = time.time()
    async with _db_lock:
        def work():
            conn = _db_connect()
            row = conn.execute("SELECT count, window_start FROM user_usage_limits WHERE user_id=? AND kind=?", (int(user_id), kind)).fetchone()
            if row:
                start = float(row["window_start"] or 0)
                count = int(row["count"] or 0)
                if not start or now - start >= window:
                    start = now
                    count = 1
                else:
                    count += 1
                conn.execute(
                    "INSERT OR REPLACE INTO user_usage_limits(user_id, kind, last_used_at, count, window_start) VALUES(?,?,?,?,?)",
                    (int(user_id), kind, now, count, start),
                )
            else:
                conn.execute(
                    "INSERT OR REPLACE INTO user_usage_limits(user_id, kind, last_used_at, count, window_start) VALUES(?,?,?,?,?)",
                    (int(user_id), kind, now, 1, now),
                )
            conn.commit(); conn.close()
        await asyncio.to_thread(work)

# =============================================================================
# AI ENGINE — Google Gemini first, OpenRouter fallback
# =============================================================================

def _google_model_name() -> str:
    model = (GEMINI_DIRECT_MODEL or "gemini-2.5-flash-lite").strip()
    if model.startswith("models/"):
        model = model.split("/", 1)[1]
    return model

async def _call_google_text(prompt: str, api_key: str, temperature: float, max_output_tokens: int) -> str:
    model = _google_model_name()
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
    payload = {
        "contents": [{"role": "user", "parts": [{"text": prompt}]}],
        "generationConfig": {
            "temperature": temperature,
            "maxOutputTokens": max_output_tokens,
        },
    }
    headers = {"x-goog-api-key": api_key, "Content-Type": "application/json"}
    async with aiohttp.ClientSession() as session:
        async with session.post(url, json=payload, headers=headers, timeout=120) as response:
            data_text = await response.text()
            if response.status != 200:
                raise Exception(f"Google HTTP {response.status}: {data_text[:300]}")
            data = json.loads(data_text)
            candidates = data.get("candidates") or []
            if not candidates:
                raise Exception(f"Google empty candidates: {data_text[:300]}")
            parts = candidates[0].get("content", {}).get("parts", [])
            text = "".join(str(part.get("text", "")) for part in parts).strip()
            if not text:
                raise Exception(f"Google empty text: {data_text[:300]}")
            return text

async def _call_google_vision(image_bytes: bytes, api_key: str) -> str:
    model = _google_model_name()
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
    base64_image = base64.b64encode(image_bytes).decode("utf-8")
    payload = {
        "contents": [{
            "role": "user",
            "parts": [
                {"text": "Extract ALL visible lecture content from this slide/page image clearly and completely. Preserve titles, headings, bullet points, tables, numbers, weeks, percentages, and medical terms. Do NOT summarize. Do NOT look only for exam questions. If the content is Arabic, extract it in Arabic. If it is English, extract it in English. Return clean readable lecture text only."},
                {"inlineData": {"mimeType": "image/jpeg", "data": base64_image}},
            ],
        }],
        "generationConfig": {"temperature": 0.1, "maxOutputTokens": 4096},
    }
    headers = {"x-goog-api-key": api_key, "Content-Type": "application/json"}
    async with aiohttp.ClientSession() as session:
        async with session.post(url, json=payload, headers=headers, timeout=120) as response:
            data_text = await response.text()
            if response.status != 200:
                raise Exception(f"Google Vision HTTP {response.status}: {data_text[:300]}")
            data = json.loads(data_text)
            candidates = data.get("candidates") or []
            if not candidates:
                raise Exception(f"Google Vision empty candidates: {data_text[:300]}")
            parts = candidates[0].get("content", {}).get("parts", [])
            return "".join(str(part.get("text", "")) for part in parts).strip()



def guess_mime_type(file_name: str, fallback: str = "application/octet-stream") -> str:
    guessed, _ = mimetypes.guess_type(file_name or "")
    return guessed or fallback

async def _call_google_inline_file(raw_bytes: bytes, mime_type: str, file_name: str, prompt: str, api_key: str, max_output_tokens: int = 8192) -> str:
    """Admin-only smart extraction for lecture files when normal text extraction fails."""
    model = _google_model_name()
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
    mime_type = (mime_type or guess_mime_type(file_name)).strip() or "application/octet-stream"
    b64 = base64.b64encode(raw_bytes).decode("utf-8")
    payload = {
        "contents": [{
            "role": "user",
            "parts": [
                {"text": prompt},
                {"inlineData": {"mimeType": mime_type, "data": b64}},
            ],
        }],
        "generationConfig": {"temperature": 0.1, "maxOutputTokens": max_output_tokens},
    }
    headers = {"x-goog-api-key": api_key, "Content-Type": "application/json"}
    async with aiohttp.ClientSession() as session:
        async with session.post(url, json=payload, headers=headers, timeout=180) as response:
            data_text = await response.text()
            if response.status != 200:
                raise Exception(f"Google file vision HTTP {response.status}: {data_text[:300]}")
            data = json.loads(data_text)
            candidates = data.get("candidates") or []
            if not candidates:
                raise Exception(f"Google file vision empty candidates: {data_text[:300]}")
            parts = candidates[0].get("content", {}).get("parts", [])
            return "".join(str(part.get("text", "")) for part in parts).strip()

async def call_gemini_admin_document_extract(raw_bytes: bytes, file_name: str, mime_type: str = "") -> str:
    """Admin-only: read lecture documents/slides/images with Gemini Vision when local extraction fails."""
    prompt = (
        "You are extracting text from a medical lecture file for caching in a Telegram study bot. "
        "Read the document/slides/images carefully and extract the lecture content in order. "
        "If it contains Arabic, keep Arabic. If it contains English, keep English. "
        "Preserve headings, bullet points, tables in plain text, important terms, and slide titles. "
        "Do not summarize. Do not answer questions. Return only clear extracted lecture text."
    )
    if GOOGLE_API_KEYS:
        last_error = None
        for idx, api_key in enumerate(GOOGLE_API_KEYS, 1):
            try:
                text = await _call_google_inline_file(raw_bytes, mime_type or guess_mime_type(file_name), file_name, prompt, api_key, max_output_tokens=8192)
                if text and len(text.strip()) >= 50:
                    logger.info("Admin smart extraction: Google direct key=%s/%s file=%s", idx, len(GOOGLE_API_KEYS), file_name)
                    return text.strip()
            except Exception as e:
                last_error = e
                logger.warning("Admin smart extraction key %s/%s failed for %s: %s", idx, len(GOOGLE_API_KEYS), file_name, str(e)[:220])
        logger.warning("Admin smart extraction failed on all Google keys for %s. Last error: %s", file_name, last_error)
    return ""



# =============================================================================
# ADMIN SMART DOCUMENT REPAIR / CONVERSION
# =============================================================================

LECTURE_IMAGE_PROMPT = (
    "Extract the full readable medical lecture content from this slide/page image. "
    "Keep the original order, headings, bullets, tables, numbers, and abbreviations. "
    "If the text is Arabic, keep Arabic. If English, keep English. "
    "Do not summarize. Return only the extracted lecture text."
)

def _has_command(cmd: str) -> bool:
    return bool(shutil.which(cmd))

def _is_legacy_office_file(file_name: str) -> bool:
    name = (file_name or "").lower()
    return name.endswith((".ppt", ".doc", ".pps", ".pot")) and not name.endswith((".pptx", ".docx"))

def _looks_like_office_file(file_name: str) -> bool:
    name = (file_name or "").lower()
    return name.endswith((".ppt", ".pptx", ".pps", ".doc", ".docx", ".odp", ".odt"))

def _convert_office_to_pdf_bytes(raw_bytes: bytes, file_name: str) -> bytes:
    """Convert old/new Office documents to PDF using LibreOffice headless when available."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        logger.warning("LibreOffice/soffice not installed; cannot convert %s", file_name)
        return b""
    safe_name = re.sub(r"[^A-Za-z0-9_.-]+", "_", os.path.basename(file_name or "lecture")) or "lecture"
    if "." not in safe_name:
        safe_name += ".ppt"
    with tempfile.TemporaryDirectory() as td:
        in_path = os.path.join(td, safe_name)
        out_dir = os.path.join(td, "out")
        os.makedirs(out_dir, exist_ok=True)
        with open(in_path, "wb") as f:
            f.write(raw_bytes)
        cmd = [soffice, "--headless", "--nologo", "--nofirststartwizard", "--convert-to", "pdf", "--outdir", out_dir, in_path]
        try:
            proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=OFFICE_CONVERT_TIMEOUT)
            if proc.returncode != 0:
                logger.warning("LibreOffice conversion failed for %s: %s", file_name, (proc.stderr or proc.stdout).decode(errors="ignore")[:500])
                return b""
            pdfs = sorted(Path(out_dir).glob("*.pdf")) or sorted(Path(td).glob("*.pdf"))
            if not pdfs:
                logger.warning("LibreOffice conversion produced no PDF for %s", file_name)
                return b""
            return pdfs[0].read_bytes()
        except Exception as e:
            logger.warning("LibreOffice conversion exception for %s: %s", file_name, e)
            return b""

def _render_pdf_pages_to_png_bytes(pdf_bytes: bytes, max_pages: int = ADMIN_VISION_MAX_PAGES) -> List[bytes]:
    """Render PDF pages to PNG using poppler pdftoppm for admin-only Vision fallback."""
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        logger.warning("poppler-utils/pdftoppm not installed; cannot render scanned PDF pages")
        return []
    with tempfile.TemporaryDirectory() as td:
        pdf_path = os.path.join(td, "input.pdf")
        prefix = os.path.join(td, "page")
        with open(pdf_path, "wb") as f:
            f.write(pdf_bytes)
        cmd = [pdftoppm, "-png", "-r", "130", "-f", "1", "-l", str(max_pages), pdf_path, prefix]
        try:
            proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=120)
            if proc.returncode != 0:
                logger.warning("pdftoppm failed: %s", (proc.stderr or proc.stdout).decode(errors="ignore")[:500])
                return []
            imgs = sorted(Path(td).glob("page-*.png"))
            return [x.read_bytes() for x in imgs[:max_pages]]
        except Exception as e:
            logger.warning("pdftoppm exception: %s", e)
            return []

async def _call_google_image_with_prompt(image_bytes: bytes, prompt: str, api_key: str, max_output_tokens: int = 4096) -> str:
    model = _google_model_name()
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
    base64_image = base64.b64encode(image_bytes).decode("utf-8")
    payload = {
        "contents": [{"role": "user", "parts": [{"text": prompt}, {"inlineData": {"mimeType": "image/png", "data": base64_image}}]}],
        "generationConfig": {"temperature": 0.1, "maxOutputTokens": max_output_tokens},
    }
    headers = {"x-goog-api-key": api_key, "Content-Type": "application/json"}
    async with aiohttp.ClientSession() as session:
        async with session.post(url, json=payload, headers=headers, timeout=120) as response:
            data_text = await response.text()
            if response.status != 200:
                raise Exception(f"Google image HTTP {response.status}: {data_text[:300]}")
            data = json.loads(data_text)
            candidates = data.get("candidates") or []
            if not candidates:
                raise Exception(f"Google image empty candidates: {data_text[:300]}")
            parts = candidates[0].get("content", {}).get("parts", [])
            return "".join(str(part.get("text", "")) for part in parts).strip()

async def _admin_vision_extract_images(image_list: List[bytes], file_name: str) -> str:
    """Admin-only Vision extraction from rendered slide/page images.
    Tries Google keys first, then OpenRouter Vision fallback for admin preload only.
    """
    texts = []
    for page_idx, img in enumerate(image_list[:ADMIN_VISION_MAX_PAGES], 1):
        page_text = ""
        last_error = None

        # 1) Google direct keys
        for key_idx, api_key in enumerate(GOOGLE_API_KEYS, 1):
            try:
                page_text = await _call_google_image_with_prompt(img, LECTURE_IMAGE_PROMPT, api_key, max_output_tokens=4096)
                if page_text and len(page_text.strip()) >= 20:
                    logger.info("Admin page vision OK via Google file=%s page=%s key=%s/%s", file_name, page_idx, key_idx, len(GOOGLE_API_KEYS))
                    break
            except Exception as e:
                last_error = e
                logger.warning("Admin page vision Google failed file=%s page=%s key=%s/%s: %s", file_name, page_idx, key_idx, len(GOOGLE_API_KEYS), str(e)[:180])

        # 2) OpenRouter fallback for admin only when Google quota/failure happens
        if not page_text and OPENROUTER_API_KEY:
            try:
                page_text = await _call_openrouter_vision(img)
                if page_text and len(page_text.strip()) >= 20:
                    logger.info("Admin page vision OK via OpenRouter file=%s page=%s", file_name, page_idx)
            except Exception as e:
                last_error = e
                logger.warning("Admin page vision OpenRouter failed file=%s page=%s: %s", file_name, page_idx, str(e)[:180])

        if page_text:
            texts.append(f"\n--- Page/Slide {page_idx} ---\n{page_text.strip()}")
        elif last_error:
            logger.warning("Admin page vision all providers failed file=%s page=%s: %s", file_name, page_idx, last_error)
    return "\n".join(texts).strip()

async def admin_repair_extract_via_conversion(raw_bytes: bytes, file_name: str, file_type: str) -> Tuple[str, str]:
    """Admin-only robust fallback: Office -> PDF -> text, then rendered page images -> Vision."""
    name = (file_name or "").lower()
    pdf_bytes = b""
    if file_type == "pdf" or name.endswith(".pdf"):
        pdf_bytes = raw_bytes
    elif _looks_like_office_file(file_name):
        pdf_bytes = _convert_office_to_pdf_bytes(raw_bytes, file_name)
        if pdf_bytes:
            text = extract_text("pdf", pdf_bytes)
            if text and len(text.strip()) >= 100:
                return text.strip(), "libreoffice_pdf_text"
    if pdf_bytes:
        images = _render_pdf_pages_to_png_bytes(pdf_bytes, ADMIN_VISION_MAX_PAGES)
        if images:
            text = await _admin_vision_extract_images(images, file_name)
            if text and len(text.strip()) >= 50:
                return text.strip(), "rendered_pdf_vision"
    return "", ""

async def _call_openrouter_text(prompt: str, temperature: float, max_output_tokens: int, model: Optional[str] = None, api_key: Optional[str] = None) -> str:
    key = api_key or (OPENROUTER_PAID_KEYS[0] if OPENROUTER_PAID_KEYS else OPENROUTER_API_KEY)
    if not key:
        raise Exception("OpenRouter API key is missing.")
    model_name = model or GEMINI_MODEL_NAME
    payload = {"model": model_name, "messages": [{"role": "user", "content": prompt}], "temperature": temperature, "max_tokens": max_output_tokens}
    headers = {"Authorization": f"Bearer {key}", "Content-Type": "application/json"}
    async with aiohttp.ClientSession() as session:
        async with session.post("https://openrouter.ai/api/v1/chat/completions", json=payload, headers=headers, timeout=120) as response:
            if response.status == 200:
                return (await response.json())["choices"][0]["message"]["content"].strip()
            err_text = await response.text()
            raise Exception(f"OpenRouter HTTP {response.status}: {err_text[:300]}")

async def _call_openrouter_free_text(prompt: str, temperature: float, max_output_tokens: int) -> str:
    if not OPENROUTER_FREE_KEYS:
        raise Exception("OPENROUTER_FREE_KEYS/OPENROUTER_API_KEY is missing.")
    last_error = None
    for idx, key in enumerate(OPENROUTER_FREE_KEYS, 1):
        try:
            result = await _call_openrouter_text(prompt, temperature, max_output_tokens, model=OPENROUTER_FREE_MODEL, api_key=key)
            logger.info("AI provider: OpenRouter FREE key=%s/%s model=%s", idx, len(OPENROUTER_FREE_KEYS), OPENROUTER_FREE_MODEL)
            return result
        except Exception as e:
            last_error = e
            logger.warning("OpenRouter FREE key %s/%s failed: %s", idx, len(OPENROUTER_FREE_KEYS), str(e)[:200])
    raise Exception(f"All OpenRouter free keys failed: {last_error}")

async def _call_groq_text(prompt: str, temperature: float, max_output_tokens: int) -> str:
    if not GROQ_API_KEYS:
        raise Exception("GROQ_API_KEYS is missing.")
    last_error = None
    payload = {"model": GROQ_MODEL, "messages": [{"role": "user", "content": prompt}], "temperature": temperature, "max_tokens": max_output_tokens}
    for idx, key in enumerate(GROQ_API_KEYS, 1):
        headers = {"Authorization": f"Bearer {key}", "Content-Type": "application/json"}
        try:
            async with aiohttp.ClientSession() as session:
                async with session.post("https://api.groq.com/openai/v1/chat/completions", json=payload, headers=headers, timeout=120) as response:
                    if response.status == 200:
                        data = await response.json()
                        logger.info("AI provider: Groq FREE key=%s/%s model=%s", idx, len(GROQ_API_KEYS), GROQ_MODEL)
                        return data["choices"][0]["message"]["content"].strip()
                    err_text = await response.text()
                    raise Exception(f"Groq HTTP {response.status}: {err_text[:300]}")
        except Exception as e:
            last_error = e
            logger.warning("Groq key %s/%s failed: %s", idx, len(GROQ_API_KEYS), str(e)[:200])
    raise Exception(f"All Groq keys failed: {last_error}")

async def _call_cloudflare_text(prompt: str, temperature: float, max_output_tokens: int) -> str:
    if not CLOUDFLARE_ACCOUNT_ID or not CLOUDFLARE_API_TOKENS:
        raise Exception("Cloudflare Workers AI credentials are missing.")
    model = CLOUDFLARE_MODEL
    payload = {
        "messages": [
            {"role": "system", "content": "You are a precise medical study assistant. Return clean useful output only."},
            {"role": "user", "content": prompt},
        ],
        "temperature": temperature,
        "max_tokens": max_output_tokens,
    }
    last_error = None
    url = f"https://api.cloudflare.com/client/v4/accounts/{CLOUDFLARE_ACCOUNT_ID}/ai/run/{model}"
    for idx, token in enumerate(CLOUDFLARE_API_TOKENS, 1):
        headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(url, json=payload, headers=headers, timeout=120) as response:
                    data_text = await response.text()
                    if response.status == 200:
                        try:
                            data = json.loads(data_text)
                        except Exception:
                            data = {}
                        result = data.get("result") if isinstance(data, dict) else None
                        text = ""
                        if isinstance(result, dict):
                            text = result.get("response") or result.get("text") or result.get("generated_text") or ""
                            if not text and isinstance(result.get("output"), list):
                                text = "\n".join(str(x) for x in result.get("output") if x)
                        elif isinstance(result, str):
                            text = result
                        if not text and isinstance(data, dict):
                            text = data.get("response") or data.get("text") or ""
                        text = str(text).strip()
                        if text:
                            logger.info("AI provider: Cloudflare FREE key=%s/%s model=%s", idx, len(CLOUDFLARE_API_TOKENS), model)
                            return text
                        raise Exception(f"Cloudflare empty response: {data_text[:300]}")
                    raise Exception(f"Cloudflare HTTP {response.status}: {data_text[:300]}")
        except Exception as e:
            last_error = e
            logger.warning("Cloudflare key %s/%s failed: %s", idx, len(CLOUDFLARE_API_TOKENS), str(e)[:200])
    raise Exception(f"All Cloudflare keys failed: {last_error}")

async def _call_openrouter_vision(image_bytes: bytes, model: Optional[str] = None, api_key: Optional[str] = None) -> str:
    key = api_key or (OPENROUTER_PAID_KEYS[0] if OPENROUTER_PAID_KEYS else OPENROUTER_API_KEY)
    if not key:
        raise Exception("OpenRouter API key is missing for vision.")
    base64_image = base64.b64encode(image_bytes).decode('utf-8')
    payload = {
        "model": model or GEMINI_MODEL_NAME,
        "messages": [{"role": "user", "content": [{"type": "text", "text": "Extract ALL visible lecture content from this slide/page image clearly and completely. Preserve titles, headings, bullet points, tables, numbers, weeks, percentages, and medical terms. Do NOT summarize. Do NOT look only for exam questions. If the content is Arabic, extract it in Arabic. If it is English, extract it in English. Return clean readable lecture text only."}, {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}]}],
        "temperature": 0.1,
    }
    headers = {"Authorization": f"Bearer {key}", "Content-Type": "application/json"}
    async with aiohttp.ClientSession() as session:
        async with session.post("https://openrouter.ai/api/v1/chat/completions", json=payload, headers=headers, timeout=120) as response:
            if response.status == 200:
                data = await response.json()
                return data["choices"][0]["message"]["content"].strip()
            err_text = await response.text()
            raise Exception(f"OpenRouter Vision HTTP {response.status}: {err_text[:300]}")

async def _call_openrouter_free_vision(image_bytes: bytes) -> str:
    if not OPENROUTER_FREE_KEYS:
        raise Exception("OPENROUTER_FREE_KEYS/OPENROUTER_API_KEY is missing for free vision.")
    last_error = None
    for idx, key in enumerate(OPENROUTER_FREE_KEYS, 1):
        try:
            text = await _call_openrouter_vision(image_bytes, model=OPENROUTER_FREE_MODEL, api_key=key)
            logger.info("Vision provider: OpenRouter FREE key=%s/%s model=%s", idx, len(OPENROUTER_FREE_KEYS), OPENROUTER_FREE_MODEL)
            return text
        except Exception as e:
            last_error = e
            logger.warning("OpenRouter FREE vision key %s/%s failed: %s", idx, len(OPENROUTER_FREE_KEYS), str(e)[:200])
    raise Exception(f"All OpenRouter free vision keys failed: {last_error}")

async def call_gemini(prompt: str, temperature: float = 0.2, max_output_tokens: int = 4096, user_id: Optional[int] = None, allow_paid: Optional[bool] = None) -> str:
    """AI routing: Google first, then free providers. Paid OpenRouter is admin/premium only unless explicitly enabled."""
    global _ACTIVE_AI_JOBS, _WAITING_AI_JOBS, _GOOGLE_KEY_INDEX
    effective_user_id = int(user_id or _AI_USER_ID.get() or 0)
    async with _AI_QUEUE_LOCK:
        _WAITING_AI_JOBS += 1
    async with _GEMINI_SEMAPHORE:
        async with _AI_QUEUE_LOCK:
            _WAITING_AI_JOBS = max(0, _WAITING_AI_JOBS - 1)
            _ACTIVE_AI_JOBS += 1
        try:
            last_error = None
            if GOOGLE_API_KEYS:
                total = len(GOOGLE_API_KEYS)
                for idx, api_key in enumerate(GOOGLE_API_KEYS):
                    try:
                        result = await _call_google_text(prompt, api_key, temperature, max_output_tokens)
                        logger.info("AI provider: Google direct key=%s/%s model=%s", idx + 1, total, _google_model_name())
                        return result
                    except Exception as e:
                        last_error = e
                        logger.warning("Google direct key %s/%s failed: %s", idx + 1, total, str(e)[:200])
                logger.warning("All Google keys failed. Trying free providers before any paid fallback. Last error: %s", last_error)

            if FREE_PROVIDER_MODE:
                try:
                    return await _call_openrouter_free_text(prompt, temperature, max_output_tokens)
                except Exception as e:
                    last_error = e
                    logger.warning("OpenRouter free provider failed: %s", str(e)[:200])
                try:
                    return await _call_groq_text(prompt, temperature, max_output_tokens)
                except Exception as e:
                    last_error = e
                    logger.warning("Groq free provider failed: %s", str(e)[:200])
                try:
                    return await _call_cloudflare_text(prompt, temperature, max_output_tokens)
                except Exception as e:
                    last_error = e
                    logger.warning("Cloudflare free provider failed: %s", str(e)[:200])

            paid_allowed = allow_paid if allow_paid is not None else await _paid_fallback_allowed(effective_user_id)
            if paid_allowed:
                result = await _call_openrouter_text(prompt, temperature, max_output_tokens, model=GEMINI_MODEL_NAME)
                logger.info("AI provider: OpenRouter PAID fallback model=%s user_id=%s", GEMINI_MODEL_NAME, effective_user_id)
                return result

            raise Exception(
                "كل مزودي الذكاء المجاني مشغولون/منتهية حصتهم الآن، والمدفوع مقفل للطلاب العاديين لحماية التكلفة. "
                "سيعمل الكاش للمحاضرات المخزنة بدون API."
            )
        finally:
            async with _AI_QUEUE_LOCK:
                _ACTIVE_AI_JOBS = max(0, _ACTIVE_AI_JOBS - 1)

async def call_gemini_vision(image_bytes: bytes, user_id: Optional[int] = None, allow_paid: Optional[bool] = None) -> str:
    """Vision routing: Google first, then free OpenRouter if configured. Paid vision is admin/premium only."""
    effective_user_id = int(user_id or _AI_USER_ID.get() or 0)
    last_error = None
    if GOOGLE_API_KEYS:
        for idx, api_key in enumerate(GOOGLE_API_KEYS, 1):
            try:
                text = await _call_google_vision(image_bytes, api_key)
                if text:
                    logger.info("Vision provider: Google direct key=%s/%s model=%s", idx, len(GOOGLE_API_KEYS), _google_model_name())
                    return text
            except Exception as e:
                last_error = e
                logger.warning("Google vision key %s/%s failed: %s", idx, len(GOOGLE_API_KEYS), str(e)[:200])
        logger.warning("All Google vision keys failed. Trying free vision before paid. Last error: %s", last_error)
    if FREE_PROVIDER_MODE:
        try:
            return await _call_openrouter_free_vision(image_bytes)
        except Exception as e:
            last_error = e
            logger.warning("OpenRouter free vision failed: %s", str(e)[:200])
    paid_allowed = allow_paid if allow_paid is not None else await _paid_fallback_allowed(effective_user_id)
    if paid_allowed:
        try:
            text = await _call_openrouter_vision(image_bytes, model=GEMINI_MODEL_NAME)
            logger.info("Vision provider: OpenRouter PAID fallback model=%s user_id=%s", GEMINI_MODEL_NAME, effective_user_id)
            return text
        except Exception as e:
            logger.error("Paid Vision API error: %s", e)
            return ""
    logger.warning("Vision failed; paid fallback disabled for user_id=%s. Last error=%s", effective_user_id, last_error)
    return ""


# =============================================================================
# TEXT EXTRACTION
# =============================================================================

def classify_document(mime: str, name: str) -> str:
    name = (name or "").lower()
    mime = (mime or "").lower()
    if name.endswith(".zip") or "zip" in mime or "compressed" in mime:
        return "zip"
    if name.endswith(".pdf") or "pdf" in mime:
        return "pdf"
    if name.endswith((".docx", ".doc")) or "word" in mime:
        return "docx"
    if name.endswith((".pptx", ".ppt")) or "presentation" in mime or "powerpoint" in mime:
        return "pptx"
    if name.endswith((".jpg", ".jpeg", ".png", ".webp")) or mime.startswith("image/"):
        return "image"
    if name.endswith(".txt") or mime.startswith("text/"):
        return "txt"
    return "unsupported"

def extract_text(file_type: str, raw_bytes: bytes) -> str:
    f = io.BytesIO(raw_bytes)
    try:
        if file_type == "pdf":
            return "\n".join([p.extract_text() or "" for p in PdfReader(f).pages]).strip()
        if file_type == "docx":
            return "\n".join([p.text for p in Document(f).paragraphs if p.text and p.text.strip()]).strip()
        if file_type == "pptx":
            parts = []
            for slide in Presentation(f).slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        parts.append(shape.text.strip())
            return "\n".join(parts).strip()
        if file_type == "txt":
            return raw_bytes.decode("utf-8", errors="ignore").strip()
        if file_type in {"image", "zip"}:
            return ""
    except Exception as e:
        logger.exception("Extraction error: %s", e)
    return ""

# =============================================================================
# QUIZ & SUMMARY LOGIC
# =============================================================================

LEVEL_LIMITS = {"basic": 50, "cases": 5, "challenge": 2}
BASIC_COUNTS = (5, 10, 20, 30, 40)
QUESTION_SCHEMA = """Return valid JSON array only. No markdown. Each item must be: { "question": "...", "options": {"A":"...", "B":"...", "C":"...", "D":"...", "E":"..."}, "correct": "A", "explanation": "..." } Exactly 5 options A-E. correct is one of A-E."""
ACADEMIC_PREFIX = "Educational medical simulation. Do not provide real medical advice. "

LEVEL_PROMPTS = {
    "basic": f"{ACADEMIC_PREFIX}Generate direct medical recall MCQs. Straightforward factual questions. No patient vignettes. {QUESTION_SCHEMA}",
    "cases": f"{ACADEMIC_PREFIX}Generate USMLE Step 2 CK level CLINICAL VIGNETTES. Write detailed patient scenarios including age, gender, chief complaint, vital signs, relevant physical exam, and lab findings. The question must ask for the 'most likely diagnosis', 'next best step in management', or 'best initial test'. Distractors must be highly plausible differential diagnoses. {QUESTION_SCHEMA}",
    "challenge": f"{ACADEMIC_PREFIX}Generate EXTREMELY HARD, multi-step clinical reasoning MCQs (USMLE Step 1/Step 2 CK Hard tier). The student must first diagnose an atypical or complex presentation from a vignette, and then answer a 2nd or 3rd order question (e.g., underlying pathophysiology, specific pharmacological mechanism of action, embryological origin, or specific genetic mutation). All 5 options must be highly homogeneous and tricky. {QUESTION_SCHEMA}",
}

SUMMARY_PROMPTS = {
    "disease_v2": f"{ACADEMIC_PREFIX}High-yield medical summary with headings: Overview, Etiology, Clinical, Diagnosis, Management, Complications, Pearls.",
    "highyield": f"{ACADEMIC_PREFIX}Concise high-yield notes, bullet points. End with Quick Recall.",
    "osce": f"{ACADEMIC_PREFIX}OSCE-style approach: History, Examination, Investigations, Management, Counseling, Red Flags.",
}

def extract_json_array(raw: str) -> Optional[list]:
    """Extract a JSON array from AI output safely."""
    raw = (raw or "").strip()
    if not raw:
        return None
    raw = raw.replace("```json", "").replace("```", "").strip()
    try:
        obj = json.loads(raw)
        if isinstance(obj, list):
            return obj
        if isinstance(obj, dict):
            for key in ("questions", "items", "data"):
                if isinstance(obj.get(key), list):
                    return obj[key]
    except Exception:
        pass
    match = re.search(r"\[[\s\S]*\]", raw)
    if match:
        try:
            obj = json.loads(match.group(0))
            return obj if isinstance(obj, list) else None
        except Exception:
            return None
    return None

def _normalize_options(options: Any) -> Dict[str, str]:
    """Convert options from dict/list formats to A-E dict."""
    letters = list("ABCDE")
    out: Dict[str, str] = {}
    if isinstance(options, dict):
        for letter in letters:
            val = options.get(letter) or options.get(letter.lower()) or options.get(str(letter))
            if val is not None and str(val).strip():
                out[letter] = str(val).strip()
    elif isinstance(options, list):
        for idx, item in enumerate(options[:5]):
            if idx >= len(letters):
                break
            if isinstance(item, dict):
                val = (
                    item.get("text")
                    or item.get("option")
                    or item.get("value")
                    or item.get("answer")
                    or item.get("label")
                    or next(iter(item.values()), "")
                )
            else:
                val = item
            if val is not None and str(val).strip():
                out[letters[idx]] = str(val).strip()
    return out

def normalize_questions(items: Any) -> List[dict]:
    """Normalize AI/cached questions so the rest of the bot never crashes on list/dict mismatch."""
    if isinstance(items, dict):
        for key in ("questions", "items", "data"):
            if isinstance(items.get(key), list):
                items = items[key]
                break
        else:
            items = [items]
    if not isinstance(items, list):
        return []

    clean: List[dict] = []
    letters = list("ABCDE")

    for q in items:
        if not isinstance(q, dict):
            continue

        question = str(q.get("question") or q.get("stem") or q.get("q") or "").strip()
        options = _normalize_options(q.get("options") or q.get("choices") or q.get("answers") or [])
        correct = str(q.get("correct") or q.get("answer") or q.get("correct_answer") or "").strip().upper()
        explanation = str(q.get("explanation") or q.get("rationale") or "Based on lecture text.").strip()

        if correct not in letters and correct:
            for letter, value in options.items():
                if correct.lower() == str(value).strip().lower():
                    correct = letter
                    break

        if not question or correct not in letters:
            continue

        if len([v for v in options.values() if str(v).strip()]) < 4:
            continue
        missing_required = False
        for letter in "ABCD":
            if letter not in options or not str(options[letter]).strip():
                missing_required = True
                break
        if missing_required:
            continue
        if "E" not in options or not str(options["E"]).strip():
            options["E"] = "None of the above"

        clean.append({
            "question": question,
            "options": {letter: options.get(letter, "") for letter in letters},
            "correct": correct,
            "explanation": explanation or "Based on lecture text.",
        })

    return clean

def shuffle_options(q: dict) -> dict:
    nq = normalize_questions([q])
    if not nq:
        return {}
    q = nq[0]
    letters = list("ABCDE")
    opts = q.get("options", {})
    correct = q.get("correct", "").upper()
    if correct not in letters or any(l not in opts for l in letters):
        return q
    pairs = [(opts[l], l == correct) for l in letters]
    random.shuffle(pairs)
    new_opts, new_correct = {}, "A"
    for idx, (text, is_correct) in enumerate(pairs):
        letter = letters[idx]
        new_opts[letter] = text
        if is_correct:
            new_correct = letter
    q = dict(q)
    q["options"] = new_opts
    q["correct"] = new_correct
    return q

def split_text(text: str, chunk_size: int = 10000) -> List[str]:
    return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]


def representative_text(text: str, max_chars: int = MAX_TEXT_CHARS_FOR_AI) -> str:
    """Take an even sample from beginning, middle, and end so big lectures are represented fairly."""
    text = (text or "").strip()
    if len(text) <= max_chars:
        return text
    if max_chars < 3000:
        return text[:max_chars]
    part = max_chars // 3
    start = text[:part]
    mid_start = max(0, (len(text) // 2) - (part // 2))
    middle = text[mid_start:mid_start + part]
    end = text[-(max_chars - len(start) - len(middle)):]
    return (
        start.rstrip()
        + "\n\n...[تم اختصار جزء من منتصف المحاضرة لتقليل التكلفة]...\n\n"
        + middle.strip()
        + "\n\n...[تم اختصار جزء من نهاية المحاضرة لتقليل التكلفة]...\n\n"
        + end.lstrip()
    )[:max_chars + 300]

def canonical_text_hash(text: str) -> str:
    norm = re.sub(r"\s+", " ", (text or "").strip().lower())
    return "text:" + hashlib.sha256(norm.encode("utf-8", errors="ignore")).hexdigest()

async def generate_question_bank(text: str, level: str) -> List[dict]:
    source_text = representative_text(text, MAX_TEXT_CHARS_FOR_AI)
    chunks, all_questions, target = split_text(source_text, chunk_size=10000), [], LEVEL_LIMITS[level]
    for chunk in chunks:
        raw = await call_gemini(f"{LEVEL_PROMPTS[level]}\nGenerate up to {target} questions only.\nReturn JSON array only. No markdown. No explanation outside JSON.\n\nSOURCE TEXT:\n{chunk}", temperature=0.15, max_output_tokens=8192)
        parsed = extract_json_array(raw)
        if parsed: all_questions.extend(normalize_questions(parsed))
        if len(all_questions) >= target: break
    return all_questions[:target]

async def get_or_create_question_bank(file_hash: str, text: str, level: str) -> List[dict]:
    lock = get_generation_lock(f"questions:{file_hash}:{level}")
    async with lock:
        cached = await db_get_questions(file_hash, level)
        if cached is not None: return cached
        questions = await generate_question_bank(text, level)
        if questions: await db_save_questions(file_hash, level, questions)
        return questions

async def get_or_create_summary(file_hash: str, text: str, style: str) -> str:
    lock = get_generation_lock(f"summary:{file_hash}:{style}")
    async with lock:
        cached = await db_get_summary(file_hash, style)
        if cached: return cached
        summary = ""
        try: summary = await call_gemini(f"{SUMMARY_PROMPTS.get(style, SUMMARY_PROMPTS['disease_v2'])}\nSOURCE:\n{representative_text(text, MAX_TEXT_CHARS_FOR_AI)}", temperature=0.2, max_output_tokens=8192)
        except Exception as e: logger.error("Summary error: %s", e)
        if summary and summary.strip(): await db_save_summary(file_hash, style, summary)
        return summary

# =============================================================================
# OSCE SIMULATOR V4 — خطوة بخطوة (Step-by-Step Clinical Assessment)
# =============================================================================

OSCE_V4_DIFFICULTY = "medium-hard"

OSCE_V4_PHASE_TIME = {
    "diagnosis":   120,
    "management":  150,
    "workup":      180,
}

OSCE_V4_LAB_OPTIONS = [
    ("CBC",          "🩸 CBC — صورة الدم الكاملة"),
    ("CMP",          "🧪 CMP — كيمياء الدم الشاملة"),
    ("LFT",          "🫀 LFT — وظائف الكبد"),
    ("COAG",         "🩺 Coagulation — تخثر الدم"),
    ("CARDIAC",      "❤️ Troponin / BNP — قلب"),
    ("THYROID",      "🦋 TSH / T4 — غدة درقية"),
    ("URINE",        "💛 Urinalysis — تحليل بول"),
    ("CULTURE",      "🦠 Culture & Sensitivity — مزرعة"),
    ("LIPASE",       "🫁 Lipase / Amylase — بنكرياس"),
    ("ABG",          "💨 ABG — غازات الدم"),
]

OSCE_V4_IMAGING_OPTIONS = [
    ("CXR",          "🫁 Chest X-Ray — صورة صدر"),
    ("AXR",          "🩻 Abdominal X-Ray — صورة بطن"),
    ("USS_ABD",      "🔊 Ultrasound Abdomen — سونار البطن"),
    ("USS_PELV",     "🔊 Ultrasound Pelvis — سونار الحوض"),
    ("ECHO",         "❤️ Echocardiogram — صدى القلب"),
    ("CT_HEAD",      "🧠 CT Head — سي تي رأس"),
    ("CT_CHEST",     "🫁 CT Chest — سي تي صدر"),
    ("CT_ABD",       "🩻 CT Abdomen/Pelvis — سي تي بطن"),
    ("MRI_BRAIN",    "🧲 MRI Brain — رنين مغناطيسي دماغ"),
    ("MRI_SPINE",    "🦴 MRI Spine — رنين مغناطيسي عمود"),
]

OSCE_V4_CASE_PROMPT = """
You are a Senior Medical Examiner. Generate ONE realistic, medium-to-hard OSCE clinical case.
Difficulty: {difficulty} — avoid trivial presentations. Use atypical features or uncommon complications.

Return ONLY valid JSON (no markdown, no explanation):

{{
  "station_title": "Short neutral title (e.g., 'Acute Abdomen in a Young Female')",
  "case_vignette": "3-5 sentence clinical scenario: patient age, sex, chief complaint, duration, key HPI details, relevant PMH, vital signs, relevant physical exam findings. Write like a real exam question.",
  "hidden_diagnosis": "Exact medical diagnosis",
  "differential_diagnoses": ["Most likely", "2nd", "3rd"],
  "key_discriminating_features": "1-2 sentences: what features point to the correct diagnosis",
  "ideal_management": "Step-by-step management: immediate stabilization, investigations, treatment, disposition. Be specific.",
  "labs": {{
    "CBC":     "Result text or 'Normal'",
    "CMP":     "Result text or 'Normal'",
    "LFT":     "Result text or 'Normal'",
    "COAG":    "Result text or 'Normal'",
    "CARDIAC": "Result text or 'Normal'",
    "THYROID": "Result text or 'Normal'",
    "URINE":   "Result text or 'Normal'",
    "CULTURE": "Result text or 'Pending'",
    "LIPASE":  "Result text or 'Normal'",
    "ABG":     "Result text or 'Normal'"
  }},
  "imaging": {{
    "CXR":      "Findings or 'No acute pathology'",
    "AXR":      "Findings or 'No acute pathology'",
    "USS_ABD":  "Findings or 'No acute pathology'",
    "USS_PELV": "Findings or 'No acute pathology'",
    "ECHO":     "Findings or 'No significant abnormality'",
    "CT_HEAD":  "Findings or 'No intracranial pathology'",
    "CT_CHEST": "Findings or 'No acute pathology'",
    "CT_ABD":   "Findings or 'No acute pathology'",
    "MRI_BRAIN":"Findings or 'No significant abnormality'",
    "MRI_SPINE":"Findings or 'No significant abnormality'"
  }},
  "key_investigations": ["List the 3-5 most important investigations to order for this case"],
  "management_mark_scheme": [
    "Immediate action 1",
    "Immediate action 2",
    "Investigation ordered",
    "Treatment given",
    "Disposition/consult"
  ],
  "teaching_pearl": "One high-yield clinical pearl about this case"
}}

SOURCE TEXT (lecture/notes to base the case on):
{text}
"""

async def generate_osce_v4_case(text: str, max_chars: int = 12000) -> dict:
    prompt = OSCE_V4_CASE_PROMPT.format(
        difficulty=OSCE_V4_DIFFICULTY,
        text=text[:max_chars]
    )
    raw = await call_gemini(prompt, temperature=0.2, max_output_tokens=4096)
    parsed = _extract_json_object_v4(raw)
    return _normalize_osce_v4_case(parsed or {})


def _extract_json_object_v4(raw: str) -> Optional[dict]:
    raw = (raw or "").strip().replace("```json", "").replace("```", "").strip()
    try:
        obj = json.loads(raw)
        return obj if isinstance(obj, dict) else None
    except Exception:
        pass
    match = re.search(r"\{.*\}", raw, flags=re.DOTALL)
    if match:
        try:
            obj = json.loads(match.group(0))
            return obj if isinstance(obj, dict) else None
        except Exception:
            pass
    return None


def _normalize_osce_v4_case(case: dict) -> dict:
    def s(key, default=""): return str(case.get(key) or default).strip()
    def lst(key, default=None): return case.get(key) if isinstance(case.get(key), list) else (default or [])
    def dct(key): return case.get(key) if isinstance(case.get(key), dict) else {}

    return {
        "station_title":              s("station_title", "Clinical OSCE Station"),
        "case_vignette":              s("case_vignette", "A patient presents with an acute complaint."),
        "hidden_diagnosis":           s("hidden_diagnosis", "Diagnosis not specified"),
        "differential_diagnoses":     lst("differential_diagnoses", ["Unknown"]),
        "key_discriminating_features":s("key_discriminating_features", "See case details."),
        "ideal_management":           s("ideal_management", "Assess, stabilize, investigate, treat."),
        "labs":                       dct("labs"),
        "imaging":                    dct("imaging"),
        "key_investigations":         lst("key_investigations", []),
        "management_mark_scheme":     lst("management_mark_scheme", ["Assess patient", "Order investigations", "Treat appropriately"]),
        "teaching_pearl":             s("teaching_pearl", "Always consider the full differential."),
    }


def create_osce_v4_session(case: dict, file_hash: str) -> dict:
    return {
        "version":           4,
        "case":              case,
        "file_hash":         file_hash,
        "phase":             "diagnosis",
        "started_at_ts":     time.time(),
        "phase_started_ts":  time.time(),
        "phase_scores": {
            "diagnosis":   None,
            "management":  None,
            "workup":      None,
        },
        "student_answers": {
            "diagnosis":   "",
            "management":  "",
        },
        "ordered_labs":    [],
        "ordered_imaging": [],
        "results_shown":   [],
        "turn_count":      0,
    }


def _osce_v4_phase_remaining(session: dict) -> int:
    phase = session.get("phase", "diagnosis")
    limit = OSCE_V4_PHASE_TIME.get(phase, 120)
    elapsed = time.time() - session.get("phase_started_ts", time.time())
    return max(0, int(limit - elapsed))


def _osce_v4_fmt_time(secs: int) -> str:
    return f"{secs // 60:02d}:{secs % 60:02d}"


async def _evaluate_osce_v4_diagnosis(student_answer: str, case: dict) -> dict:
    prompt = f"""
You are a medical examiner. Evaluate the student's diagnosis answer.

Correct diagnosis: {case.get('hidden_diagnosis', '')}
Differentials: {json.dumps(case.get('differential_diagnoses', []), ensure_ascii=False)}
Key features: {case.get('key_discriminating_features', '')}

Student answered: "{student_answer}"

Return ONLY valid JSON:
{{
  "score": <0-100>,
  "is_correct": <true/false>,
  "feedback_ar": "2-3 جمل تقييم بالعربي: هل أصاب؟ ما الذي فاته؟ لماذا التشخيص الصحيح هو كذا؟",
  "correct_diagnosis": "{case.get('hidden_diagnosis', '')}",
  "key_point": "نقطة مهمة واحدة يجب أن يتذكرها الطالب"
}}
"""
    raw = await call_gemini(prompt, temperature=0.1, max_output_tokens=500)
    result = _extract_json_object_v4(raw)
    if not result:
        return {
            "score": 50,
            "is_correct": False,
            "feedback_ar": "تعذّر التقييم التلقائي. التشخيص الصحيح: " + case.get('hidden_diagnosis', ''),
            "correct_diagnosis": case.get('hidden_diagnosis', ''),
            "key_point": case.get('key_discriminating_features', ''),
        }
    return result


async def _evaluate_osce_v4_management(student_answer: str, case: dict) -> dict:
    mark_scheme = case.get("management_mark_scheme", [])
    prompt = f"""
You are a medical examiner. Evaluate the student's management plan.

Ideal management: {case.get('ideal_management', '')}
Mark scheme: {json.dumps(mark_scheme, ensure_ascii=False)}

Student answered: "{student_answer}"

Return ONLY valid JSON:
{{
  "score": <0-100>,
  "covered_points": ["list of mark scheme points the student covered"],
  "missed_points": ["list of mark scheme points the student missed"],
  "feedback_ar": "3-4 جمل تقييم بالعربي: ما الذي ذكره الطالب صح؟ ما الذي فاته؟",
  "ideal_summary": "ملخص الخطة المثالية في جملتين"
}}
"""
    raw = await call_gemini(prompt, temperature=0.1, max_output_tokens=600)
    result = _extract_json_object_v4(raw)
    if not result:
        return {
            "score": 50,
            "covered_points": [],
            "missed_points": mark_scheme,
            "feedback_ar": "تعذّر التقييم. الخطة المثالية: " + case.get('ideal_management', ''),
            "ideal_summary": case.get('ideal_management', '')[:200],
        }
    return result


def _evaluate_osce_v4_workup(session: dict) -> dict:
    case = session.get("case", {})
    key_inv = set(case.get("key_investigations", []))
    ordered = set(session.get("ordered_labs", []) + session.get("ordered_imaging", []))

    def normalize(s): return s.upper().replace(" ", "").replace("-", "")
    key_norm = {normalize(k) for k in key_inv}
    ordered_norm = {normalize(o) for o in ordered}
    hit = key_norm & ordered_norm
    missed = key_norm - ordered_norm

    if len(key_norm) == 0:
        score = 70
    else:
        score = min(100, int(len(hit) / len(key_norm) * 100))

    unnecessary = ordered - key_norm
    score = max(0, score - len(unnecessary) * 5)

    return {
        "score":        score,
        "key_inv":      list(key_inv),
        "ordered":      list(ordered),
        "hit":          list(hit),
        "missed":       list(missed),
        "unnecessary":  list(unnecessary),
    }


def _osce_v4_vignette_text(session: dict) -> str:
    case = session["case"]
    rem = _osce_v4_phase_remaining(session)
    return (
        f"🎭 OSCE Station — {case.get('station_title', '')}\n"
        f"{'━' * 35}\n\n"
        f"📋 الحالة السريرية:\n{case.get('case_vignette', '')}\n\n"
        f"{'━' * 35}\n"
        f"🧠 السؤال 1/3 — التشخيص\n\n"
        f"بناءً على الحالة أعلاه، اكتب:\n"
        f"• التشخيص الأرجح\n"
        f"• واذكر تشخيصين تفريقيين محتملين\n\n"
        f"⏱ الوقت المتبقي: {_osce_v4_fmt_time(rem)}"
    )


def _osce_v4_management_prompt_text(session: dict, diag_feedback: dict) -> str:
    case = session["case"]
    score = diag_feedback.get("score", 0)
    is_correct = diag_feedback.get("is_correct", False)
    feedback = diag_feedback.get("feedback_ar", "")
    correct_dx = diag_feedback.get("correct_diagnosis", "")
    key_point = diag_feedback.get("key_point", "")
    rem = _osce_v4_phase_remaining(session)

    emoji = "✅" if is_correct else "⚠️"
    return (
        f"{emoji} تقييم التشخيص — {score}/100\n"
        f"{'━' * 35}\n\n"
        f"📝 {feedback}\n\n"
        f"✅ التشخيص الصحيح: {correct_dx}\n"
        f"💡 نقطة مهمة: {key_point}\n\n"
        f"{'━' * 35}\n"
        f"💊 السؤال 2/3 — خطة الإدارة\n\n"
        f"اكتب خطة الإدارة الكاملة لهذه الحالة:\n"
        f"• الإجراء الفوري (ABC, stabilization)\n"
        f"• التحاليل والصور الأساسية\n"
        f"• العلاج المحدد\n"
        f"• التوجيه والاستشارة\n\n"
        f"⏱ الوقت المتبقي: {_osce_v4_fmt_time(rem)}"
    )


def _osce_v4_workup_keyboard(session: dict) -> InlineKeyboardMarkup:
    ordered_labs    = set(session.get("ordered_labs", []))
    ordered_imaging = set(session.get("ordered_imaging", []))
    all_ordered     = ordered_labs | ordered_imaging

    rows = []
    rows.append([InlineKeyboardButton("🧪 ── التحاليل المخبرية ──", callback_data="osce4:noop")])

    lab_row = []
    for key, label in OSCE_V4_LAB_OPTIONS:
        done = key in ordered_labs
        btn_label = f"✅ {label}" if done else label
        lab_row.append(InlineKeyboardButton(btn_label, callback_data=f"osce4:lab:{key}"))
        if len(lab_row) == 2:
            rows.append(lab_row)
            lab_row = []
    if lab_row:
        rows.append(lab_row)

    rows.append([InlineKeyboardButton("🩻 ── التصوير الشعاعي ──", callback_data="osce4:noop")])

    img_row = []
    for key, label in OSCE_V4_IMAGING_OPTIONS:
        done = key in ordered_imaging
        btn_label = f"✅ {label}" if done else label
        img_row.append(InlineKeyboardButton(btn_label, callback_data=f"osce4:img:{key}"))
        if len(img_row) == 2:
            rows.append(img_row)
            img_row = []
    if img_row:
        rows.append(img_row)

    rows.append([
        InlineKeyboardButton(
            f"📊 مراجعة المطلوب ({len(all_ordered)})",
            callback_data="osce4:review_workup"
        ),
        InlineKeyboardButton("🏁 إنهاء وتقرير", callback_data="osce4:finish"),
    ])
    rows.append([InlineKeyboardButton("⬅️ القائمة الرئيسية", callback_data="back:mode")])

    return InlineKeyboardMarkup(rows)


def _osce_v4_workup_prompt_text(session: dict, mgmt_feedback: dict) -> str:
    score       = mgmt_feedback.get("score", 0)
    feedback    = mgmt_feedback.get("feedback_ar", "")
    covered     = mgmt_feedback.get("covered_points", [])
    missed      = mgmt_feedback.get("missed_points", [])
    ideal       = mgmt_feedback.get("ideal_summary", "")
    rem         = _osce_v4_phase_remaining(session)

    covered_str = "\n".join(f"  ✅ {p}" for p in covered) or "  (لا شيء)"
    missed_str  = "\n".join(f"  ❌ {p}" for p in missed)  or "  ✅ لا يوجد"

    return (
        f"💊 تقييم خطة الإدارة — {score}/100\n"
        f"{'━' * 35}\n\n"
        f"📝 {feedback}\n\n"
        f"✅ ما ذكرته:\n{covered_str}\n\n"
        f"❌ ما فاتك:\n{missed_str}\n\n"
        f"💡 الخطة المثالية: {ideal}\n\n"
        f"{'━' * 35}\n"
        f"🔬 السؤال 3/3 — التحاليل والتصوير\n\n"
        f"اختر من الأزرار أدناه التحاليل والصور التي تريد طلبها.\n"
        f"يمكنك الضغط على أكثر من خيار.\n"
        f"عند الانتهاء اضغط: 🏁 إنهاء وتقرير\n\n"
        f"⏱ الوقت المتبقي: {_osce_v4_fmt_time(rem)}"
    )


async def _osce_v4_final_report(session: dict) -> str:
    case          = session["case"]
    phase_scores  = session.get("phase_scores", {})
    answers       = session.get("student_answers", {})
    workup_eval   = _evaluate_osce_v4_workup(session)

    dx_score   = phase_scores.get("diagnosis",  {}).get("score",  0) if isinstance(phase_scores.get("diagnosis"),  dict) else 0
    mgmt_score = phase_scores.get("management", {}).get("score",  0) if isinstance(phase_scores.get("management"), dict) else 0
    inv_score  = workup_eval.get("score", 0)

    total_score = round((dx_score * 0.35) + (mgmt_score * 0.35) + (inv_score * 0.30))

    ordered_labs    = session.get("ordered_labs", [])
    ordered_imaging = session.get("ordered_imaging", [])

    def lab_name(key):
        for k, lbl in OSCE_V4_LAB_OPTIONS:
            if k == key: return lbl
        return key

    def img_name(key):
        for k, lbl in OSCE_V4_IMAGING_OPTIONS:
            if k == key: return lbl
        return key

    ordered_str = ""
    if ordered_labs:
        ordered_str += "  تحاليل: " + ", ".join(lab_name(k) for k in ordered_labs) + "\n"
    if ordered_imaging:
        ordered_str += "  صور:     " + ", ".join(img_name(k) for k in ordered_imaging) + "\n"
    if not ordered_str:
        ordered_str = "  (لم يطلب شيئاً)\n"

    missed_inv_str = ""
    if workup_eval.get("missed"):
        missed_inv_str = "\n  ❌ فاتك: " + ", ".join(workup_eval["missed"])

    grade = (
        "🏆 ممتاز"      if total_score >= 85 else
        "✅ جيد جداً"   if total_score >= 70 else
        "⚠️ مقبول"     if total_score >= 55 else
        "❌ يحتاج مراجعة"
    )

    elapsed = int(time.time() - session.get("started_at_ts", time.time()))

    pearl_prompt = (
        f"Give ONE short Arabic clinical teaching pearl (max 2 sentences) "
        f"about: {case.get('hidden_diagnosis', '')}. "
        f"Focus on what a student commonly misses. Reply in Arabic only."
    )
    try:
        pearl = await call_gemini(pearl_prompt, temperature=0.2, max_output_tokens=150)
    except Exception:
        pearl = case.get("teaching_pearl", "راجع الحالة مرة أخرى.")

    report = (
        f"📋 التقرير النهائي — OSCE V4\n"
        f"{'━' * 38}\n\n"
        f"🏛 المحطة: {case.get('station_title', '')}\n"
        f"⏱ الوقت الكلي: {elapsed // 60} دقيقة {elapsed % 60} ثانية\n\n"

        f"{'━' * 38}\n"
        f"📊 النتيجة الإجمالية: {total_score}/100 — {grade}\n"
        f"{'━' * 38}\n\n"

        f"🧠 1. التشخيص ({dx_score}/100)\n"
        f"   إجابتك: {answers.get('diagnosis', '(فارغ)')[:200]}\n"
        f"   ✅ التشخيص الصحيح: {case.get('hidden_diagnosis', '')}\n\n"

        f"💊 2. الإدارة ({mgmt_score}/100)\n"
        f"   إجابتك: {answers.get('management', '(فارغ)')[:200]}\n"
        f"   💡 الخطة المثالية:\n   {case.get('ideal_management', '')[:400]}\n\n"

        f"🔬 3. التحاليل والصور ({inv_score}/100)\n"
        f"{ordered_str}"
        f"   ✅ الأهم كان:\n   " + "\n   ".join(f"• {i}" for i in case.get("key_investigations", [])) +
        missed_inv_str + "\n\n"

        f"{'━' * 38}\n"
        f"📖 النقطة التعليمية:\n{pearl}\n\n"

        f"{'━' * 38}\n"
        f"🎓 التشخيصات التفريقية الكاملة:\n" +
        "\n".join(f"  • {d}" for d in case.get("differential_diagnoses", [])) +
        f"\n\n🔑 مفتاح التمييز: {case.get('key_discriminating_features', '')}"
    )

    return report


async def handle_osce_v4_text(update: Update, context: ContextTypes.DEFAULT_TYPE) -> bool:
    session = context.user_data.get("osce_v4_session")
    if not session:
        return False

    text  = (update.message.text or "").strip()
    phase = session.get("phase", "diagnosis")

    if text.lower() in {"انهاء", "إنهاء", "end", "finish", "stop osce", "خلاص"}:
        await update.message.reply_text("⏳ جاري إصدار التقرير النهائي...")
        await _osce_v4_send_final_report(update, context, session)
        return True

    if phase == "diagnosis":
        session["student_answers"]["diagnosis"] = text
        await update.message.chat.send_action(action="typing")

        diag_feedback = await _evaluate_osce_v4_diagnosis(text, session["case"])
        session["phase_scores"]["diagnosis"] = diag_feedback

        session["phase"]            = "management"
        session["phase_started_ts"] = time.time()

        prompt_text = _osce_v4_management_prompt_text(session, diag_feedback)
        await update.message.reply_text(prompt_text)
        return True

    if phase == "management":
        session["student_answers"]["management"] = text
        await update.message.chat.send_action(action="typing")

        mgmt_feedback = await _evaluate_osce_v4_management(text, session["case"])
        session["phase_scores"]["management"] = mgmt_feedback

        session["phase"]            = "workup"
        session["phase_started_ts"] = time.time()

        prompt_text = _osce_v4_workup_prompt_text(session, mgmt_feedback)
        kb = _osce_v4_workup_keyboard(session)
        await update.message.reply_text(prompt_text, reply_markup=kb)
        return True

    if phase == "workup":
        await update.message.reply_text(
            "💡 استخدم الأزرار أدناه لاختيار التحاليل والصور.\n"
            "عند الانتهاء اضغط: 🏁 إنهاء وتقرير",
            reply_markup=_osce_v4_workup_keyboard(session),
        )
        return True

    return False


async def osce_v4_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()

    session = context.user_data.get("osce_v4_session")
    if not session:
        await query.edit_message_text(
            "⚠️ انتهت جلسة OSCE أو لا توجد جلسة نشطة.\n"
            "أرسل /start للبدء من جديد."
        )
        return

    data   = query.data
    parts  = data.split(":")
    action = parts[1] if len(parts) > 1 else ""

    if action == "noop":
        return

    if action == "lab" and len(parts) == 3:
        key = parts[2]
        labs_ordered = session.setdefault("ordered_labs", [])

        if key in labs_ordered:
            labs_ordered.remove(key)
            await query.answer(f"❌ تم إلغاء: {key}", show_alert=False)
        else:
            labs_ordered.append(key)
            result_text = session["case"].get("labs", {}).get(key, "النتيجة غير متوفرة.")
            await query.answer(f"🧪 {key}: {result_text[:60]}", show_alert=True)

        try:
            await query.edit_message_reply_markup(
                reply_markup=_osce_v4_workup_keyboard(session)
            )
        except Exception:
            pass
        return

    if action == "img" and len(parts) == 3:
        key = parts[2]
        imgs_ordered = session.setdefault("ordered_imaging", [])

        if key in imgs_ordered:
            imgs_ordered.remove(key)
            await query.answer(f"❌ تم إلغاء: {key}", show_alert=False)
        else:
            imgs_ordered.append(key)
            result_text = session["case"].get("imaging", {}).get(key, "النتيجة غير متوفرة.")
            await query.answer(f"🩻 {key}: {result_text[:60]}", show_alert=True)

        try:
            await query.edit_message_reply_markup(
                reply_markup=_osce_v4_workup_keyboard(session)
            )
        except Exception:
            pass
        return

    if action == "review_workup":
        ordered_labs    = session.get("ordered_labs", [])
        ordered_imaging = session.get("ordered_imaging", [])
        lines = []
        if ordered_labs:
            for k in ordered_labs:
                res = session["case"].get("labs", {}).get(k, "N/A")
                lines.append(f"🧪 {k}: {res}")
        if ordered_imaging:
            for k in ordered_imaging:
                res = session["case"].get("imaging", {}).get(k, "N/A")
                lines.append(f"🩻 {k}: {res}")
        summary = "\n".join(lines) if lines else "لم تطلب أي تحاليل أو صور بعد."
        await query.answer(summary[:200], show_alert=True)
        return

    if action == "finish":
        try:
            await query.edit_message_text("⏳ الممتحن يراجع أداءك... يرجى الانتظار.")
        except Exception:
            pass
        await _osce_v4_send_final_report(update, context, session, query=query)
        return


async def _osce_v4_send_final_report(update: Update, context: ContextTypes.DEFAULT_TYPE, session: dict, query=None):
    report = await _osce_v4_final_report(session)
    context.user_data.pop("osce_v4_session", None)

    user = update.effective_user
    if user:
        phase_scores = session.get("phase_scores", {})
        dx_score   = phase_scores.get("diagnosis",  {}).get("score",  50) if isinstance(phase_scores.get("diagnosis"),  dict) else 50
        mgmt_score = phase_scores.get("management", {}).get("score",  50) if isinstance(phase_scores.get("management"), dict) else 50
        inv_eval   = _evaluate_osce_v4_workup(session)
        inv_score  = inv_eval.get("score", 50)
        total      = round((dx_score * 0.35) + (mgmt_score * 0.35) + (inv_score * 0.30))
        xp_points  = 40 + int(total * 0.5)

        progress = await db_add_xp(user.id, xp_points, "OSCE V4 completed", osce_done=True)
        report += (
            f"\n\n{'━' * 38}\n"
            f"⭐ +{xp_points} XP | 🎖 Level {progress.get('level', 1)} "
            f"| 🔥 Streak {progress.get('streak', 0)} يوم"
        )
        if progress.get("leveled_up"):
            report += "\n🎉 رفعت مستوى! استمر على هذا المستوى."

    chat_id = update.effective_chat.id

    chunks = [report[i:i + 3800] for i in range(0, len(report), 3800)]
    back_kb = InlineKeyboardMarkup([[InlineKeyboardButton("⬅️ القائمة الرئيسية", callback_data="back:mode")]])

    for i, chunk in enumerate(chunks):
        if i == len(chunks) - 1:
            await context.bot.send_message(chat_id, chunk, reply_markup=back_kb)
        else:
            await context.bot.send_message(chat_id, chunk)

# =============================================================================
# ADMIN STAGING COMMANDS (THE PIPELINE)
# =============================================================================

async def stage_start_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id not in ADMIN_USER_IDS:
        return await update.message.reply_text(f"⛔ عذراً، هذا الأمر للإدارة فقط!\n\nرقم הـ ID الخاص بك هو:\n{update.effective_user.id}", parse_mode="Markdown")
    if len(context.args) < 2: return await update.message.reply_text("استخدم: /stage_start [uni] [subject]")
    context.user_data["active_staging"] = {"uni": context.args[0].lower(), "subject": context.args[1].lower()}
    await update.message.reply_text("📥 بدأ التجميع! أي نص/صورة/ملف ترسله سيُخزن للمراجعة.\nللتصدير: /stage_dump", parse_mode="Markdown")

async def stage_dump_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if user_id not in ADMIN_USER_IDS or not context.user_data.get("active_staging"): return
    st = context.user_data["active_staging"]
    uni, subject = st["uni"], st["subject"]

    def get_dump():
        conn = _db_connect()
        rows = conn.execute("SELECT raw_content FROM admin_staging WHERE user_id=? AND uni=? AND subject=? ORDER BY created_at ASC", (user_id, uni, subject)).fetchall()
        content = [r["raw_content"] for r in rows]
        conn.execute("DELETE FROM admin_staging WHERE user_id=? AND uni=? AND subject=?", (user_id, uni, subject))
        conn.commit()
        conn.close()
        return content

    async with _db_lock:
        content_list = await asyncio.to_thread(get_dump)

    if not content_list: return await update.message.reply_text("المستودع فارغ.")

    txt = "\n\n========================================\n".join(content_list)
    f = io.BytesIO(txt.encode('utf-8')); f.name = f"dump_{uni}_{subject}.txt"
    context.user_data.pop("active_staging", None)
    await context.bot.send_document(update.effective_chat.id, f, caption="📦 المستودع تم تصديره وتفريغه. قم بتنظيفه بـ ChatGPT ثم استخدم /upload_final")


async def finish_past_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id if update.effective_user else 0
    if not is_admin_user(user_id):
        return await update.message.reply_text("⚠️ هذا الأمر للأدمن فقط.")
    mem = await build_admin_staging_zip(user_id)
    if not mem:
        return await update.message.reply_text("📦 لا يوجد شيء في صندوق السنوات السابقة حالياً.")
    await update.message.reply_document(document=mem, filename=mem.name, caption="✅ هذا ZIP مراجعة السنوات السابقة. أرسله للمراجعة قبل إدخاله النهائي للبوت.")

async def upload_final_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id not in ADMIN_USER_IDS: return
    if len(context.args) < 2: return await update.message.reply_text("استخدم: /upload_final [uni] [subject] وارفق الـ JSON")
    context.user_data["pending_final"] = {"uni": context.args[0].lower(), "subject": context.args[1].lower()}
    await update.message.reply_text("🔄 مستعد. أرسل الآن ملف الـ JSON النهائي لرفعه للطلبة.")

async def handle_admin_acc(update: Update, context: ContextTypes.DEFAULT_TYPE, text: str):
    st = context.user_data.get("active_staging")
    if not st: return
    async with _db_lock:
        def work():
            conn = _db_connect()
            conn.execute("INSERT INTO admin_staging (user_id, uni, subject, raw_content, created_at) VALUES (?, ?, ?, ?, ?)",
                         (update.effective_user.id, st["uni"], st["subject"], text, datetime.utcnow().isoformat()))
            conn.commit()
            conn.close()
        await asyncio.to_thread(work)
    await update.message.reply_text("✅ [المستودع المؤقت]: تم التقاط وحفظ المحتوى بنجاح.")


# =============================================================================
# ADMIN STAGING ORIGINAL FILE STORAGE
# =============================================================================

def _safe_stage_filename(name: str) -> str:
    name = (name or "file").strip().replace("\\", "_").replace("/", "_")
    name = re.sub(r"[^A-Za-z0-9._\-\u0600-\u06FF ]+", "_", name)
    name = re.sub(r"\s+", "_", name).strip("._")
    return name[:120] or "file"

def _staging_files_root() -> str:
    base_dir = os.path.dirname(DATABASE_PATH) or "."
    return os.path.join(base_dir, "admin_staging_files")

async def admin_save_staging_original(user_id: int, st: dict, filename: str, raw_bytes: bytes) -> str:
    """Save original admin staging upload on disk and return a relative path for the dump."""
    if not st or not raw_bytes:
        return ""
    uni = _safe_stage_filename(str(st.get("uni", "unknown")))
    subject = _safe_stage_filename(str(st.get("subject", "unknown")))
    folder = os.path.join(_staging_files_root(), str(user_id), uni, subject)
    safe_name = _safe_stage_filename(filename or "uploaded_file")
    final_name = f"{datetime.utcnow().strftime('%Y%m%d_%H%M%S_%f')}_{safe_name}"
    full_path = os.path.join(folder, final_name)

    def work():
        os.makedirs(folder, exist_ok=True)
        with open(full_path, "wb") as f:
            f.write(raw_bytes)

    await asyncio.to_thread(work)
    return os.path.relpath(full_path, os.path.dirname(DATABASE_PATH) or ".")

def _admin_staging_file_paths(user_id: int, uni: str, subject: str) -> List[Tuple[str, str]]:
    folder = os.path.join(_staging_files_root(), str(user_id), _safe_stage_filename(uni), _safe_stage_filename(subject))
    out: List[Tuple[str, str]] = []
    if not os.path.isdir(folder):
        return out
    for root, _, files in os.walk(folder):
        for fn in sorted(files):
            full = os.path.join(root, fn)
            rel = os.path.relpath(full, folder)
            out.append((full, rel))
    return out

async def build_admin_staging_zip(user_id: int) -> Optional[io.BytesIO]:
    def fetch_rows():
        conn = _db_connect()
        rows = conn.execute(
            "SELECT uni, subject, raw_content, created_at FROM admin_staging WHERE user_id=? ORDER BY uni, subject, created_at",
            (user_id,),
        ).fetchall()
        conn.close()
        return [dict(r) for r in rows]

    async with _db_lock:
        rows = await asyncio.to_thread(fetch_rows)

    grouped: Dict[Tuple[str, str], List[dict]] = {}
    for r in rows:
        grouped.setdefault((r.get("uni", "unknown"), r.get("subject", "unknown")), []).append(r)

    # include groups that only have original files, even if raw text is empty
    root = os.path.join(_staging_files_root(), str(user_id))
    if os.path.isdir(root):
        for uni in os.listdir(root):
            uni_dir = os.path.join(root, uni)
            if not os.path.isdir(uni_dir):
                continue
            for subject in os.listdir(uni_dir):
                subject_dir = os.path.join(uni_dir, subject)
                if os.path.isdir(subject_dir):
                    grouped.setdefault((uni, subject), [])

    if not grouped:
        return None

    mem = io.BytesIO()
    with zipfile.ZipFile(mem, "w", zipfile.ZIP_DEFLATED) as z:
        readme = (
            "هذا ZIP صادر من بوت MedMCQ.\n"
            "الآن يحتوي على النص المستخرج + الملفات الأصلية.\n\n"
            "المسارات المهمة:\n"
            "- raw_questions.txt: النصوص التي جمعها البوت.\n"
            "- original_files/: ملفات PDF/Word/PPT/صور الأصلية كما أرسلها الأدمن.\n"
            "- metadata.json: معلومات مختصرة عن محتوى الحزمة.\n\n"
            "أرسل هذا ZIP إلى ChatGPT لتنظيف الأسئلة وتحويلها إلى JSON array بالشكل:\n"
            "question/options/correct/explanation\n"
            "ثم ارفع الناتج باستخدام /upload_final uni subject\n"
        )
        z.writestr("README.txt", readme)
        meta = {"created_at": datetime.utcnow().isoformat(), "user_id": user_id, "groups": []}
        for (uni, subject), items in sorted(grouped.items()):
            raw_parts = []
            for i, item in enumerate(items, 1):
                raw_parts.append(
                    f"===== ITEM {i} | {item.get('created_at','')} =====\n{item.get('raw_content','')}"
                )
            z.writestr(f"{uni}/{subject}/raw_questions.txt", "\n\n".join(raw_parts))

            files = _admin_staging_file_paths(user_id, uni, subject)
            for full, rel in files:
                z.write(full, f"{uni}/{subject}/original_files/{rel}")

            meta["groups"].append({
                "uni": uni,
                "subject": subject,
                "raw_items": len(items),
                "original_files": len(files),
            })
        z.writestr("metadata.json", json.dumps(meta, ensure_ascii=False, indent=2))
    mem.seek(0)
    mem.name = f"admin_staging_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.zip"
    return mem

# =============================================================================
# KEYBOARDS & GENERAL UI
# =============================================================================



async def enqueue_admin_preload_jobs(chat_id: int, user_id: int, file_hash: str, text: str, file_name: str) -> Tuple[int, str]:
    """Queue one background preload job that creates summaries + question banks for a lecture."""
    try:
        from queue_jobs import create_job, queue_size
        job_id = create_job(
            "preload_file",
            chat_id,
            user_id,
            {"file_hash": file_hash, "text": text, "file_name": file_name, "styles": list(SUMMARY_PROMPTS.keys()), "levels": list(LEVEL_LIMITS.keys())},
        )
        return queue_size(), job_id
    except Exception:
        logger.exception("admin preload queue failed")
        return -1, ""

async def admin_extract_lecture_text(raw_bytes: bytes, file_name: str, mime_type: str, file_type: str) -> Tuple[str, str]:
    """Return (text, method). Admin preload tries local extraction, LibreOffice repair, PDF rendering + Vision, then Gemini file fallback."""
    text = ""
    method = "local"
    if file_type not in {"unsupported", "image", "zip"}:
        text = extract_text(file_type, raw_bytes)
    if text and len(text.strip()) >= 100:
        return text.strip(), method

    # Admin-only robust repair:
    # - old .ppt/.doc -> LibreOffice PDF -> text
    # - scanned/image PDF or converted slides -> render pages -> Gemini Vision
    repaired, repair_method = await admin_repair_extract_via_conversion(raw_bytes, file_name, file_type)
    if repaired and len(repaired.strip()) >= 50:
        return repaired.strip(), repair_method or "conversion_vision"

    # Direct images: read by Gemini Vision with lecture-specific prompt.
    if file_type == "image":
        img_text = await _admin_vision_extract_images([raw_bytes], file_name)
        if img_text and len(img_text.strip()) >= 50:
            return img_text.strip(), "image_vision"

    # Final fallback: ask Gemini to read the original file bytes directly when supported.
    smart = await call_gemini_admin_document_extract(raw_bytes, file_name, mime_type or guess_mime_type(file_name))
    if smart and len(smart.strip()) >= 50:
        return smart.strip(), "smart_vision"
    return (text or "").strip(), method

async def admin_preload_one(update: Update, context: ContextTypes.DEFAULT_TYPE, raw_bytes: bytes, file_name: str, mime_type: str) -> Dict[str, Any]:
    """Store one admin lecture and enqueue cache generation. No Telegram editing here; returns a report."""
    file_type = classify_document(mime_type or "", file_name or "")
    if file_type == "zip":
        return {"ok": False, "file": file_name, "reason": "zip_inside_zip"}
    if len(raw_bytes) > MAX_FILE_SIZE_BYTES:
        return {"ok": False, "file": file_name, "reason": "too_large"}
    if file_type == "unsupported":
        # Still try smart extraction for admin if the API can read the file type, but mark it as smart-only.
        text, method = await admin_extract_lecture_text(raw_bytes, file_name, mime_type, file_type)
    else:
        text, method = await admin_extract_lecture_text(raw_bytes, file_name, mime_type, file_type)
    if not text or len(text.strip()) < 50:
        return {"ok": False, "file": file_name, "reason": "no_text"}

    raw_hash = sha256_bytes(raw_bytes)
    content_hash = canonical_text_hash(text)
    if not await db_get_file(content_hash):
        await db_save_file(content_hash, file_name or "lecture", len(raw_bytes), mime_type or guess_mime_type(file_name), text)
    if not await db_get_file(raw_hash):
        await db_save_file(raw_hash, file_name or "lecture", len(raw_bytes), mime_type or guess_mime_type(file_name), text)

    qsize, job_id = await enqueue_admin_preload_jobs(
        update.effective_chat.id,
        update.effective_user.id if update.effective_user else 0,
        content_hash,
        text,
        file_name or "lecture",
    )
    if not job_id:
        # Direct fallback if Redis is unavailable.
        for style in SUMMARY_PROMPTS:
            await get_or_create_summary(content_hash, text, style)
        for level in LEVEL_LIMITS:
            await get_or_create_question_bank(content_hash, text, level)
    return {"ok": True, "file": file_name, "method": method, "job_id": job_id, "queue_size": qsize, "text_len": len(text)}

async def handle_admin_preload_file(update: Update, context: ContextTypes.DEFAULT_TYPE, raw_bytes: bytes, file_name: str, mime_type: str, status_message=None):
    """Admin-only semester lecture preload: store text and pre-generate cached outputs."""
    status = status_message or await update.message.reply_text("📚 جاري تخزين محاضرة الفصل...")
    file_type = classify_document(mime_type or "", file_name or "")
    if file_type == "zip":
        return await handle_admin_preload_zip(update, context, raw_bytes, file_name, status)
    result = await admin_preload_one(update, context, raw_bytes, file_name or "lecture", mime_type or "")
    if not result.get("ok"):
        reason = result.get("reason")
        if reason == "too_large":
            return await status.edit_text("⚠️ الملف أكبر من 30MB. قسّمه ثم أعد الإرسال.")
        if reason == "no_text":
            return await status.edit_text(
                "⚠️ لم أستطع استخراج نص واضح من هذا الملف حتى بالدعم الذكي.\n"
                "جرّب تحويله إلى PDF أو PPTX واضح ثم أعد إرساله."
            )
        return await status.edit_text("⚠️ لم أستطع معالجة هذا الملف للتخزين.")
    method_label = "استخراج عادي" if result.get("method") == "local" else "دعم ذكي Vision للأدمن"
    if result.get("job_id"):
        return await status.edit_text(
            "✅ تم تخزين الملف في مخزون الفصل.\n"
            f"📄 الملف: {file_name}\n"
            f"🧠 طريقة القراءة: {method_label}\n"
            f"📝 طول النص المستخرج: {result.get('text_len', 0)} حرف\n"
            f"⚙️ النص المستخدم للـ AI: حتى {MAX_TEXT_CHARS_FOR_AI} حرف بتوزيع بداية/وسط/نهاية.\n"
            f"📥 دخل تجهيز الملخصات والأسئلة في الخلفية.\n"
            f"📌 Job ID: {result.get('job_id')}\n"
            f"📊 الطابور الآن: {result.get('queue_size')}"
        )
    return await status.edit_text(f"✅ تم تجهيز الكاش مباشرة للملف: {file_name}\n🧠 طريقة القراءة: {method_label}")

async def handle_admin_preload_zip(update: Update, context: ContextTypes.DEFAULT_TYPE, raw_bytes: bytes, file_name: str, status_message=None):
    """Admin can upload one ZIP containing many lecture files. Each file inside is processed separately."""
    status = status_message or await update.message.reply_text("📦 جاري فتح ZIP وتخزين المحاضرات...")
    if len(raw_bytes) > ADMIN_PRELOAD_ZIP_MAX_TOTAL_BYTES:
        return await status.edit_text(f"⚠️ ملف ZIP كبير جدًا. الحد الحالي {ADMIN_PRELOAD_ZIP_MAX_TOTAL_BYTES // (1024*1024)}MB للـ ZIP كامل.")
    try:
        zf = zipfile.ZipFile(io.BytesIO(raw_bytes))
    except Exception:
        return await status.edit_text("⚠️ ملف ZIP غير صالح أو تالف.")

    infos = []
    total_uncompressed = 0
    for info in zf.infolist():
        if info.is_dir():
            continue
        name = info.filename.replace("\\", "/")
        base = os.path.basename(name)
        if not base or name.startswith("__MACOSX/") or base.startswith("."):
            continue
        if len(infos) >= ADMIN_PRELOAD_ZIP_MAX_FILES:
            break
        total_uncompressed += int(info.file_size or 0)
        if total_uncompressed > ADMIN_PRELOAD_ZIP_MAX_TOTAL_BYTES:
            return await status.edit_text(f"⚠️ محتوى ZIP بعد الفك أكبر من الحد {ADMIN_PRELOAD_ZIP_MAX_TOTAL_BYTES // (1024*1024)}MB.")
        infos.append(info)

    if not infos:
        return await status.edit_text("⚠️ لم أجد ملفات محاضرات داخل ZIP.")

    await status.edit_text(f"📦 وجدت {len(infos)} ملف داخل ZIP. سأخزنها الآن واحدة واحدة...")
    ok_count = 0
    fail_count = 0
    smart_count = 0
    reports = []
    for idx, info in enumerate(infos, 1):
        base = os.path.basename(info.filename.replace("\\", "/"))
        try:
            if int(info.file_size or 0) > MAX_FILE_SIZE_BYTES:
                fail_count += 1
                reports.append(f"❌ {base}: أكبر من 30MB")
                continue
            data = zf.read(info)
            mime = guess_mime_type(base)
            result = await admin_preload_one(update, context, data, base, mime)
            if result.get("ok"):
                ok_count += 1
                if result.get("method") == "smart_vision":
                    smart_count += 1
                reports.append(f"✅ {base}: تم التخزين ({'Vision' if result.get('method') == 'smart_vision' else 'Text'})")
            else:
                fail_count += 1
                reports.append(f"❌ {base}: {result.get('reason', 'failed')}")
        except Exception as e:
            fail_count += 1
            logger.exception("admin zip preload failed for %s", base)
            reports.append(f"❌ {base}: خطأ أثناء المعالجة")
        if idx % 3 == 0 or idx == len(infos):
            try:
                await status.edit_text(f"📦 جاري التخزين... {idx}/{len(infos)}\n✅ نجح: {ok_count}\n❌ فشل: {fail_count}")
            except Exception:
                pass

    report_text = "\n".join(reports[-20:])
    return await status.edit_text(
        "✅ انتهى تخزين محاضرات ZIP.\n\n"
        f"📦 إجمالي الملفات: {len(infos)}\n"
        f"✅ تم تخزينها: {ok_count}\n"
        f"🧠 قرئت بالدعم الذكي Vision: {smart_count}\n"
        f"❌ فشلت: {fail_count}\n\n"
        "آخر النتائج:\n" + report_text[:2500]
    )

def lang_keyboard() -> InlineKeyboardMarkup: return InlineKeyboardMarkup([[InlineKeyboardButton("🇸🇦 العربية", callback_data="lang:ar"), InlineKeyboardButton("🇬🇧 English", callback_data="lang:en")]])
def mode_keyboard(user_data: dict) -> InlineKeyboardMarkup:
    rows = [
        [InlineKeyboardButton(t(user_data, "mode_quiz"), callback_data="mode:quiz"), InlineKeyboardButton(t(user_data, "mode_summary"), callback_data="mode:summary")],
        [InlineKeyboardButton(t(user_data, "mode_osce"), callback_data="mode:osce"), InlineKeyboardButton(t(user_data, "mode_past_papers"), callback_data="mode:past_papers")],
        [InlineKeyboardButton(t(user_data, "mode_progress"), callback_data="panel:progress"), InlineKeyboardButton(t(user_data, "mode_leaderboard"), callback_data="panel:leaderboard")],
        [InlineKeyboardButton(t(user_data, "mode_subscription"), callback_data="panel:subscription"), InlineKeyboardButton("📞 الدعم", callback_data="panel:support")],
    ]
    if user_data.get("is_admin"):
        rows.append([InlineKeyboardButton("🛠️ لوحة الأدمن", callback_data="admin:home")])
    return InlineKeyboardMarkup(rows)
def summary_style_keyboard(user_data: dict) -> InlineKeyboardMarkup: return InlineKeyboardMarkup([[InlineKeyboardButton(t(user_data, "style_disease"), callback_data="style:disease_v2")], [InlineKeyboardButton(t(user_data, "style_highyield"), callback_data="style:highyield")], [InlineKeyboardButton(t(user_data, "style_osce"), callback_data="style:osce")], [InlineKeyboardButton("⬅️ رجوع", callback_data="back:mode")]])
def level_keyboard(setup_id: str, user_data: dict) -> InlineKeyboardMarkup: return InlineKeyboardMarkup([[InlineKeyboardButton(t(user_data, "level_basic"), callback_data=f"level:{setup_id}:basic")], [InlineKeyboardButton(t(user_data, "level_cases"), callback_data=f"level:{setup_id}:cases")], [InlineKeyboardButton(t(user_data, "level_challenge"), callback_data=f"level:{setup_id}:challenge")], [InlineKeyboardButton("⬅️ رجوع", callback_data="back:mode")]])
def basic_count_keyboard(setup_id: str) -> InlineKeyboardMarkup: return InlineKeyboardMarkup([[InlineKeyboardButton("5", callback_data=f"count:{setup_id}:5"), InlineKeyboardButton("10", callback_data=f"count:{setup_id}:10"), InlineKeyboardButton("20", callback_data=f"count:{setup_id}:20")], [InlineKeyboardButton("30", callback_data=f"count:{setup_id}:30"), InlineKeyboardButton("40", callback_data=f"count:{setup_id}:40")], [InlineKeyboardButton("⬅️ رجوع", callback_data=f"back:levels:{setup_id}")]])
def answer_keyboard(quiz_id: str, idx: int) -> InlineKeyboardMarkup: return InlineKeyboardMarkup([[InlineKeyboardButton(f" {l} ", callback_data=f"ans:{quiz_id}:{idx}:{l}") for l in "ABCDE"]])
def post_quiz_keyboard(user_data: dict, has_wrong: bool) -> InlineKeyboardMarkup:
    rows = []
    if has_wrong: rows.append([InlineKeyboardButton(t(user_data, "review"), callback_data="review:start")])
    rows.extend([[InlineKeyboardButton(t(user_data, "stats"), callback_data="review:stats")], [InlineKeyboardButton(t(user_data, "new_quiz"), callback_data="review:new")]])
    return InlineKeyboardMarkup(rows)
def review_keyboard(current: int, total: int) -> InlineKeyboardMarkup:
    buttons = []
    if current < total - 1: buttons.append(InlineKeyboardButton("Next ▶", callback_data="review:next"))
    buttons.append(InlineKeyboardButton("Done ✖", callback_data="review:exit"))
    return InlineKeyboardMarkup([buttons])

def main_menu_text(user_data: dict) -> str:
    lang = get_lang(user_data)
    if lang == "ar": return "🩺 MedMCQ AI Academy\n━━━━━━━━━━━━━━\nمنصتك الذكية للتدريب الطبي\n\nاختر الخدمة التي تريدها:\n📝 أسئلة تفاعلية من محاضراتك\n📚 ملخصات طبية\n🎭 محاكي OSCE\n📂 أسئلة سنوات سابقة\n📊 تقدم ونقاط\n🏆 لوحة الصدارة\n💎 الاشتراك والدفع"
    return "🩺 MedMCQ AI Academy\n━━━━━━━━━━━━━━\nYour smart medical learning platform.\n\nChoose a service:\n📝 Interactive quizzes\n📚 AI summaries\n🎭 Smart OSCE\n📂 Past papers\n📊 XP, levels\n🏆 Leaderboard\n💎 Subscription"

def format_question(q: dict, idx: int, total: int, user_data: dict) -> str:
    nq = normalize_questions([q])
    q = nq[0] if nq else {"question": "", "options": {}, "correct": "", "explanation": ""}
    opts = _normalize_options(q.get("options", {}))
    for letter in "ABCDE":
        opts.setdefault(letter, "")
    return (
        f"{progress_bar(idx, total)}\n{t(user_data, 'q_meta', i=idx, n=total)}\n\n"
        f"{escape_md(q.get('question', ''))}\n\n"
        f"A. {escape_md(opts.get('A', ''))}\n"
        f"B. {escape_md(opts.get('B', ''))}\n"
        f"C. {escape_md(opts.get('C', ''))}\n"
        f"D. {escape_md(opts.get('D', ''))}\n"
        f"E. {escape_md(opts.get('E', ''))}"
    )

def format_review(entry: dict, idx: int, total: int, user_data: dict) -> str:
    q = entry.get("question", {}) if isinstance(entry, dict) else {}
    chosen = entry.get("chosen", "") if isinstance(entry, dict) else ""
    nq = normalize_questions([q])
    q = nq[0] if nq else {"question": "", "options": {}, "correct": "", "explanation": ""}
    opts = _normalize_options(q.get("options", {}))
    correct = q.get("correct", "")
    correct_text = opts.get(correct, "")
    return (
        f"📋 Review {idx}/{total}\n\n{escape_md(q.get('question', ''))}\n\n"
        f"A. {escape_md(opts.get('A', ''))}\n"
        f"B. {escape_md(opts.get('B', ''))}\n"
        f"C. {escape_md(opts.get('C', ''))}\n"
        f"D. {escape_md(opts.get('D', ''))}\n"
        f"E. {escape_md(opts.get('E', ''))}\n\n"
        f"Your answer: {chosen}\n"
        f"Correct answer: {correct}. {escape_md(correct_text)}\n\n"
        f"{t(user_data, 'explanation')} {escape_md(q.get('explanation', ''))}"
    )

# =============================================================================
# SENDERS & COMMANDS
# =============================================================================

async def safe_send(bot, chat_id: int, text: str, **kwargs) -> bool:
    for _ in range(2):
        try: await bot.send_message(chat_id=chat_id, text=text, **kwargs); return True
        except Exception as e: logger.warning("send failed: %s", e); await asyncio.sleep(0.5)
    return False

async def send_current_question(chat_id: int, context: ContextTypes.DEFAULT_TYPE) -> None:
    session = context.user_data.get("session")
    if not isinstance(session, dict):
        return
    idx = int(session.get("current", 0))
    questions = normalize_questions(session.get("questions", []))
    session["questions"] = questions
    if idx >= len(questions):
        await send_final_score(chat_id, context)
        return
    await safe_send(
        context.bot,
        chat_id,
        format_question(questions[idx], idx + 1, len(questions), context.user_data),
        parse_mode="Markdown",
        reply_markup=answer_keyboard(session["quiz_id"], idx),
    )

async def send_final_score(chat_id: int, context: ContextTypes.DEFAULT_TYPE) -> None:
    session = context.user_data.get("session")
    if not isinstance(session, dict):
        return
    questions = normalize_questions(session.get("questions", []))
    session["questions"] = questions
    total = len(questions)
    score = int(session.get("score", 0))
    wrong = session.get("wrong", [])
    update_stats(context.user_data, score, total)
    xp_points = 25 + (score * 10) + (20 if total and (score / total) >= 0.8 else 0)
    progress = await db_add_xp(chat_id, xp_points, "Quiz completed", correct=score, total=total, quiz_done=True)
    text = t(context.user_data, "complete", score=score, total=total, pct=round(score / total * 100) if total else 0) + f"\n\n⭐ +{xp_points} XP | 🎖 Level {progress.get('level', 1)} | 🔥 Streak {progress.get('streak', 0)}"
    if progress.get("leveled_up"):
        text += "\n🎉 Level up! ممتاز، استمر."
    await context.bot.send_message(chat_id=chat_id, text=text, reply_markup=post_quiz_keyboard(context.user_data, bool(wrong)))
    context.user_data["review_pool"] = wrong
    context.user_data.pop("session", None)


async def send_welcome_gif(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Send welcome GIF if WELCOME_GIF is set in .env.
    WELCOME_GIF can be a Telegram file_id, direct URL, or local path such as welcome.gif.
    """
    gif = (WELCOME_GIF or "").strip()
    if not gif:
        return
    try:
        if os.path.exists(gif):
            with open(gif, "rb") as f:
                await update.message.reply_animation(animation=f)
        else:
            await update.message.reply_animation(animation=gif)
    except Exception as e:
        logger.warning("Welcome GIF failed: %s", e)

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await db_register_user(update.effective_user)
    await send_welcome_gif(update, context)
    await update.message.reply_text("🩺 MedMCQ AI Academy\n━━━━━━━━━━━━━━\n🎓 Smart Medical Learning Platform\n\n📚 AI Summaries\n🧠 Smart Quizzes\n🎭 Interactive OSCE Cases\n📂 Past Papers\n🔥 Streak, XP & Levels\n🏆 Leaderboard\n💎 Subscription\n\nاختر لغتك للبدء", reply_markup=lang_keyboard())

async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE): await update.message.reply_text(t(context.user_data, "help"))
async def language_command(update: Update, context: ContextTypes.DEFAULT_TYPE): await update.message.reply_text(t(context.user_data, "choose_lang"), reply_markup=lang_keyboard())

async def stop_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    had = any(k in context.user_data for k in ("session", "pending_file", "review", "osce_session", "osce_v4_session"))
    for k in ("session", "pending_file", "pending_mode", "pending_summary_style", "review", "review_pool", "busy", "osce_session", "osce_v4_session"): context.user_data.pop(k, None)
    await update.message.reply_text(t(context.user_data, "stopped" if had else "no_active"))

async def stats_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    stats = get_stats(context.user_data); total = stats["total"]
    await update.message.reply_text(t(context.user_data, "stats_text", quizzes=stats["quizzes"], correct=stats["correct"], total=total, accuracy=round(stats["correct"] / total * 100) if total else 0, best=stats["best"]))

async def review_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    pool = context.user_data.get("review_pool", [])
    if not pool: await update.message.reply_text(t(context.user_data, "review_none")); return
    context.user_data["review"] = {"pool": pool, "current": 0}
    await send_review(update.effective_chat.id, context)

async def support_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text("📞 مركز الدعم\n\nاختر نوع الدعم الذي تحتاجه:", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("💳 دعم عبر رصيد المدار", callback_data="support:balance")], [InlineKeyboardButton("💬 تواصل مع الدعم", callback_data="support:contact")], [InlineKeyboardButton("📢 أبلغ عن مشكلة", callback_data="support:problem")], [InlineKeyboardButton("⬅️ رجوع", callback_data="back:mode")]]))

async def osce_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data["pending_mode"] = "osce"
    await update.message.reply_text(t(context.user_data, "mode_osce_prompt"))

async def testbot_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not update.effective_user or update.effective_user.id not in ADMIN_USER_IDS: return await update.message.reply_text("Unauthorized.")
    try:
        provider = "Google direct" if GOOGLE_API_KEYS else "OpenRouter"
        await update.message.reply_text(f"✅ AI test ({provider} first): {await call_gemini('Reply with exactly: OK', temperature=0.0, max_output_tokens=10)}")
    except Exception as e:
        await update.message.reply_text(f"❌ Failed: {e}")

# =============================================================================
# CALLBACKS — STABLE ROUTER + PAST PAPERS + ADMIN PANEL
# =============================================================================

PAST_SUBJECTS = [
    ("forensic", "⚖️ طب شرعي (Forensic)"),
    ("obgyn", "🤰 نساء وولادة (Ob/Gyn)"),
    ("ophthalmology", "👁️ عيون (Ophthalmology)"),
    ("radiology", "🩻 أشعة (Radiology)"),
    ("community", "👥 مجتمع وأسرة (Community)"),
]
SUBJECT_LABELS = dict(PAST_SUBJECTS)
UNI_LABELS = {"ajdabiya": "جامعة أجدابيا", "benghazi": "جامعة بنغازي"}


def _safe_callback_data(data: str) -> str:
    """Normalize old and new callback codes so old Telegram messages keep working."""
    if not data:
        return ""
    replacements = {
        "pp_uni:": "pp:uni:",
        "pp-uni:": "pp:uni:",
        "pp-uni2:": "pp:uni:",
        "pp_uni2:": "pp:uni:",
        "pp_sub:": "pp:subject:",
        "pp-sub:": "pp:subject:",
        "pp_menu:": "pp:subject:",
        "pp-menu:": "pp:subject:",
        "pp_mode:": "pp:mode:",
        "pp-mode:": "pp:mode:",
        "pp_quiz_start:": "pp:quiz:",
        "pp-quiz-start:": "pp:quiz:",
        "pp_quiz_next:": "pp:quiznext:",
        "pp-quiz-next:": "pp:quiznext:",
        "pp_wrong:": "pp:wrong:",
        "pp-wrong:": "pp:wrong:",
        "pp_progress:": "pp:progress:",
        "pp-progress:": "pp:progress:",
        "pp_ans:": "pp:ans:",
        "pp-ans:": "pp:ans:",
        "pp-past:home": "pp:home",
    }
    for old, new in replacements.items():
        if data.startswith(old):
            return data.replace(old, new, 1)
    if data == "mode:past_papers" or data == "mode:past-papers":
        return "pp:home"
    if data == "back:mode" or data == "back:main":
        return "main:home"
    if data == "back:support":
        return "panel:support"
    return data

async def safe_edit(query, text: str, reply_markup=None, **kwargs):
    try:
        return await query.edit_message_text(text, reply_markup=reply_markup, **kwargs)
    except Exception as e:
        logger.warning("edit_message_text failed; sending new message: %s", e)
        try:
            return await query.message.reply_text(text, reply_markup=reply_markup, **kwargs)
        except Exception:
            logger.exception("safe_edit failed completely")
            return None

async def safe_answer_query(query, text: str = None, show_alert: bool = False):
    try:
        await query.answer(text=text, show_alert=show_alert)
    except Exception:
        pass

def past_home_keyboard() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("🏛️ جامعة أجدابيا", callback_data="pp:uni:ajdabiya")],
        [InlineKeyboardButton("🏛️ جامعة بنغازي", callback_data="pp:uni:benghazi")],
        [InlineKeyboardButton("⬅️ رجوع للرئيسية", callback_data="main:home")],
    ])

def past_subjects_keyboard(uni: str) -> InlineKeyboardMarkup:
    rows = [[InlineKeyboardButton(label, callback_data=f"pp:subject:{uni}:{sub}")] for sub, label in PAST_SUBJECTS]
    rows.append([InlineKeyboardButton("⬅️ رجوع للجامعات", callback_data="pp:home")])
    return InlineKeyboardMarkup(rows)

def pp_subject_menu_keyboard(uni: str, subject: str, has_wrong: bool = False) -> InlineKeyboardMarkup:
    rows = [
        [InlineKeyboardButton("🎯 تدريب عشوائي", callback_data=f"pp:mode:random:{uni}:{subject}")],
        [InlineKeyboardButton("📚 حل بالترتيب", callback_data=f"pp:mode:seq:{uni}:{subject}")],
        [InlineKeyboardButton("📝 اختبار 20 سؤال", callback_data=f"pp:quiz:{uni}:{subject}")],
        [InlineKeyboardButton("❌ راجع أخطائي", callback_data=f"pp:wrong:{uni}:{subject}")],
        [InlineKeyboardButton("📊 تقدمي", callback_data=f"pp:progress:{uni}:{subject}")],
        [InlineKeyboardButton("⬅️ رجوع للمواد", callback_data=f"pp:uni:{uni}")],
    ]
    return InlineKeyboardMarkup(rows)


_PP_ARABIC_DIACRITICS_RE = re.compile(r"[\u064B-\u065F\u0670\u06D6-\u06ED]")
_PP_NON_WORD_RE = re.compile(r"[^0-9a-zA-Z\u0600-\u06FF]+")

def _pp_clean_text(value: object) -> str:
    """Normalize question text so duplicated past-paper questions are detected."""
    text = "" if value is None else str(value)
    text = text.lower()
    text = _PP_ARABIC_DIACRITICS_RE.sub("", text)
    # normalize common Arabic letter variants
    text = text.replace("أ", "ا").replace("إ", "ا").replace("آ", "ا")
    text = text.replace("ى", "ي").replace("ؤ", "و").replace("ئ", "ي").replace("ة", "ه")
    text = _PP_NON_WORD_RE.sub(" ", text)
    return re.sub(r"\s+", " ", text).strip()

def _pp_question_fingerprint(q: dict) -> str:
    """Stable fingerprint for duplicate detection, tolerant to spaces/punctuation/option order."""
    if not isinstance(q, dict):
        return ""
    question = _pp_clean_text(q.get("question", ""))
    options = q.get("options", {})
    option_values = []
    if isinstance(options, dict):
        option_values = [_pp_clean_text(v) for v in options.values()]
    elif isinstance(options, list):
        option_values = [_pp_clean_text(v) for v in options]
    option_values = sorted([v for v in option_values if v])
    base = question + "||" + "|".join(option_values)
    return hashlib.sha256(base.encode("utf-8")).hexdigest()

def _pp_normalize_question(q: dict) -> Optional[dict]:
    if not isinstance(q, dict):
        return None
    question = str(q.get("question") or q.get("q") or q.get("stem") or "").strip()
    options = q.get("options") or q.get("choices") or {}
    answer = str(q.get("correct") or q.get("correct_answer") or q.get("answer") or q.get("correctAnswer") or "").strip()
    explanation = str(q.get("explanation") or q.get("explain") or "").strip()
    if isinstance(options, list):
        letters = list("ABCDE")
        options = {letters[i]: str(v) for i, v in enumerate(options[:5])}
    if not isinstance(options, dict):
        return None
    norm_opts = {}
    for l in "ABCDE":
        val = options.get(l) or options.get(l.lower())
        norm_opts[l] = str(val or "").strip()
    if not question or not any(norm_opts.values()):
        return None
    if not norm_opts["E"]:
        norm_opts["E"] = "None of the above"
    answer = answer.upper().strip()
    if answer not in list("ABCDE"):
        for l, val in norm_opts.items():
            if answer and answer.lower() == val.lower():
                answer = l
                break
    if answer not in list("ABCDE"):
        answer = "A"
    return {"question": question, "options": norm_opts, "correct": answer, "explanation": explanation or "من أسئلة السنوات السابقة."}

async def pp_load_questions(uni: str, subject: str) -> List[dict]:
    def work():
        conn = _db_connect()
        rows = conn.execute("SELECT questions_json FROM past_papers WHERE lower(uni)=lower(?) AND lower(subject)=lower(?) ORDER BY id ASC", (uni, subject)).fetchall()
        conn.close()
        out = []
        seen = set()
        for row in rows:
            try:
                arr = json.loads(row["questions_json"])
                if not isinstance(arr, list):
                    continue
                for item in arr:
                    nq = _pp_normalize_question(item)
                    if not nq:
                        continue
                    fp = _pp_question_fingerprint(nq)
                    if fp in seen:
                        continue
                    seen.add(fp)
                    out.append(nq)
            except Exception:
                continue
        return out
    return await asyncio.to_thread(work)

async def pp_filter_new_questions(uni: str, subject: str, incoming: list) -> Tuple[List[dict], int, int]:
    """Return (new_questions, duplicate_count, invalid_count) for an upload_final JSON list."""
    if not isinstance(incoming, list):
        return [], 0, 1

    existing_questions = await pp_load_questions(uni, subject)
    existing_fps = {_pp_question_fingerprint(q) for q in existing_questions if q}
    batch_fps = set()

    new_questions: List[dict] = []
    duplicate_count = 0
    invalid_count = 0

    for item in incoming:
        nq = _pp_normalize_question(item)
        if not nq:
            invalid_count += 1
            continue

        fp = _pp_question_fingerprint(nq)
        if not fp:
            invalid_count += 1
            continue

        if fp in existing_fps or fp in batch_fps:
            duplicate_count += 1
            continue

        batch_fps.add(fp)
        new_questions.append(nq)

    return new_questions, duplicate_count, invalid_count

async def pp_admin_stats_text() -> str:
    def work():
        conn = _db_connect()
        rows = conn.execute("SELECT uni, subject, questions_json FROM past_papers ORDER BY uni, subject, id").fetchall()
        conn.close()
        buckets: Dict[Tuple[str, str], set] = {}
        raw_counts: Dict[Tuple[str, str], int] = {}
        for row in rows:
            key = (str(row["uni"]).lower(), str(row["subject"]).lower())
            buckets.setdefault(key, set())
            raw_counts.setdefault(key, 0)
            try:
                arr = json.loads(row["questions_json"] or "[]")
            except Exception:
                arr = []
            if not isinstance(arr, list):
                continue
            raw_counts[key] += len(arr)
            for item in arr:
                nq = _pp_normalize_question(item)
                if nq:
                    buckets[key].add(_pp_question_fingerprint(nq))
        return [(k[0], k[1], len(v), raw_counts.get(k, 0)) for k, v in buckets.items()]

    rows = await asyncio.to_thread(work)
    if not rows:
        return "📚 بنك السنوات السابقة فارغ حالياً."
    lines = ["📊 إحصائيات بنك السنوات السابقة\n"]
    total_unique = 0
    total_raw = 0
    for uni, subject, unique_count, raw_count in rows:
        total_unique += unique_count
        total_raw += raw_count
        dup = max(0, raw_count - unique_count)
        lines.append(
            f"🏛️ {UNI_LABELS.get(uni, uni)} | 📚 {SUBJECT_LABELS.get(subject, subject)}\n"
            f"   ✅ فريدة: {unique_count} | ♻️ مكررة داخل المخزن: {dup} | خام: {raw_count}"
        )
    lines.append(f"\nالمجموع الفريد: {total_unique}\nالمجموع الخام: {total_raw}")
    return "\n".join(lines)

async def past_stats_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not update.effective_user or update.effective_user.id not in ADMIN_USER_IDS:
        return await update.message.reply_text("⛔ هذا الأمر للإدارة فقط.")
    text = await pp_admin_stats_text()
    for chunk in [text[i:i + 3800] for i in range(0, len(text), 3800)]:
        await update.message.reply_text(chunk)

async def pp_get_progress(user_id: int, uni: str, subject: str) -> dict:
    def work():
        conn = _db_connect()
        row = conn.execute("SELECT * FROM past_progress WHERE user_id=? AND uni=? AND subject=?", (user_id, uni, subject)).fetchone()
        conn.close()
        if not row:
            return {"attempts": 0, "correct": 0, "wrong": 0, "wrong_json": "[]"}
        return dict(row)
    return await asyncio.to_thread(work)

async def pp_record_answer(user_id: int, uni: str, subject: str, question: dict, is_correct: bool):
    fp = _pp_question_fingerprint(question)
    async with _db_lock:
        def work():
            conn = _db_connect()
            row = conn.execute("SELECT * FROM past_progress WHERE user_id=? AND uni=? AND subject=?", (user_id, uni, subject)).fetchone()
            if row:
                attempts = int(row["attempts"] or 0)
                correct = int(row["correct"] or 0)
                wrong = int(row["wrong"] or 0)
                try:
                    wrong_list = json.loads(row["wrong_json"] or "[]")
                except Exception:
                    wrong_list = []
            else:
                attempts = correct = wrong = 0
                wrong_list = []
            attempts += 1
            if is_correct:
                correct += 1
                wrong_list = [x for x in wrong_list if x.get("fp") != fp]
            else:
                wrong += 1
                item = dict(question)
                item["fp"] = fp
                wrong_list = [x for x in wrong_list if x.get("fp") != fp]
                wrong_list.append(item)
                wrong_list = wrong_list[-200:]
            conn.execute(
                "INSERT OR REPLACE INTO past_progress(user_id, uni, subject, attempts, correct, wrong, wrong_json, updated_at) VALUES(?,?,?,?,?,?,?,?)",
                (user_id, uni, subject, attempts, correct, wrong, json.dumps(wrong_list, ensure_ascii=False), datetime.utcnow().isoformat()),
            )
            conn.commit()
            conn.close()
        await asyncio.to_thread(work)

async def show_main_menu(query, context: ContextTypes.DEFAULT_TYPE):
    context.user_data["is_admin"] = is_admin_user(query.from_user.id if query.from_user else None)
    await safe_edit(query, main_menu_text(context.user_data), reply_markup=mode_keyboard(context.user_data))

async def lang_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    lang = query.data.split(":", 1)[1]
    if lang in SUPPORTED_LANGS:
        context.user_data["lang"] = lang
    context.user_data["is_admin"] = is_admin_user(query.from_user.id if query.from_user else None)
    await show_main_menu(query, context)

async def mode_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    data = _safe_callback_data(query.data or "")
    if data == "pp:home":
        return await pp_home_callback(update, context)
    mode = data.split(":", 1)[1] if ":" in data else ""
    if mode not in ("quiz", "summary", "osce"):
        return await show_main_menu(query, context)
    context.user_data["pending_mode"] = mode
    try:
        await query.edit_message_reply_markup(reply_markup=None)
    except Exception:
        pass
    if mode == "summary":
        await context.bot.send_message(update.effective_chat.id, t(context.user_data, "summary_style_prompt") + "\n\n" + ai_queue_status_text(), reply_markup=summary_style_keyboard(context.user_data))
    elif mode == "osce":
        await context.bot.send_message(update.effective_chat.id, t(context.user_data, "mode_osce_prompt") + "\n\n" + ai_queue_status_text())
    else:
        await context.bot.send_message(update.effective_chat.id, t(context.user_data, "mode_quiz_prompt") + "\n\n" + ai_queue_status_text())

async def pp_home_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    await safe_edit(query, "📂 بنك الامتحانات (السنوات السابقة)\n\nاختر الجامعة:", reply_markup=past_home_keyboard())

async def pp_uni_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    data = _safe_callback_data(query.data or "")
    parts = data.split(":")
    uni = parts[2] if len(parts) >= 3 else "ajdabiya"
    title = UNI_LABELS.get(uni, uni)
    await safe_edit(query, f"🏛️ {title}\n\nاختر المادة:", reply_markup=past_subjects_keyboard(uni))

async def pp_sub_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    data = _safe_callback_data(query.data or "")
    parts = data.split(":")
    if len(parts) < 4:
        return await pp_home_callback(update, context)
    uni, subject = parts[2], parts[3]
    questions = await pp_load_questions(uni, subject)
    user_id = query.from_user.id if query.from_user else update.effective_chat.id
    prog = await pp_get_progress(user_id, uni, subject)
    attempts = int(prog.get("attempts") or 0)
    correct = int(prog.get("correct") or 0)
    wrong = int(prog.get("wrong") or 0)
    pct = round(correct / attempts * 100) if attempts else 0
    text = (
        f"🏛️ {UNI_LABELS.get(uni, uni)}\n"
        f"📚 {SUBJECT_LABELS.get(subject, subject)}\n\n"
        f"📊 عدد الأسئلة المحفوظة: {len(questions)}\n"
        f"✅ الصحيحة: {correct}\n"
        f"❌ الخاطئة: {wrong}\n"
        f"📈 النسبة: {pct}%\n\n"
        f"اختر طريقة التدريب:"
    )
    await safe_edit(query, text, reply_markup=pp_subject_menu_keyboard(uni, subject))

async def pp_start_session(update: Update, context: ContextTypes.DEFAULT_TYPE, uni: str, subject: str, mode: str):
    query = update.callback_query
    questions = await pp_load_questions(uni, subject)
    if not questions:
        return await safe_edit(query, f"🚧 لا توجد أسئلة محفوظة حالياً في مادة:\n{SUBJECT_LABELS.get(subject, subject)}\n\nارفع أسئلة من لوحة الأدمن ثم جرّب مرة ثانية.", reply_markup=pp_subject_menu_keyboard(uni, subject))
    if mode == "random":
        selected = [random.choice(questions)]
    elif mode == "seq":
        selected = questions
    elif mode == "quiz":
        selected = random.sample(questions, min(20, len(questions)))
    else:
        selected = questions
    context.user_data["session"] = {
        "quiz_id": str(uuid.uuid4())[:8],
        "questions": [shuffle_options(q) for q in selected],
        "current": 0,
        "score": 0,
        "wrong": [],
        "level": f"past:{uni}:{subject}:{mode}",
    }
    await safe_edit(query, f"✅ بدأ التدريب\n🏛️ {UNI_LABELS.get(uni, uni)}\n📚 {SUBJECT_LABELS.get(subject, subject)}\n📝 عدد الأسئلة: {len(selected)}")
    await asyncio.sleep(0.3)
    await send_current_question(update.effective_chat.id, context)

async def pp_mode_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    data = _safe_callback_data(query.data or "")
    parts = data.split(":")
    if len(parts) < 5:
        return await pp_home_callback(update, context)
    mode, uni, subject = parts[2], parts[3], parts[4]
    await pp_start_session(update, context, uni, subject, mode)

async def pp_quiz_start_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    data = _safe_callback_data(query.data or "")
    parts = data.split(":")
    if len(parts) < 4:
        return await pp_home_callback(update, context)
    await pp_start_session(update, context, parts[2], parts[3], "quiz")

async def pp_wrong_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    data = _safe_callback_data(query.data or "")
    parts = data.split(":")
    if len(parts) < 4:
        return await pp_home_callback(update, context)
    uni, subject = parts[2], parts[3]
    user_id = query.from_user.id if query.from_user else update.effective_chat.id
    prog = await pp_get_progress(user_id, uni, subject)
    try:
        wrong_list = json.loads(prog.get("wrong_json") or "[]")
    except Exception:
        wrong_list = []
    wrong_list = [_pp_normalize_question(q) for q in wrong_list]
    wrong_list = [q for q in wrong_list if q]
    if not wrong_list:
        return await safe_edit(query, "✅ لا توجد أسئلة خاطئة محفوظة للمراجعة في هذه المادة.", reply_markup=pp_subject_menu_keyboard(uni, subject))
    context.user_data["session"] = {
        "quiz_id": str(uuid.uuid4())[:8],
        "questions": [shuffle_options(q) for q in wrong_list],
        "current": 0,
        "score": 0,
        "wrong": [],
        "level": f"past:{uni}:{subject}:wrong_review",
    }
    await safe_edit(query, f"❌ مراجعة الأخطاء\nعدد الأسئلة: {len(wrong_list)}")
    await asyncio.sleep(0.3)
    await send_current_question(update.effective_chat.id, context)

async def pp_progress_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    data = _safe_callback_data(query.data or "")
    parts = data.split(":")
    if len(parts) < 4:
        return await pp_home_callback(update, context)
    uni, subject = parts[2], parts[3]
    user_id = query.from_user.id if query.from_user else update.effective_chat.id
    prog = await pp_get_progress(user_id, uni, subject)
    attempts = int(prog.get("attempts") or 0)
    correct = int(prog.get("correct") or 0)
    wrong = int(prog.get("wrong") or 0)
    pct = round(correct / attempts * 100) if attempts else 0
    await safe_edit(query, f"📊 تقدمي في المادة\n\n🏛️ {UNI_LABELS.get(uni, uni)}\n📚 {SUBJECT_LABELS.get(subject, subject)}\n\n📝 محاولات الحل: {attempts}\n✅ الصحيحة: {correct}\n❌ الخاطئة: {wrong}\n📈 النسبة: {pct}%", reply_markup=pp_subject_menu_keyboard(uni, subject))

# Placeholder for old callback compatibility
async def pp_answer_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await answer_callback(update, context)

async def pp_quiz_next_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await pp_quiz_start_callback(update, context)

async def summary_style_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query; await safe_answer_query(query); style = query.data.split(":", 1)[1]
    if style in SUMMARY_PROMPTS:
        context.user_data.update({"pending_summary_style": style, "pending_mode": "summary"})
        try: await query.edit_message_reply_markup(reply_markup=None)
        except Exception: pass
        await context.bot.send_message(update.effective_chat.id, t(context.user_data, "mode_summary_prompt"))

async def level_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    parts = query.data.split(":")
    if len(parts) != 3:
        return

    setup_id, level = parts[1], parts[2]
    pending = context.user_data.get("pending_file")
    if not isinstance(pending, dict) or pending.get("setup_id") != setup_id:
        if update.effective_user:
            pending = load_pending_quiz_from_redis(update.effective_user.id, setup_id)
            if isinstance(pending, dict):
                context.user_data["pending_file"] = pending

    if not isinstance(pending, dict) or pending.get("setup_id") != setup_id:
        return await safe_edit(query, t(context.user_data, "stale"))

    if level not in LEVEL_LIMITS:
        return

    cached_for_level = await db_get_questions(pending["file_hash"], level)
    needs_ai_for_level = cached_for_level is None
    if needs_ai_for_level and not await check_daily_limit(update, context, level):
        return

    pending["level"] = level
    context.user_data["_needs_ai_limit_for_level"] = level if needs_ai_for_level else ""

    if level == "basic":
        await safe_edit(query, t(context.user_data, "basic_count_prompt"), reply_markup=basic_count_keyboard(setup_id))
        return

    # إذا كان بنك الأسئلة غير موجود، أدخله للطابور. لو Redis غير متاح، يرجع للتوليد المباشر.
    if USE_QUEUE:
        try:
            cached = await db_get_questions(pending["file_hash"], level)
            if cached is None:
                from queue_jobs import create_job, queue_size, save_pending_quiz
                save_pending_quiz(update.effective_user.id if update.effective_user else 0, setup_id, pending)
                job_id = create_job(
                    "question_bank",
                    update.effective_chat.id,
                    update.effective_user.id if update.effective_user else 0,
                    {
                        "file_hash": pending["file_hash"],
                        "text": pending["text"],
                        "level": level,
                        "setup_id": setup_id,
                    },
                )
                await increment_daily_usage(context.user_data, level, update.effective_user.id if update.effective_user else 0)
                context.user_data["_needs_ai_limit_for_level"] = ""
                return await safe_edit(
                    query,
                    "✅ تم استلام طلب توليد الأسئلة.\n"
                    "📥 دخل الطلب في الطابور.\n"
                    f"📌 رقم الطلب: {job_id}\n"
                    f"📊 عدد الطلبات في الطابور: {queue_size()}\n\n"
                    "سأرسل لك زر بدء الاختبار عندما يصبح بنك الأسئلة جاهزاً."
                )
        except Exception:
            logger.exception("queue question_bank failed; falling back to direct generation")

    await safe_edit(query, t(context.user_data, "generating_pool"))
    await start_quiz_from_bank(update, context, setup_id, level, LEVEL_LIMITS[level])




def load_pending_quiz_from_redis(user_id: int, setup_id: str):
    try:
        from queue_jobs import load_pending_quiz
        return load_pending_quiz(user_id, setup_id)
    except Exception:
        return None


async def count_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query; await safe_answer_query(query); parts = query.data.split(":")
    if len(parts) != 3: return
    try: count = int(parts[2])
    except Exception: return
    if count not in BASIC_COUNTS: return
    pending = context.user_data.get("pending_file")
    if not isinstance(pending, dict):
        pending = None
    if (not pending or pending.get("setup_id") != parts[1]) and update.effective_user:
        pending = load_pending_quiz_from_redis(update.effective_user.id, parts[1])
        if isinstance(pending, dict):
            context.user_data["pending_file"] = pending
    if not isinstance(pending, dict) or pending.get("setup_id") != parts[1]:
        return await safe_edit(query, t(context.user_data, "stale"))

    cached_for_basic = await db_get_questions(pending["file_hash"], "basic")
    needs_ai_for_basic = cached_for_basic is None
    if needs_ai_for_basic and not await check_daily_limit(update, context, "basic"):
        return
    context.user_data["_needs_ai_limit_for_level"] = "basic" if needs_ai_for_basic else ""

    if USE_QUEUE:
        try:
            cached = await db_get_questions(pending["file_hash"], "basic")
            if cached is None:
                from queue_jobs import create_job, queue_size, save_pending_quiz
                save_pending_quiz(update.effective_user.id if update.effective_user else 0, parts[1], pending)
                job_id = create_job(
                    "question_bank",
                    update.effective_chat.id,
                    update.effective_user.id if update.effective_user else 0,
                    {
                        "file_hash": pending["file_hash"],
                        "text": pending["text"],
                        "level": "basic",
                        "setup_id": parts[1],
                        "requested_count": count,
                    },
                )
                await increment_daily_usage(context.user_data, "basic", update.effective_user.id if update.effective_user else 0)
                context.user_data["_needs_ai_limit_for_level"] = ""
                return await safe_edit(
                    query,
                    "✅ تم استلام طلب توليد أسئلة Basic.\n"
                    "📥 دخل الطلب في الطابور.\n"
                    f"📌 رقم الطلب: {job_id}\n"
                    f"📊 عدد الطلبات في الطابور: {queue_size()}\n\n"
                    "سأرسل لك زر بدء الاختبار عندما يصبح بنك الأسئلة جاهزاً."
                )
        except Exception:
            logger.exception("queue basic question_bank failed; falling back to direct generation")

    await safe_edit(query, t(context.user_data, "generating_pool"))
    await start_quiz_from_bank(update, context, parts[1], "basic", count)

async def start_quiz_from_bank(update: Update, context: ContextTypes.DEFAULT_TYPE, setup_id: str, level: str, requested_count: int):
    """Start quiz directly without worker/Redis queue. Stable version."""
    query = update.callback_query
    pending = context.user_data.get("pending_file")
    if (not isinstance(pending, dict) or pending.get("setup_id") != setup_id) and update.effective_user:
        pending = load_pending_quiz_from_redis(update.effective_user.id, setup_id)
        if isinstance(pending, dict):
            context.user_data["pending_file"] = pending
    if not isinstance(pending, dict) or pending.get("setup_id") != setup_id:
        return await safe_edit(query, t(context.user_data, "stale"))

    try:
        await safe_edit(query, t(context.user_data, "generating_pool"))

        questions = await get_or_create_question_bank(pending["file_hash"], pending["text"], level)
        questions = normalize_questions(questions)

        if not questions:
            return await safe_edit(query, "⚠️ لم أستطع توليد أسئلة من هذا الملف. جرب ملفاً أوضح أو اختر ملخص أولاً ثم أعد المحاولة.")

        pool = [shuffle_options(q) for q in questions]
        pool = [q for q in pool if q]
        random.shuffle(pool)
        selected = pool[:min(requested_count, len(pool))]

        if len(selected) < requested_count:
            await context.bot.send_message(update.effective_chat.id, t(context.user_data, "not_enough", n=len(selected)))

        context.user_data["session"] = {
            "quiz_id": str(uuid.uuid4())[:8],
            "questions": selected,
            "current": 0,
            "score": 0,
            "wrong": [],
            "level": level,
        }
        context.user_data.pop("pending_file", None)

        await safe_edit(query, t(context.user_data, "quiz_ready", n=len(selected)))
        await send_current_question(update.effective_chat.id, context)
        if context.user_data.pop("_needs_ai_limit_for_level", "") == level:
            await increment_daily_usage(context.user_data, level, update.effective_user.id if update.effective_user else 0)

    except Exception as e:
        logger.exception("start_quiz_from_bank error")
        await safe_edit(query, t(context.user_data, "error", err=str(e)))

async def answer_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    parts = query.data.split(":")
    if len(parts) != 4:
        return
    quiz_id, chosen = parts[1], parts[3]
    try:
        idx = int(parts[2])
    except Exception:
        return

    session = context.user_data.get("session")
    if not isinstance(session, dict) or session.get("quiz_id") != quiz_id or int(session.get("current", -1)) != idx:
        try:
            return await query.edit_message_reply_markup(reply_markup=None)
        except Exception:
            return

    questions = normalize_questions(session.get("questions", []))
    session["questions"] = questions
    if idx < 0 or idx >= len(questions):
        return

    q = questions[idx]
    opts = _normalize_options(q.get("options", {}))
    correct_letter = q.get("correct", "")
    is_correct = chosen == correct_letter

    if is_correct:
        session["score"] = int(session.get("score", 0)) + 1
        res = t(context.user_data, "correct")
    else:
        res = t(context.user_data, "wrong", chosen=chosen)
        wrong = session.setdefault("wrong", [])
        if isinstance(wrong, list):
            wrong.append({"question": q, "chosen": chosen})

    level = session.get("level", "")
    if isinstance(level, str) and level.startswith("past:"):
        try:
            _, uni, subject, *_ = level.split(":")
            user_id = query.from_user.id if query.from_user else update.effective_chat.id
            await pp_record_answer(user_id, uni, subject, q, is_correct)
        except Exception:
            logger.exception("failed to record past paper progress")

    feedback = (
        f"{format_question(q, idx + 1, len(questions), context.user_data)}\n\n"
        f"{escape_md(res)}\n"
        f"{t(context.user_data, 'answer')} {correct_letter}. {escape_md(opts.get(correct_letter, ''))}\n\n"
        f"{t(context.user_data, 'explanation')} {escape_md(q.get('explanation', ''))}"
    )
    chunks = [feedback[i:i + 4000] for i in range(0, len(feedback), 4000)]
    try:
        await query.edit_message_text(chunks[0], parse_mode="Markdown")
    except Exception:
        try:
            await query.edit_message_reply_markup(reply_markup=None)
        except Exception:
            pass
        await context.bot.send_message(update.effective_chat.id, chunks[0].replace("`", ""))
    for c in chunks[1:]:
        await context.bot.send_message(update.effective_chat.id, c.replace("`", ""))

    session["current"] = int(session.get("current", 0)) + 1
    await asyncio.sleep(0.25)
    if session["current"] >= len(questions):
        await send_final_score(update.effective_chat.id, context)
    else:
        await send_current_question(update.effective_chat.id, context)

async def review_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query; await safe_answer_query(query); action = query.data.split(":", 1)[1]
    if action == "start":
        if not context.user_data.get("review_pool"): return await safe_edit(query, t(context.user_data, "review_none"))
        context.user_data["review"] = {"pool": context.user_data["review_pool"], "current": 0}
        try: await query.edit_message_reply_markup(reply_markup=None)
        except Exception: pass
        await send_review(update.effective_chat.id, context)
    elif action == "next" and context.user_data.get("review"):
        context.user_data["review"]["current"] += 1
        try: await query.edit_message_reply_markup(reply_markup=None)
        except Exception: pass
        await send_review(update.effective_chat.id, context)
    elif action == "exit":
        context.user_data.pop("review", None)
        try: await query.edit_message_reply_markup(reply_markup=None)
        except Exception: pass
        await context.bot.send_message(update.effective_chat.id, t(context.user_data, "review_done"))
    elif action == "stats":
        try: await query.edit_message_reply_markup(reply_markup=None)
        except Exception: pass
        stats = get_stats(context.user_data)
        await context.bot.send_message(update.effective_chat.id, t(context.user_data, "stats_text", quizzes=stats["quizzes"], correct=stats["correct"], total=stats["total"], accuracy=round(stats["correct"] / stats["total"] * 100) if stats["total"] else 0, best=stats["best"]))
    elif action == "new":
        try: await query.edit_message_reply_markup(reply_markup=None)
        except Exception: pass
        await context.bot.send_message(update.effective_chat.id, t(context.user_data, "choose_mode"), reply_markup=mode_keyboard(context.user_data))

async def send_review(chat_id: int, context: ContextTypes.DEFAULT_TYPE):
    review = context.user_data.get("review")
    if not review: return
    pool, current = review["pool"], review["current"]
    if current >= len(pool):
        context.user_data.pop("review", None); return await context.bot.send_message(chat_id, t(context.user_data, "review_done"))
    text = format_review(pool[current], current + 1, len(pool), context.user_data)
    chunks = [text[i:i + 4000] for i in range(0, len(text), 4000)]
    for i, c in enumerate(chunks):
        try:
            if i == len(chunks) - 1: await context.bot.send_message(chat_id, c, parse_mode="Markdown", reply_markup=review_keyboard(current, len(pool)))
            else: await context.bot.send_message(chat_id, c, parse_mode="Markdown")
        except Exception:
            if i == len(chunks) - 1: await context.bot.send_message(chat_id, c.replace("`", ""), reply_markup=review_keyboard(current, len(pool)))
            else: await context.bot.send_message(chat_id, c.replace("`", ""))

async def panel_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query; await safe_answer_query(query); await db_register_user(update.effective_user); action = query.data.split(":", 1)[1]
    if action == "support":
        return await safe_edit(query, "📞 مركز الدعم\n\n📱 رقم المدار: 0918874659\n👤 التواصل: عبر رسالة الدعم داخل البوت\n\nاختر:", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("💳 دعم عبر رصيد المدار", callback_data="support:balance")], [InlineKeyboardButton("💬 تواصل مع الدعم", callback_data="support:contact")], [InlineKeyboardButton("📢 أبلغ عن مشكلة", callback_data="support:problem")], [InlineKeyboardButton("⬅️ رجوع", callback_data="main:home")]]))
    if action == "progress":
        p = await db_get_progress(update.effective_user.id if update.effective_user else update.effective_chat.id)
        msg = f"📊 تقدمي ونقاطي\n\n⭐ XP: {p.get('xp', 0)}\n🎖 Level: {p.get('level', 1)}\n🔥 Streak: {p.get('streak', 0)} يوم\n📝 Quizzes: {p.get('quizzes', 0)}\n🎯 Accuracy: {round(int(p.get('correct') or 0) / int(p.get('total') or 1) * 100) if p.get('total') else 0}%\n🎭 OSCE completed: {p.get('osce_done', 0)}\n📚 Summaries: {p.get('summaries_done', 0)}\n\nكل اختبار/ملخص/OSCE يعطيك XP ويرفع مستواك."
        await safe_edit(query, msg, reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("⬅️ رجوع", callback_data="main:home")]]))
    elif action == "leaderboard":
        rows = await db_get_leaderboard(10)
        lines = [f"{['🥇','🥈','🥉'][i] if i < 3 else f'{i+1}.'} {r.get('full_name') or r.get('username') or 'Student'} — ⭐ {r.get('xp', 0)} XP | 🎖 L{r.get('level', 1)} | 🔥 {r.get('streak', 0)}" for i, r in enumerate(rows)]
        await safe_edit(query, "🏆 لوحة الصدارة\n\n" + ("\n".join(lines) if lines else "لا توجد نتائج بعد."), reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("⬅️ رجوع", callback_data="main:home")]]))
    elif action == "subscription":
        uid = update.effective_user.id if update.effective_user else 0
        info = await db_get_user_plan_info(uid)
        plan = info.get("plan", "free")
        exp = float(info.get("plan_expires_at") or 0)
        exp_txt = "لا يوجد" if not exp else datetime.fromtimestamp(exp).strftime("%Y-%m-%d")
        msg = (
            "💎 الاشتراك Premium\n\n"
            f"خطة حسابك الحالية: {PLAN_LABELS.get(plan, plan)}\n"
            f"تنتهي في: {exp_txt}\n"
            f"User ID الخاص بك: `{uid}`\n\n"
            "Free:\n"
            "• ملخص واحد يومياً\n"
            "• Basic مرة يومياً\n"
            "• Cases و Challenge كل يومين\n"
            "• OSCE مرة أسبوعياً\n\n"
            "💎 Premium شهري — 15 د.ل\n"
            "• 3 ملخصات يومياً\n"
            "• 3 Basic يومياً\n"
            "• 2 Cases + 2 Challenge يومياً\n"
            "• OSCE مرتين بالأسبوع\n\n"
            "🔥 Premium فصل كامل — 40 د.ل / 4 شهور\n"
            "نفس مزايا Premium بسعر أوفر.\n\n"
            "🚀 Premium Plus فصل كامل — 60 د.ل / 4 شهور\n"
            "• 5 ملخصات يومياً\n"
            "• 5 Basic يومياً\n"
            "• 3 Cases + 3 Challenge يومياً\n"
            "• OSCE يومياً\n\n"
            "للتفعيل: أرسل إثبات الدفع مع User ID من زر الدعم."
        )
        await safe_edit(query, msg, parse_mode="Markdown", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🧾 طلب تفعيل / إرسال إثبات", callback_data="payment:request")], [InlineKeyboardButton("📞 الدعم", callback_data="panel:support")], [InlineKeyboardButton("⬅️ رجوع", callback_data="main:home")]]))

async def payment_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query; await safe_answer_query(query)
    if update.effective_user:
        await db_register_user(update.effective_user)
        await db_create_payment_request(update.effective_user.id, "User opened payment request")
        await send_admin_log(context, f"💎 طلب اشتراك\nUser ID: {update.effective_user.id}\nName: {update.effective_user.full_name}")
    await safe_edit(query, "🧾 تم فتح طلب اشتراك مبدئي.\nللتفعيل اليدوي: اضغط الدعم وأرسل إثبات الدفع مع User ID الخاص بك.", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("📞 الدعم", callback_data="panel:support")], [InlineKeyboardButton("⬅️ رجوع", callback_data="main:home")]]))

async def my_plan_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id if update.effective_user else 0
    info = await db_get_user_plan_info(user_id)
    plan = info.get("plan", "free")
    exp = float(info.get("plan_expires_at") or 0)
    exp_txt = "لا يوجد" if not exp else datetime.fromtimestamp(exp).strftime("%Y-%m-%d")
    await update.message.reply_text(
        f"💎 خطة حسابك\n\n"
        f"User ID: `{user_id}`\n"
        f"الخطة: {PLAN_LABELS.get(plan, plan)}\n"
        f"تنتهي في: {exp_txt}",
        parse_mode="Markdown"
    )

async def set_plan_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id if update.effective_user else 0
    if not is_admin_user(user_id):
        return await update.message.reply_text("⚠️ هذا الأمر للأدمن فقط.")
    if len(context.args) < 2:
        return await update.message.reply_text(
            "الاستخدام:\n"
            "/set_plan USER_ID premium 30\n"
            "/set_plan USER_ID premium 120\n"
            "/set_plan USER_ID premium_plus 120\n"
            "/set_plan USER_ID free\n\n"
            "الأسعار المعتمدة:\n"
            "Premium شهري: 15 د.ل\n"
            "Premium فصل كامل: 40 د.ل\n"
            "Premium Plus فصل كامل: 60 د.ل"
        )
    try:
        target_id = int(context.args[0])
    except Exception:
        return await update.message.reply_text("⚠️ USER_ID غير صالح.")
    plan = context.args[1].strip().lower()
    if plan not in {"free", "premium", "premium_plus", "admin"}:
        return await update.message.reply_text("⚠️ الخطة لازم تكون: free أو premium أو premium_plus أو admin")
    days = None
    if len(context.args) >= 3:
        try:
            days = int(context.args[2])
        except Exception:
            days = None
    if plan == "premium" and days is None:
        days = 30
    if plan == "premium_plus" and days is None:
        days = 120
    await db_set_user_plan(target_id, plan, days)
    info = await db_get_user_plan_info(target_id)
    exp = float(info.get("plan_expires_at") or 0)
    exp_txt = "لا يوجد" if not exp else datetime.fromtimestamp(exp).strftime("%Y-%m-%d")
    await update.message.reply_text(f"✅ تم ضبط خطة المستخدم {target_id} إلى: {PLAN_LABELS.get(plan, plan)}\nتنتهي في: {exp_txt}")

async def remove_premium_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id if update.effective_user else 0
    if not is_admin_user(user_id):
        return await update.message.reply_text("⚠️ هذا الأمر للأدمن فقط.")
    if not context.args:
        return await update.message.reply_text("الاستخدام: /remove_premium USER_ID")
    try:
        target_id = int(context.args[0])
    except Exception:
        return await update.message.reply_text("⚠️ USER_ID غير صالح.")
    await db_set_user_plan(target_id, "free")
    await update.message.reply_text(f"✅ تم إلغاء Premium للمستخدم {target_id}")

async def ai_policy_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id if update.effective_user else 0
    if not is_admin_user(user_id):
        return await update.message.reply_text("⚠️ هذا الأمر للأدمن فقط.")
    msg = (
        "🛡️ سياسة مزودي الذكاء الحالية\n\n"
        f"Google keys: {len(GOOGLE_API_KEYS)}\n"
        f"OpenRouter free keys: {len(OPENROUTER_FREE_KEYS)} | model: {OPENROUTER_FREE_MODEL}\n"
        f"Groq keys: {len(GROQ_API_KEYS)} | model: {GROQ_MODEL}\n"
        f"Cloudflare keys: {len(CLOUDFLARE_API_TOKENS)} | model: {CLOUDFLARE_MODEL}\n"
        f"Paid OpenRouter keys: {len(OPENROUTER_PAID_KEYS)} | model: {GEMINI_MODEL_NAME}\n\n"
        f"FREE_PROVIDER_MODE: {FREE_PROVIDER_MODE}\n"
        f"PREMIUM_ENABLED: {PREMIUM_ENABLED}\n"
        f"Premium prices: 15 / 40 / 60 LYD\n"
        f"ALLOW_PAID_FALLBACK_FOR_USERS: {ALLOW_PAID_FALLBACK_FOR_USERS}\n"
        f"ALLOW_PAID_FALLBACK_FOR_PREMIUM: {ALLOW_PAID_FALLBACK_FOR_PREMIUM}\n"
        f"ALLOW_PAID_FALLBACK_FOR_ADMIN: {ALLOW_PAID_FALLBACK_FOR_ADMIN}\n\n"
        "الترتيب: كاش → Google → OpenRouter Free → Groq Free → Cloudflare Free → Paid فقط للأدمن/بريميوم حسب الإعداد."
    )
    await update.message.reply_text(msg)

async def back_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    data = _safe_callback_data(query.data or "")
    if data == "main:home":
        return await show_main_menu(query, context)
    parts = query.data.split(":")
    target = parts[1] if len(parts) > 1 else "mode"
    if target in ("mode", "main"):
        await show_main_menu(query, context)
    elif target == "support":
        await safe_edit(
            query,
            "📞 الدعم\n\n💳 رقم المدار: 0918874659\n👤 التواصل: عبر رسالة الدعم داخل البوت",
            reply_markup=InlineKeyboardMarkup([
                [InlineKeyboardButton("💳 دعم عبر رصيد المدار", callback_data="support:balance")],
                [InlineKeyboardButton("💬 تواصل مع الدعم", callback_data="support:contact")],
                [InlineKeyboardButton("⬅️ القائمة الرئيسية", callback_data="main:home")]
            ])
        )
    elif target == "levels" and len(parts) == 3:
        pending = context.user_data.get("pending_file")
        if not pending or pending.get("setup_id") != parts[2]: return await safe_edit(query, t(context.user_data, "stale"))
        await safe_edit(query, t(context.user_data, "file_ready_quiz"), reply_markup=level_keyboard(parts[2], context.user_data))
    else:
        await show_main_menu(query, context)

async def support_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query; await safe_answer_query(query)
    if query.data == "support:balance": await safe_edit(query, "💳 دعم عبر رصيد المدار\n\nرقم المدار:\n0918874659\n\nبعد إرسال الرصيد، اضغط تواصل مع الدعم.", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("💬 تواصل مع الدعم", callback_data="support:contact")], [InlineKeyboardButton("⬅️ رجوع", callback_data="panel:support")]]))
    elif query.data == "support:contact": context.user_data["awaiting_support_message"] = "contact"; await safe_edit(query, "💬 اكتب رسالتك الآن.\n\nمثال:\nاسمي أحمد، أرسلت رصيد، وأريد تفعيل الخدمة.")
    elif query.data == "support:problem": context.user_data["awaiting_support_message"] = "problem"; await safe_edit(query, "📢 اكتب المشكلة التي تواجهك الآن.")

async def admin_panel_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    if not is_admin_user(query.from_user.id if query.from_user else None):
        return await safe_edit(query, "⛔ هذه اللوحة للإدارة فقط.", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("⬅️ رجوع", callback_data="main:home")]]))
    data = query.data or "admin:home"
    parts = data.split(":")
    action = parts[1] if len(parts) > 1 else "home"
    if action == "home":
        kb = InlineKeyboardMarkup([
            [InlineKeyboardButton("➕ إضافة أسئلة للمخزون", callback_data="admin:add")],
            [InlineKeyboardButton("📚 تخزين محاضرات الفصل", callback_data="admin:preload")],
            [InlineKeyboardButton("⚡ جمع سنوات سابقة سريع", callback_data="admin:quickpast")],
            [InlineKeyboardButton("📦 أعطيني المخزون ZIP", callback_data="admin:dumpzip")],
            [InlineKeyboardButton("🧹 تفريغ المخزون المؤقت", callback_data="admin:clear")],
            [InlineKeyboardButton("❌ إلغاء التجميع الحالي", callback_data="admin:cancel")],
            [InlineKeyboardButton("⏹️ إيقاف تخزين محاضرات الفصل", callback_data="admin:preloadstop")],
            [InlineKeyboardButton("⬅️ رجوع للرئيسية", callback_data="main:home")],
        ])
        st = context.user_data.get("active_staging")
        active = f"\n\n📥 التجميع الحالي: {st['uni']} / {st['subject']}" if st else "\n\nلا يوجد تجميع نشط الآن."
        return await safe_edit(query, "🛠️ لوحة الأدمن\n" + active, reply_markup=kb)
    if action == "preload":
        context.user_data.pop("active_staging", None)
        context.user_data["admin_preload"] = True
        text = (
            "📚 وضع تخزين محاضرات الفصل مفعّل.\n\n"
            "أرسل الآن ملفات المحاضرات واحداً واحداً، أو ارفع ZIP واحد يحتوي كل المحاضرات.\n"
            "الحد الأقصى لكل ملف داخل ZIP: 30MB.\n"
            "للأدمن فقط: لو الملف PPT قديم أو PDF مصور أو صورة، سأحاول قراءته بالدعم الذكي Vision ثم أخزنه.\n"
            "سأخزن النص وأجهز الملخصات وبنوك الأسئلة في الخلفية حتى يستفيد الطلاب لاحقاً بدون API.\n\n"
            "للإيقاف اضغط: ⏹️ إيقاف تخزين محاضرات الفصل"
        )
        return await safe_edit(
            query,
            text,
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("⏹️ إيقاف تخزين محاضرات الفصل", callback_data="admin:preloadstop")], [InlineKeyboardButton("🛠️ رجوع للوحة الأدمن", callback_data="admin:home")]])
        )
    if action == "preloadstop":
        context.user_data.pop("admin_preload", None)
        return await safe_edit(query, "✅ تم إيقاف وضع تخزين محاضرات الفصل.", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🛠️ رجوع للوحة الأدمن", callback_data="admin:home")]]))
    if action == "quickpast":
        kb = InlineKeyboardMarkup([[InlineKeyboardButton("🏛️ أجدابيا", callback_data="admin:quickuni:ajdabiya")], [InlineKeyboardButton("🏛️ بنغازي", callback_data="admin:quickuni:benghazi")], [InlineKeyboardButton("⬅️ رجوع", callback_data="admin:home")]])
        return await safe_edit(query, "⚡ جمع سنوات سابقة سريع\n\nاختر الجامعة مرة واحدة، بعدها أرسل أو Forward أي صور/ملفات/رسائل.\nعند الانتهاء اكتب /finish_past أو اضغط أعطيني ZIP.", reply_markup=kb)
    if action == "quickuni" and len(parts) >= 3:
        uni = parts[2]
        kb = InlineKeyboardMarkup([[InlineKeyboardButton(label, callback_data=f"admin:quicksubject:{uni}:{sub}")] for sub, label in PAST_SUBJECTS] + [[InlineKeyboardButton("⬅️ رجوع", callback_data="admin:quickpast")]])
        return await safe_edit(query, f"⚡ جمع سنوات سابقة سريع\n🏛️ {UNI_LABELS.get(uni, uni)}\n\nاختر المادة:", reply_markup=kb)
    if action == "quicksubject" and len(parts) >= 4:
        uni, subject = parts[2], parts[3]
        context.user_data.pop("admin_preload", None)
        context.user_data["active_staging"] = {"uni": uni, "subject": subject, "quick": True}
        return await safe_edit(query, f"✅ فتح صندوق الجمع السريع\n\n🏛️ {UNI_LABELS.get(uni, uni)}\n📚 {SUBJECT_LABELS.get(subject, subject)}\n\nالآن فقط أرسل أو Forward أي شيء عندك: صور، PDF، Word، PPT، نصوص أو سكرينات.\nسأحفظ الأصل والنص للمراجعة.\n\nعند الانتهاء اكتب: /finish_past", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("📦 أعطيني ZIP الآن", callback_data="admin:dumpzip")], [InlineKeyboardButton("🛠️ رجوع للوحة الأدمن", callback_data="admin:home")]]))
    if action == "add":
        kb = InlineKeyboardMarkup([[InlineKeyboardButton("🏛️ أجدابيا", callback_data="admin:adduni:ajdabiya")], [InlineKeyboardButton("🏛️ بنغازي", callback_data="admin:adduni:benghazi")], [InlineKeyboardButton("⬅️ رجوع", callback_data="admin:home")]])
        return await safe_edit(query, "➕ إضافة أسئلة\n\nاختر الجامعة:", reply_markup=kb)
    if action == "adduni" and len(parts) >= 3:
        uni = parts[2]
        kb = InlineKeyboardMarkup([[InlineKeyboardButton(label, callback_data=f"admin:addsubject:{uni}:{sub}")] for sub, label in PAST_SUBJECTS] + [[InlineKeyboardButton("⬅️ رجوع", callback_data="admin:add")]])
        return await safe_edit(query, f"➕ إضافة أسئلة\n🏛️ {UNI_LABELS.get(uni, uni)}\n\nاختر المادة:", reply_markup=kb)
    if action == "addsubject" and len(parts) >= 4:
        uni, subject = parts[2], parts[3]
        context.user_data.pop("admin_preload", None)
        context.user_data["active_staging"] = {"uni": uni, "subject": subject}
        return await safe_edit(query, f"✅ بدأ التجميع المؤقت\n\n🏛️ {UNI_LABELS.get(uni, uni)}\n📚 {SUBJECT_LABELS.get(subject, subject)}\n\nأرسل الآن أي سؤال كنص، صورة، PDF، Word أو PPT.\nكل شيء ترسله سيتخزن في المخزون المؤقت.\n\nبعد الانتهاء اضغط من لوحة الأدمن: 📦 أعطيني المخزون ZIP", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🛠️ رجوع للوحة الأدمن", callback_data="admin:home")]]))
    if action == "dumpzip":
        await admin_send_staging_zip(update, context)
        return
    if action == "clear":
        user_id = query.from_user.id
        async with _db_lock:
            def work():
                conn = _db_connect(); conn.execute("DELETE FROM admin_staging WHERE user_id=?", (user_id,)); conn.commit(); conn.close()
            await asyncio.to_thread(work)
        context.user_data.pop("active_staging", None)
        return await safe_edit(query, "🧹 تم تفريغ المخزون المؤقت.", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🛠️ رجوع للوحة الأدمن", callback_data="admin:home")]]))
    if action == "cancel":
        context.user_data.pop("active_staging", None)
        return await safe_edit(query, "❌ تم إلغاء وضع التجميع الحالي فقط.\nالمخزون المحفوظ لم يتم حذفه.", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🛠️ رجوع للوحة الأدمن", callback_data="admin:home")]]))

async def admin_send_staging_zip(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    user_id = query.from_user.id
    async with _db_lock:
        def work():
            conn = _db_connect()
            rows = conn.execute("SELECT uni, subject, raw_content, created_at FROM admin_staging WHERE user_id=? ORDER BY uni, subject, created_at", (user_id,)).fetchall()
            conn.close()
            return [dict(r) for r in rows]
        rows = await asyncio.to_thread(work)
    if not rows:
        return await safe_edit(query, "📦 المخزون المؤقت فارغ.", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🛠️ رجوع للوحة الأدمن", callback_data="admin:home")]]))
    grouped: Dict[Tuple[str, str], List[dict]] = {}
    for r in rows:
        grouped.setdefault((r["uni"], r["subject"]), []).append(r)
    mem = io.BytesIO()
    with zipfile.ZipFile(mem, "w", zipfile.ZIP_DEFLATED) as z:
        readme = "هذا ZIP صادر من بوت MedMCQ.\nنظف الأسئلة وحولها إلى JSON array بالشكل: question/options/correct/explanation.\nثم ارفع الناتج باستخدام /upload_final uni subject\n"
        z.writestr("README.txt", readme)
        for (uni, subject), items in grouped.items():
            text = []
            for i, item in enumerate(items, 1):
                text.append(f"===== ITEM {i} | {item.get('created_at','')} =====\n{item.get('raw_content','')}")
            z.writestr(f"{uni}/{subject}/raw_questions.txt", "\n\n".join(text))
    mem.seek(0)
    mem.name = f"admin_staging_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.zip"
    await context.bot.send_document(chat_id=update.effective_chat.id, document=mem, caption="📦 هذا مخزون الأسئلة بصيغة ZIP. أرسله لـ ChatGPT لتنظيفه وتحويله إلى JSON.")
    await safe_edit(query, "✅ تم إرسال ZIP.\nالمخزون لم يُحذف. إذا تريد حذفه اضغط: تفريغ المخزون المؤقت.", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("🧹 تفريغ المخزون المؤقت", callback_data="admin:clear")], [InlineKeyboardButton("🛠️ رجوع للوحة الأدمن", callback_data="admin:home")]]))


# =============================================================================
# OVERRIDES: ADMIN STAGING ZIP WITH ORIGINAL FILES
# =============================================================================

async def stage_dump_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id if update.effective_user else 0
    if user_id not in ADMIN_USER_IDS:
        return
    mem = await build_admin_staging_zip(user_id)
    if not mem:
        return await update.message.reply_text("المستودع فارغ. أرسل أسئلة أو ملفات أولاً.")
    await context.bot.send_document(
        chat_id=update.effective_chat.id,
        document=mem,
        caption="📦 هذا ZIP يحتوي النص المستخرج + الملفات الأصلية. أرسله لـ ChatGPT للتنظيف والتحويل إلى JSON."
    )
    await update.message.reply_text("✅ تم تصدير ZIP. لم يتم حذف المخزون المؤقت.")

async def admin_send_staging_zip(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await safe_answer_query(query)
    user_id = query.from_user.id if query and query.from_user else 0
    if user_id not in ADMIN_USER_IDS:
        return await safe_edit(query, "⛔ هذه الميزة للإدارة فقط.")
    mem = await build_admin_staging_zip(user_id)
    if not mem:
        return await safe_edit(query, "المخزون فارغ. أرسل أسئلة أو ملفات أولاً.")
    await context.bot.send_document(
        chat_id=update.effective_chat.id,
        document=mem,
        caption="📦 هذا ZIP يحتوي النص المستخرج + الملفات الأصلية. أرسله لـ ChatGPT لتنظيفه وتحويله إلى JSON."
    )
    await safe_edit(
        query,
        "✅ تم إرسال ZIP.\nالمخزون لم يُحذف. إذا تريد حذفه اضغط: تفريغ المخزون المؤقت.",
        reply_markup=InlineKeyboardMarkup([
            [InlineKeyboardButton("🧹 تفريغ المخزون المؤقت", callback_data="admin:clear")],
            [InlineKeyboardButton("🛠️ رجوع للوحة الأدمن", callback_data="admin:home")],
        ])
    )

async def universal_callback_router(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    raw_data = query.data or ""
    data = _safe_callback_data(raw_data)
    # removed invalid CallbackQuery.data assignment
    try:
        if data.startswith("lang:"):
            await lang_callback(update, context)
        elif data == "main:home" or data.startswith("back:"):
            await back_callback(update, context)
        elif data.startswith("admin:"):
            await admin_panel_callback(update, context)
        elif data == "pp:home":
            await pp_home_callback(update, context)
        elif data.startswith("pp:uni:"):
            await pp_uni_callback(update, context)
        elif data.startswith("pp:subject:"):
            await pp_sub_callback(update, context)
        elif data.startswith("pp:mode:"):
            await pp_mode_callback(update, context)
        elif data.startswith("pp:quiz:"):
            await pp_quiz_start_callback(update, context)
        elif data.startswith("pp:wrong:"):
            await pp_wrong_callback(update, context)
        elif data.startswith("pp:progress:"):
            await pp_progress_callback(update, context)
        elif data.startswith("mode:"):
            await mode_callback(update, context)
        elif data.startswith("panel:"):
            await panel_callback(update, context)
        elif data.startswith("style:"):
            await summary_style_callback(update, context)
        elif data.startswith("level:"):
            await level_callback(update, context)
        elif data.startswith("count:"):
            await count_callback(update, context)
        elif data.startswith("ans:"):
            await answer_callback(update, context)
        elif data.startswith("review:"):
            await review_callback(update, context)
        elif data.startswith("payment:"):
            await payment_callback(update, context)
        elif data.startswith("support:"):
            await support_callback(update, context)
        elif data.startswith("osce4:"):
            await osce_v4_callback(update, context)
        else:
            await safe_answer_query(query)
            await safe_edit(query, f"⚠️ زر غير معروف أو قديم.\n\nالكود: {raw_data}\n\nافتح /start من جديد.", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("⬅️ رجوع للرئيسية", callback_data="main:home")]]))
    except ApplicationHandlerStop:
        raise
    except Exception as e:
        logger.exception("callback router error for data=%s", raw_data)
        await safe_answer_query(query)
        await safe_edit(query, f"⚠️ حدث خطأ تلقائي وتم احتواؤه.\n\nالزر: {raw_data}\nالخطأ: {str(e)[:300]}\n\nاضغط رجوع للرئيسية وجرب مرة ثانية.", reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("⬅️ رجوع للرئيسية", callback_data="main:home")]]))
    finally:
        raise ApplicationHandlerStop

async def send_admin_log(context: ContextTypes.DEFAULT_TYPE, text: str):
    try: await context.bot.send_message(chat_id=ADMIN_LOG_CHAT_ID, text=text)
    except Exception: pass

# =============================================================================
# FILE & PHOTO HANDLING
# =============================================================================

async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    set_ai_user_id(update.effective_user.id if update.effective_user else 0)
    user_id = update.effective_user.id
    if user_id not in ADMIN_USER_IDS:
        await update.message.reply_text("⚠️ عذراً، ميزة قراءة الصور وتحليلها مخصصة للإدارة فقط لتوفير استهلاك الـ API.\nيرجى إرسال ملفات نصية (PDF, Word).")
        return
    if context.user_data.get("admin_preload"):
        status = await update.message.reply_text("🩻 [الأدمن]: صورة محاضرة مكتشفة. جاري قراءة الصورة وتخزينها...")
        photo_file = await update.message.photo[-1].get_file()
        raw_bytes = bytes(await photo_file.download_as_bytearray())
        extracted_text = await call_gemini_vision(raw_bytes)
        if not extracted_text or len(extracted_text.strip()) < 50:
            return await status.edit_text("⚠️ لم أستطع استخراج نص واضح من الصورة.")
        fake_bytes = extracted_text.encode("utf-8")
        await handle_admin_preload_file(update, context, fake_bytes, f"photo_lecture_{int(time.time())}.txt", "text/plain", status)
        return
    if context.user_data.get("active_staging"):
        status = await update.message.reply_text("🩻 [الأدمن]: صورة ورقة امتحان مكتشفة! جاري تشغيل الـ Vision مجاناً...")
        photo_file = await update.message.photo[-1].get_file()
        raw_bytes = bytes(await photo_file.download_as_bytearray())
        st = context.user_data.get("active_staging")
        original_path = await admin_save_staging_original(user_id, st, f"photo_{int(time.time())}.jpg", raw_bytes)
        extracted_text = await call_gemini_vision(raw_bytes)
        await status.delete()
        await handle_admin_acc(update, context, f"[صورة امتحان]\n[original_file: {original_path}]\n\n{extracted_text}")

async def handle_file(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    user_id = user.id if user else 0
    set_ai_user_id(user_id)
    doc = update.message.document
    if not doc: return await update.message.reply_text("⚠️ أرسل ملفاً (PDF, DOCX, PPTX, TXT).")

    # 1. مسار حقن الـ JSON للأدمن
    pf = context.user_data.get("pending_final")
    if user_id in ADMIN_USER_IDS and pf:
        tg_file = await doc.get_file()
        raw_bytes = bytes(await tg_file.download_as_bytearray())
        try:
            json_text = raw_bytes.decode('utf-8', errors='ignore').strip()
            json_text = json_text.replace("```json", "").replace("```", "").strip()
            parsed_data = json.loads(json_text)
            if isinstance(parsed_data, dict):
                for key in ("questions", "items", "data"):
                    if isinstance(parsed_data.get(key), list):
                        parsed_data = parsed_data[key]
                        break

            if isinstance(parsed_data, list):
                new_questions, duplicate_count, invalid_count = await pp_filter_new_questions(pf["uni"], pf["subject"], parsed_data)

                if new_questions:
                    async with _db_lock:
                        def w():
                            conn = _db_connect()
                            conn.execute(
                                "INSERT INTO past_papers (uni, subject, questions_json, created_at) VALUES (?, ?, ?, ?)",
                                (pf["uni"], pf["subject"], json.dumps(new_questions, ensure_ascii=False), datetime.utcnow().isoformat())
                            )
                            conn.commit(); conn.close()
                        await asyncio.to_thread(w)

                total = len(parsed_data)
                saved = len(new_questions)
                context.user_data.pop("pending_final", None)

                await update.message.reply_text(
                    "🚀 **تمت معالجة أسئلة السنوات السابقة**\n\n"
                    f"🏛️ الجامعة: `{pf['uni']}`\n"
                    f"📚 المادة: `{pf['subject']}`\n\n"
                    f"✅ الجديد المحفوظ: `{saved}`\n"
                    f"♻️ المكرر المتجاهل: `{duplicate_count}`\n"
                    f"⚠️ غير صالح/غير مقروء: `{invalid_count}`\n"
                    f"📦 إجمالي الملف: `{total}`\n\n"
                    "إذا كان السؤال موجوداً سابقاً فلن يتم حفظه مرة ثانية.",
                    parse_mode="Markdown"
                )
                return
        except Exception as e:
            return await update.message.reply_text(f"❌ هيكلية الـ JSON غير صالحة. التفاصيل: {e}")

    # 2. مسار التجميع للأدمن
    st = context.user_data.get("active_staging")
    if user_id in ADMIN_USER_IDS and st:
        status = await update.message.reply_text("🔄 جاري حفظ الملف الأصلي واستخراج النص وإيداعه في المستودع...")
        tg_file = await doc.get_file()
        raw_bytes = bytes(await tg_file.download_as_bytearray())
        original_path = await admin_save_staging_original(user_id, st, doc.file_name or "uploaded_file", raw_bytes)
        file_type = classify_document(doc.mime_type or "", doc.file_name or "")
        if file_type != "unsupported":
            extracted_text = extract_text(file_type, raw_bytes)
        else:
            extracted_text = "[نوع الملف غير مدعوم للاستخراج النصي، لكن تم حفظ الملف الأصلي داخل original_files.]"
        await status.delete()
        await handle_admin_acc(update, context, f"[محتوى ملف مستند مجمع: {doc.file_name}]\n[original_file: {original_path}]\n\n{extracted_text}")
        return

    # 2.5 مسار تخزين محاضرات الفصل للأدمن
    if user_id in ADMIN_USER_IDS and context.user_data.get("admin_preload"):
        status = await update.message.reply_text("📚 جاري استلام وتخزين محاضرة الفصل...")
        tg_file = await doc.get_file()
        raw_bytes = bytes(await tg_file.download_as_bytearray())
        await handle_admin_preload_file(update, context, raw_bytes, doc.file_name or "lecture", doc.mime_type or "", status)
        return

    # 3. مسار الطلبة العادي
    if user and await cooldown_guard(update, context.user_data, user.id): return
    if context.user_data.get("busy"): return await update.message.reply_text(t(context.user_data, "busy"))

    context.user_data["busy"] = True
    try:
        file_type = classify_document(doc.mime_type or "", doc.file_name or "")
        if file_type == "unsupported": return await update.message.reply_text(t(context.user_data, "file_unsupported"))
        if doc.file_size and doc.file_size > MAX_FILE_SIZE_BYTES: return await update.message.reply_text(t(context.user_data, "file_too_large"))

        status = await update.message.reply_text(t(context.user_data, "file_received"))
        tg_file = await doc.get_file()
        raw_bytes = bytes(await tg_file.download_as_bytearray())
        file_hash = sha256_bytes(raw_bytes)

        cached_file = await db_get_file(file_hash)
        if cached_file:
            text = cached_file["extracted_text"]
            # If this exact file was preloaded, switch to the canonical text hash where AI outputs are cached.
            content_hash = canonical_text_hash(text)
            if await db_get_file(content_hash):
                file_hash = content_hash
            await status.edit_text(t(context.user_data, "cached_file")); await asyncio.sleep(0.4)
        else:
            text = extract_text(file_type, raw_bytes)
            if not text or len(text.strip()) < 100: return await status.edit_text(t(context.user_data, "file_empty"))
            content_hash = canonical_text_hash(text)
            content_cached = await db_get_file(content_hash)
            if content_cached:
                file_hash = content_hash
                text = content_cached["extracted_text"]
                await status.edit_text(t(context.user_data, "cached_file")); await asyncio.sleep(0.4)
            else:
                file_hash = content_hash
                await db_save_file(file_hash, doc.file_name or "document", len(raw_bytes), doc.mime_type or "", text)

        mode = context.user_data.get("pending_mode", "quiz")

        # --- OSCE V4 integration ---
        if mode == "osce":
            if not await check_daily_limit(update, context, "osce"):
                return
            await status.edit_text(t(context.user_data, "osce_generating"))
            case = await generate_osce_v4_case(text)
            session = create_osce_v4_session(case, file_hash)
            context.user_data["osce_v4_session"] = session
            context.user_data.pop("pending_mode", None)
            await increment_daily_usage(context.user_data, "osce", user.id if user else 0)
            if user:
                await db_add_xp(user.id, 10, "OSCE V4 case started")
            vignette_text = _osce_v4_vignette_text(session)
            await status.edit_text(vignette_text)
            return

        if mode == "summary":
            style = context.user_data.get("pending_summary_style", "disease_v2")
            cached_summary = await db_get_summary(file_hash, style)
            if cached_summary:
                await status.delete()
                full_text = t(context.user_data, "summary_header") + cached_summary
                for c in [full_text[i:i + 4000] for i in range(0, len(full_text), 4000)]:
                    await update.message.reply_text(c)
                return

            if not await check_daily_limit(update, context, "summary"):
                return

            if USE_QUEUE:
                try:
                    from queue_jobs import create_job, queue_size
                    job_id = create_job(
                        "summary",
                        update.effective_chat.id,
                        update.effective_user.id if update.effective_user else 0,
                        {
                            "file_hash": file_hash,
                            "text": text,
                            "style": style,
                            "lang": get_lang(context.user_data),
                        },
                    )
                    await increment_daily_usage(context.user_data, "summary", user.id if user else 0)
                    await status.edit_text(
                        "✅ تم استلام طلب الملخص.\n"
                        "📥 دخل الطلب في الطابور.\n"
                        f"📌 رقم الطلب: {job_id}\n"
                        f"📊 عدد الطلبات في الطابور: {queue_size()}\n\n"
                        "سأرسل الملخص تلقائياً عند الانتهاء."
                    )
                    return
                except Exception:
                    logger.exception("queue summary failed; falling back to direct summary")

            await status.edit_text(t(context.user_data, "generating_summary"))
            summary = await get_or_create_summary(file_hash, text, style)
            await status.delete()
            if summary:
                full_text = t(context.user_data, "summary_header") + summary
                chunks = [full_text[i:i + 4000] for i in range(0, len(full_text), 4000)]
                for c in chunks:
                    await update.message.reply_text(c)
                await increment_daily_usage(context.user_data, "summary", user.id if user else 0)
                if user:
                    p = await db_add_xp(user.id, 15, "Summary generated", summary_done=True)
                    await update.message.reply_text(f"⭐ +15 XP | 🎖 Level {p.get('level', 1)} | 🔥 Streak {p.get('streak', 0)}")
            else:
                await update.message.reply_text(t(context.user_data, "summary_failed"))
            return

        setup_id = str(uuid.uuid4())[:8]
        context.user_data["pending_mode"] = "quiz"
        pending_file = {"setup_id": setup_id, "file_hash": file_hash, "file_name": doc.file_name or "document", "text": text}
        context.user_data["pending_file"] = pending_file
        if USE_QUEUE and user:
            try:
                from queue_jobs import save_pending_quiz
                save_pending_quiz(user.id, setup_id, pending_file)
            except Exception:
                logger.exception("save pending quiz failed")
        await status.edit_text(t(context.user_data, "file_ready_quiz"), reply_markup=level_keyboard(setup_id, context.user_data))

    except Exception:
        logger.exception("handle_file error"); await update.message.reply_text("❌ حدث خطأ أثناء معالجة الملف.")
    finally: context.user_data["busy"] = False

async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    text = update.message.text

    # OSCE V4 text handler replaces old handler
    if await handle_osce_v4_text(update, context): return

    if user.id in ADMIN_USER_IDS and context.user_data.get("active_staging") and not text.startswith('/'):
        return await handle_admin_acc(update, context, text)

    support_type = context.user_data.get("awaiting_support_message")
    if support_type:
        context.user_data.pop("awaiting_support_message", None)
        kind = "💬 تواصل مع الدعم" if support_type == "contact" else "📢 بلاغ مشكلة"
        await send_admin_log(context, f"{kind}\n\nUser ID: {user.id if user else 'unknown'}\nName: {user.full_name if user else 'unknown'}\n\nMessage:\n{text}")
        return await update.message.reply_text("✅ تم إرسال رسالتك للدعم.")

    if context.user_data.get("session"): return await update.message.reply_text(t(context.user_data, "quiz_in_progress"))
    await update.message.reply_text(t(context.user_data, "send_file_hint"))

async def error_handler(update: object, context: ContextTypes.DEFAULT_TYPE):
    logger.error("Unhandled exception", exc_info=context.error)
    if isinstance(update, Update) and update.effective_message:
        try: await update.effective_message.reply_text(t(context.user_data or {}, "error", err=str(context.error)))
        except Exception: pass

# =============================================================================
# MAIN
# =============================================================================

BOT_COMMANDS = [
    BotCommand("start", "الرئيسية"),
    BotCommand("help", "مساعدة"),
    BotCommand("stop", "إنهاء الاختبار الحالي"),
    BotCommand("stats", "إحصائياتي"),
    BotCommand("review", "مراجعة الأخطاء"),
    BotCommand("language", "اللغة"),
    BotCommand("support", "الدعم"),
    BotCommand("osce", "بدء الأوسكي"),
    BotCommand("stage_start", "تجميع (أدمن)"),
    BotCommand("stage_dump", "تصدير (أدمن)"),
    BotCommand("upload_final", "حقن (أدمن)"),
    BotCommand("past_stats", "إحصائيات السنوات السابقة"),
    BotCommand("my_plan", "خطة حسابي"),
    BotCommand("finish_past", "إنهاء جمع السنوات السابقة (أدمن)"),
    BotCommand("ai_policy", "سياسة AI (أدمن)"),
    BotCommand("set_plan", "تفعيل خطة مستخدم (أدمن)"),
    BotCommand("remove_premium", "إلغاء بريميوم (أدمن)")
]

async def post_init(application: Application):
    init_db()
    await cleanup_old_cache(days=30)
    await application.bot.set_my_commands(BOT_COMMANDS)
    logger.info("Bot ready. Model=%s DB=%s", GEMINI_MODEL_NAME, DATABASE_PATH)

def main():
    if not TOKEN: raise RuntimeError("TELEGRAM_BOT_TOKEN is missing.")
    application = ApplicationBuilder().token(TOKEN).concurrent_updates(False).post_init(post_init).build()
    application.add_error_handler(error_handler)

    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("help", help_command))
    application.add_handler(CommandHandler("stop", stop_command))
    application.add_handler(CommandHandler("stats", stats_command))
    application.add_handler(CommandHandler("review", review_command))
    application.add_handler(CommandHandler("language", language_command))
    application.add_handler(CommandHandler("support", support_command))
    application.add_handler(CommandHandler("osce", osce_command))
    application.add_handler(CommandHandler("testbot", testbot_command))
    application.add_handler(CommandHandler("stage_start", stage_start_cmd))
    application.add_handler(CommandHandler("stage_dump", stage_dump_cmd))
    application.add_handler(CommandHandler("upload_final", upload_final_cmd))
    application.add_handler(CommandHandler("past_stats", past_stats_cmd))
    application.add_handler(CommandHandler("my_plan", my_plan_cmd))
    application.add_handler(CommandHandler("finish_past", finish_past_cmd))
    application.add_handler(CommandHandler("ai_policy", ai_policy_cmd))
    application.add_handler(CommandHandler("set_plan", set_plan_cmd))
    application.add_handler(CommandHandler("remove_premium", remove_premium_cmd))

    # Unified callback router: all buttons pass through one protected handler.
    application.add_handler(CallbackQueryHandler(universal_callback_router, pattern=r".*"), group=-100)

    application.add_handler(MessageHandler(filters.PHOTO, handle_photo))
    application.add_handler(MessageHandler(filters.Document.ALL, handle_file))
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))

    logger.info("Starting bot polling...")
    application.run_polling(allowed_updates=Update.ALL_TYPES)

if __name__ == "__main__":
    main()
